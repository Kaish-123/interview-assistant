import sys
import subprocess

import sounddevice as sd
import numpy as np
import wave
import threading
from openai import OpenAI
import tkinter as tk
from tkinter import ttk, font, filedialog, simpledialog, messagebox
import re
import time
import queue
import textract
import pyautogui
from PIL import Image
import io
import base64
import pyperclip
from PIL import ImageGrab
from pynput import keyboard
import tempfile
import json
import os
import Quartz
from dotenv import load_dotenv
import sqlite3
import datetime

# Resolve all data files relative to this script's folder — works no matter
# which directory you launch from (e.g. python3 /full/path/interview_assistant_v2.py)
SCRIPT_DIR = os.path.dirname(os.path.abspath(__file__))

# Optional: faster-whisper for local transcription (graceful fallback)
try:
    from faster_whisper import WhisperModel as FasterWhisperModel
    _FASTER_WHISPER_AVAILABLE = True
except ImportError:
    _FASTER_WHISPER_AVAILABLE = False

# Optional: macOS stealth mode (AppKit may already be imported via Quartz)
try:
    import AppKit as _AppKit
    _APPKIT_AVAILABLE = True
except ImportError:
    _APPKIT_AVAILABLE = False


# ============================================================================
# HELPER FUNCTIONS
# ============================================================================

def estimate_tokens_for_messages(messages: list, optimization_mode: bool = True) -> int:
    """Estimate token count for a list of messages."""
    tokens = 0
    for msg in messages:
        content = msg.get("content", "")
        if isinstance(content, str):
            tokens += len(content) // 4
        elif isinstance(content, list):
            for item in content:
                if isinstance(item, dict):
                    if item.get("type") == "text":
                        tokens += len(item.get("text", "")) // 4
                    elif item.get("type") == "image_url":
                        tokens += 85 if optimization_mode else 765
    return tokens


_QUESTION_PATTERNS = {
    "behavioral": re.compile(
        r"\b(tell me about|describe a time|give me an example|walk me through|"
        r"situation where|how did you handle|conflict with|leadership|teamwork|"
        r"challenging project|greatest weakness|strength)\b", re.I),
    "coding": re.compile(
        r"\b(write|implement|code|algorithm|function|leetcode|data structure|"
        r"complexity|big.?o|array|linked list|tree|graph|dynamic programming|"
        r"recursion|sort|search)\b", re.I),
    "system_design": re.compile(
        r"\b(design|architect|scale|distributed|microservice|database schema|"
        r"api design|load balanc|cache|cdn|message queue|kafka|redis|"
        r"high availability|fault toleran)\b", re.I),
    "technical": re.compile(
        r"\b(explain|what is|how does|difference between|when would you|"
        r"pros and cons|compare|define|concept)\b", re.I),
}

_QUESTION_TYPE_INSTRUCTIONS = {
    "behavioral": (
        "This is a behavioral question. Structure your answer using the STAR method: "
        "Situation, Task, Action, Result. Be specific and concise."
    ),
    "coding": (
        "This is a coding/algorithm question. Provide: 1) Clarifying questions, "
        "2) Approach & complexity, 3) Clean code solution, 4) Edge cases."
    ),
    "system_design": (
        "This is a system design question. Cover: 1) Requirements clarification, "
        "2) High-level design, 3) Key components, 4) Scalability considerations."
    ),
    "technical": (
        "This is a technical concept question. Give a clear, structured explanation "
        "with examples. Keep it interview-appropriate (2-3 minutes spoken)."
    ),
    "general": "",
}

def classify_question(text: str) -> str:
    """Classify question type via regex — no API call."""
    for qtype, pattern in _QUESTION_PATTERNS.items():
        if pattern.search(text):
            return qtype
    return "general"


# ============================================================================
# MATH RENDERER — converts LaTeX notation to readable plain-text / Unicode
# ============================================================================

def _parse_braced(s: str, start: int):
    """Extract balanced-brace content starting at s[start] (must be '{').
    Returns (content_str, absolute_end_pos_after_closing_brace)."""
    depth = 0
    i = start
    while i < len(s):
        if s[i] == '{':
            depth += 1
        elif s[i] == '}':
            depth -= 1
            if depth == 0:
                return s[start + 1:i], i + 1
        i += 1
    return s[start + 1:], len(s)


_LATEX_SYMBOLS = sorted([
    # Arithmetic / relations
    (r'\times', '×'), (r'\cdot', '·'), (r'\div', '÷'),
    (r'\leq', '≤'), (r'\geq', '≥'), (r'\neq', '≠'),
    (r'\approx', '≈'), (r'\equiv', '≡'), (r'\pm', '±'), (r'\mp', '∓'),
    (r'\infty', '∞'), (r'\partial', '∂'), (r'\nabla', '∇'),
    (r'\int', '∫'), (r'\sum', 'Σ'), (r'\prod', 'Π'),
    # Arrows
    (r'\rightarrow', '→'), (r'\leftarrow', '←'), (r'\Rightarrow', '⟹'),
    (r'\Leftarrow', '⟸'), (r'\leftrightarrow', '↔'), (r'\to', '→'),
    # Sets / logic
    (r'\in', '∈'), (r'\notin', '∉'), (r'\subset', '⊂'), (r'\subseteq', '⊆'),
    (r'\cup', '∪'), (r'\cap', '∩'), (r'\forall', '∀'), (r'\exists', '∃'),
    (r'\neg', '¬'), (r'\land', '∧'), (r'\lor', '∨'),
    # Greek lowercase
    (r'\alpha', 'α'), (r'\beta', 'β'), (r'\gamma', 'γ'), (r'\delta', 'δ'),
    (r'\epsilon', 'ε'), (r'\varepsilon', 'ε'), (r'\zeta', 'ζ'), (r'\eta', 'η'),
    (r'\theta', 'θ'), (r'\vartheta', 'θ'), (r'\iota', 'ι'), (r'\kappa', 'κ'),
    (r'\lambda', 'λ'), (r'\mu', 'μ'), (r'\nu', 'ν'), (r'\xi', 'ξ'),
    (r'\varpi', 'ϖ'), (r'\pi', 'π'), (r'\rho', 'ρ'), (r'\varrho', 'ϱ'),
    (r'\sigma', 'σ'), (r'\varsigma', 'ς'), (r'\tau', 'τ'), (r'\upsilon', 'υ'),
    (r'\varphi', 'φ'), (r'\phi', 'φ'), (r'\chi', 'χ'), (r'\psi', 'ψ'),
    (r'\omega', 'ω'),
    # Greek uppercase
    (r'\Gamma', 'Γ'), (r'\Delta', 'Δ'), (r'\Theta', 'Θ'), (r'\Lambda', 'Λ'),
    (r'\Xi', 'Ξ'), (r'\Pi', 'Π'), (r'\Sigma', 'Σ'), (r'\Upsilon', 'Υ'),
    (r'\Phi', 'Φ'), (r'\Psi', 'Ψ'), (r'\Omega', 'Ω'),
    # Misc
    (r'\ldots', '…'), (r'\cdots', '⋯'), (r'\vdots', '⋮'), (r'\ddots', '⋱'),
    (r'\langle', '⟨'), (r'\rangle', '⟩'),
    (r'\quad', '  '), (r'\qquad', '    '),
    (r'\left', ''), (r'\right', ''),   # strip delimiter modifiers; keep actual char
    (r'\sqrt', '√'),
    (r'\%', '%'), (r'\,', ' '), (r'\;', ' '), (r'\:', ' '), (r'\!', ''),
], key=lambda x: -len(x[0]))  # longest-first so \leftarrow beats \left

_SUP_TRANS = str.maketrans('0123456789+-=()n', '⁰¹²³⁴⁵⁶⁷⁸⁹⁺⁻⁼⁽⁾ⁿ')
_SUB_TRANS = str.maketrans('0123456789+-=()', '₀₁₂₃₄₅₆₇₈₉₊₋₌₍₎')


def _process_math_expr(s: str) -> str:
    """Recursively convert a LaTeX math expression to readable Unicode text."""
    result = []
    i = 0
    n = len(s)
    while i < n:
        c = s[i]

        # \frac{numerator}{denominator}
        if s[i:i+5] == r'\frac' and i + 5 < n and s[i + 5] == '{':
            num_raw, after_num = _parse_braced(s, i + 5)
            if after_num < n and s[after_num] == '{':
                den_raw, after_den = _parse_braced(s, after_num)
                nr = _process_math_expr(num_raw)
                dr = _process_math_expr(den_raw)
                result.append(f"({nr}) / ({dr})")
                i = after_den
                continue

        # \sqrt{...}  (with optional [n] which we ignore)
        if s[i:i+5] == r'\sqrt':
            j = i + 5
            if j < n and s[j] == '[':        # optional root index, skip it
                end_bracket = s.find(']', j)
                j = end_bracket + 1 if end_bracket != -1 else j
            if j < n and s[j] == '{':
                inner, after = _parse_braced(s, j)
                result.append(f"√({_process_math_expr(inner)})")
                i = after
                continue

        # \text{}, \textbf{}, \mathrm{}, \mathbf{}, \operatorname{}, etc.
        _TEXT_CMD = re.match(
            r'\\(?:text(?:bf|it|rm|sf|tt)?|math(?:bf|it|rm|sf|tt|cal|bb)?'
            r'|operatorname\*?|mbox|hbox)\{', s[i:])
        if _TEXT_CMD:
            inner, after = _parse_braced(s, i + _TEXT_CMD.end() - 1)
            result.append(inner)    # keep text content as-is
            i = after
            continue

        # ^{...} or ^x  — superscript
        if c == '^' and i + 1 < n:
            if s[i + 1] == '{':
                inner, after = _parse_braced(s, i + 1)
                inner_p = _process_math_expr(inner)
                sup = inner_p.translate(_SUP_TRANS) if all(
                    ch in '0123456789+-=()n' for ch in inner_p) else f"^({inner_p})"
                result.append(sup)
                i = after
            else:
                ch = s[i + 1]
                sup = ch.translate(_SUP_TRANS) if ch in '0123456789+-=()n' else f"^{ch}"
                result.append(sup)
                i += 2
            continue

        # _{...} or _x  — subscript
        if c == '_' and i + 1 < n:
            if s[i + 1] == '{':
                inner, after = _parse_braced(s, i + 1)
                inner_p = _process_math_expr(inner)
                sub = inner_p.translate(_SUB_TRANS) if all(
                    ch in '0123456789+-=()' for ch in inner_p) else f"_({inner_p})"
                result.append(sub)
                i = after
            else:
                ch = s[i + 1]
                sub = ch.translate(_SUB_TRANS) if ch in '0123456789+-=()' else f"_{ch}"
                result.append(sub)
                i += 2
            continue

        # \\ — line break inside aligned / multiline math
        if c == '\\' and i + 1 < n and s[i + 1] == '\\':
            result.append('\n    ')
            i += 2
            continue

        # Backslash — symbol map or unknown command
        if c == '\\':
            matched = False
            for latex_sym, uni_sym in _LATEX_SYMBOLS:
                if s[i:].startswith(latex_sym):
                    end_pos = i + len(latex_sym)
                    # Don't match if this is a prefix of a longer alpha command
                    if latex_sym[-1].isalpha() and end_pos < n and s[end_pos].isalpha():
                        continue
                    result.append(uni_sym)
                    i = end_pos
                    matched = True
                    break
            if not matched:
                # Skip unknown \command
                j = i + 1
                while j < n and s[j].isalpha():
                    j += 1
                i = j
            continue

        # Balanced braces — recurse into content
        if c == '{':
            inner, after = _parse_braced(s, i)
            result.append(_process_math_expr(inner))
            i = after
            continue

        # Stray closing brace — skip
        if c == '}':
            i += 1
            continue

        result.append(c)
        i += 1

    return ''.join(result)


def format_math_for_display(text: str) -> str:
    """Post-process a GPT response, converting LaTeX math blocks to readable plain text.

    Handles:
      - Display math: \\[...\\] and $$...$$  → boxed block
      - Inline math:  $...$                  → rendered inline
      - \\frac{}{}, \\sqrt{}, ^{}, _{}, Greek letters, all common symbols
    Code fences are left completely untouched.
    """
    if '\\' not in text and '$' not in text:
        return text  # fast path — nothing to do

    # -- 1. Protect code blocks from modification --
    code_blocks: dict = {}

    def _save_code(m):
        key = f"\x00CODE{len(code_blocks)}\x00"
        code_blocks[key] = m.group(0)
        return key

    text = re.sub(r'```.*?```', _save_code, text, flags=re.DOTALL)
    text = re.sub(r'`[^`\n]+`', _save_code, text)

    # -- 2. Display math  \[...\]  and  $$...$$  → boxed block --
    def _display_math(m):
        raw = m.group(1).strip()
        # Split on \\ (line-break inside aligned blocks)
        lines = [ln.strip() for ln in re.split(r'\\\\', raw) if ln.strip()]
        if not lines:
            return ''
        converted = [_process_math_expr(ln) for ln in lines]
        width = max(len(c) for c in converted) + 2
        bar = '─' * width
        box = [f'\n  ┌{bar}┐']
        for c in converted:
            box.append(f'  │  {c}')
        box.append(f'  └{bar}┘\n')
        return '\n'.join(box)

    text = re.sub(r'\\\[(.*?)\\\]', _display_math, text, flags=re.DOTALL)
    text = re.sub(r'\$\$(.*?)\$\$', _display_math, text, flags=re.DOTALL)

    # -- 3. Inline math  $...$  → rendered inline --
    def _inline_math(m):
        return _process_math_expr(m.group(1))

    text = re.sub(r'\$(?!\$)([^$\n]+?)\$', _inline_math, text)

    # -- 4. Restore code blocks --
    for key, val in code_blocks.items():
        text = text.replace(key, val)

    return text


# ============================================================================
# IMAGE OPTIMIZATION HELPERS (Safe - doesn't delete anything, just compresses)
# ============================================================================

def compress_image_for_api(image: Image.Image, max_size: int = 1024, quality: int = 85) -> str:
    """
    Compress an image for API transmission while preserving visual quality.
    - Resizes if larger than max_size (keeps aspect ratio)
    - Compresses to JPEG for smaller payload
    - Returns base64 string
    
    This reduces payload by 60-80% without losing meaningful detail for GPT.
    """
    # Resize if too large (keeping aspect ratio)
    width, height = image.size
    if width > max_size or height > max_size:
        ratio = min(max_size / width, max_size / height)
        new_size = (int(width * ratio), int(height * ratio))
        image = image.resize(new_size, Image.Resampling.LANCZOS)
    
    # Convert to RGB if necessary (for JPEG)
    if image.mode in ('RGBA', 'P'):
        background = Image.new('RGB', image.size, (255, 255, 255))
        if image.mode == 'RGBA':
            background.paste(image, mask=image.split()[3])
        else:
            background.paste(image)
        image = background
    
    # Compress to JPEG
    buffer = io.BytesIO()
    image.save(buffer, format="JPEG", quality=quality, optimize=True)
    return base64.b64encode(buffer.getvalue()).decode("utf-8")


def compress_image_png(image: Image.Image, max_size: int = 1024) -> str:
    """
    Compress image as PNG (for screenshots with text that need sharpness).
    """
    width, height = image.size
    if width > max_size or height > max_size:
        ratio = min(max_size / width, max_size / height)
        new_size = (int(width * ratio), int(height * ratio))
        image = image.resize(new_size, Image.Resampling.LANCZOS)
    
    buffer = io.BytesIO()
    image.save(buffer, format="PNG", optimize=True)
    return base64.b64encode(buffer.getvalue()).decode("utf-8")




# ============================================================================
# ROBUST FILE TEXT EXTRACTION (no textract dependency for common types)
# ============================================================================

# Extensions that can be read directly as plain text / code
_TEXT_EXTS = {
    '.txt', '.py', '.js', '.ts', '.jsx', '.tsx', '.java', '.c', '.cpp', '.h',
    '.hpp', '.cs', '.go', '.rs', '.rb', '.php', '.swift', '.kt', '.sh', '.bat',
    '.ps1', '.sql', '.r', '.lua', '.pl', '.yaml', '.yml', '.json', '.xml',
    '.html', '.htm', '.css', '.scss', '.sass', '.md', '.markdown', '.rst',
    '.csv', '.tsv', '.ini', '.cfg', '.toml', '.env', '.gitignore', '.makefile',
    '.log', '.conf', '.properties',
}

# Extensions to skip entirely (binaries that yield no useful text)
_SKIP_EXTS = {
    '.png', '.jpg', '.jpeg', '.gif', '.bmp', '.svg', '.ico', '.webp',
    '.mp3', '.mp4', '.wav', '.avi', '.mov', '.mkv',
    '.zip', '.tar', '.gz', '.7z', '.rar',
    '.exe', '.dll', '.so', '.dylib', '.bin', '.o', '.class',
    '.pyc', '.pyo',
    '.DS_Store', '.db', '.sqlite',
}


def extract_text_from_file(file_path: str) -> str:
    """
    Robustly extract readable text from any file type.
    Tries native pure-Python readers first; falls back to textract as last resort.
    Raises RuntimeError if nothing can read the file.
    """
    ext = os.path.splitext(file_path)[1].lower()
    basename = os.path.basename(file_path)

    # Skip known binary/media files early
    if ext in _SKIP_EXTS:
        raise RuntimeError(f"Skipped binary/media file: {basename}")

    # ── Plain text / source code ──────────────────────────────────────────
    if ext in _TEXT_EXTS or ext == '':
        for enc in ('utf-8', 'latin-1', 'cp1252'):
            try:
                with open(file_path, 'r', encoding=enc) as f:
                    return f.read()
            except UnicodeDecodeError:
                continue
        # Binary fallback: decode with replacement
        with open(file_path, 'rb') as fb:
            return fb.read().decode('utf-8', errors='replace')

    # ── PDF ───────────────────────────────────────────────────────────────
    if ext == '.pdf':
        try:
            import pdfplumber
            with pdfplumber.open(file_path) as pdf:
                pages = [page.extract_text() or '' for page in pdf.pages]
                text = '\n'.join(pages).strip()
                if text:
                    return text
        except Exception:
            pass
        try:
            from pypdf import PdfReader
            reader = PdfReader(file_path)
            text = '\n'.join(p.extract_text() or '' for p in reader.pages).strip()
            if text:
                return text
        except Exception:
            pass

    # ── DOCX (Word) ───────────────────────────────────────────────────────
    if ext == '.docx':
        # Primary: python-docx (best paragraph/table handling)
        try:
            from docx import Document
            doc = Document(file_path)
            return '\n'.join(p.text for p in doc.paragraphs)
        except ImportError:
            pass            # python-docx not installed — fall through to stdlib
        except Exception:
            pass            # corrupt / edge-case — still try the zipfile path

        # Fallback: pure stdlib zipfile + xml.etree (zero extra dependencies)
        try:
            import zipfile
            import xml.etree.ElementTree as ET
            WNS = '{http://schemas.openxmlformats.org/wordprocessingml/2006/main}'
            with zipfile.ZipFile(file_path, 'r') as z:
                with z.open('word/document.xml') as fxml:
                    root = ET.parse(fxml).getroot()
            paras = []
            for para in root.iter(f'{WNS}p'):
                runs = [t.text for t in para.iter(f'{WNS}t') if t.text]
                paras.append(''.join(runs))
            return '\n'.join(paras)
        except Exception as e:
            raise RuntimeError(f"DOCX read failed for {basename}: {e}")

    # ── XLSX / XLS (Excel) ────────────────────────────────────────────────
    if ext in ('.xlsx', '.xlsm', '.xls'):
        try:
            import openpyxl
            wb = openpyxl.load_workbook(file_path, read_only=True, data_only=True)
            rows = []
            for sheet in wb.worksheets:
                rows.append(f"[Sheet: {sheet.title}]")
                for row in sheet.iter_rows(values_only=True):
                    rows.append('\t'.join(str(c) if c is not None else '' for c in row))
            return '\n'.join(rows)
        except Exception as e:
            raise RuntimeError(f"Excel read failed for {basename}: {e}")

    # ── PPTX (PowerPoint) ─────────────────────────────────────────────────
    if ext == '.pptx':
        try:
            from pptx import Presentation
            prs = Presentation(file_path)
            lines = []
            for i, slide in enumerate(prs.slides, 1):
                lines.append(f"[Slide {i}]")
                for shape in slide.shapes:
                    if hasattr(shape, 'text') and shape.text.strip():
                        lines.append(shape.text)
            return '\n'.join(lines)
        except Exception as e:
            raise RuntimeError(f"PPTX read failed for {basename}: {e}")

    # ── Fallback: textract ────────────────────────────────────────────────
    try:
        import textract as _textract
        result = _textract.process(file_path)
        return result.decode('utf-8', errors='replace')
    except Exception as e:
        raise RuntimeError(f"Could not extract text from {basename}: {e}")


def collect_files_from_folder(folder_path: str) -> list:
    """
    Recursively collect all readable files from a folder.
    Skips hidden directories, build artifacts, and known binary extensions.
    Returns a list of absolute file paths.
    """
    SKIP_DIRS = {
        '.git', '__pycache__', 'node_modules', '.venv', 'venv', 'env',
        '.idea', '.vscode', 'dist', 'build', '.mypy_cache', '.pytest_cache',
    }
    MAX_FILE_SIZE = 5 * 1024 * 1024  # 5 MB per file
    paths = []
    for root, dirs, files in os.walk(folder_path):
        # Prune dirs in-place to avoid descending into them
        dirs[:] = sorted(
            d for d in dirs
            if d not in SKIP_DIRS and not d.startswith('.')
        )
        for fname in sorted(files):
            if fname.startswith('.'):
                continue
            fpath = os.path.join(root, fname)
            try:
                if os.path.getsize(fpath) <= MAX_FILE_SIZE:
                    paths.append(fpath)
            except OSError:
                pass
    return paths


# ============================================================================

class UIPreferences:
    FILE = os.path.join(SCRIPT_DIR, "ui_prefs.json")

    @staticmethod
    def load():
        try:
            if os.path.exists(UIPreferences.FILE):
                with open(UIPreferences.FILE, "r") as f:
                    return json.load(f)
        except Exception as e:
            print("UI Prefs load error:", e)
        return {}

    @staticmethod
    def save(data: dict):
        try:
            with open(UIPreferences.FILE, "w") as f:
                json.dump(data, f, indent=2)
        except Exception as e:
            print("UI Prefs save error:", e)

    
class PromptManager:
    def __init__(self, tabs_file_path=None, prompts_file_path=None):
        tabs_file_path = tabs_file_path or os.path.join(SCRIPT_DIR, "tabs.json")
        prompts_file_path = prompts_file_path or os.path.join(SCRIPT_DIR, "prompts.json")
        self.tabs_file_path = tabs_file_path
        self.prompts_file_path = prompts_file_path
        self.data = {"tabs": []}
        self.load_tabs()

    def load_tabs(self):
        # Load tabs data from tabs.json
        if os.path.exists(self.tabs_file_path):
            try:
                with open(self.tabs_file_path, 'r') as f:
                    self.data = json.load(f)
            except Exception as e:
                print(f"Error loading tabs data: {str(e)}")
                self.data = {"tabs": []}
    
    def save_tabs(self):
        # Save tabs data to tabs.json
        try:
            with open(self.tabs_file_path, 'w') as f:
                json.dump(self.data, f, indent=2)
        except Exception as e:
            print(f"Error saving tabs data: {str(e)}")

    def add_tab(self, name):
        # Adds a new tab and saves it
        self.data["tabs"].append({"name": name, "subTabs": []})
        self.save_tabs()
        return len(self.data["tabs"]) - 1
    
    def add_subtab(self, tab_index, name, prompt="", text_input=""):
        # Adds a new subtab, including prompt and text_input, and saves it
        if 0 <= tab_index < len(self.data["tabs"]):
            subtab_data = {
                "name": name,
                "prompt": prompt,  # Save the prompt
                "text_input": text_input  # Save the text input
            }
            self.data["tabs"][tab_index]["subTabs"].append(subtab_data)
            self.save_tabs()  # Ensure changes are saved immediately
            return len(self.data["tabs"][tab_index]["subTabs"]) - 1
        return -1


    def get_tab_count(self):
        return len(self.data["tabs"])
    
    def get_tab_name(self, index):
        if 0 <= index < len(self.data["tabs"]):
            return self.data["tabs"][index]["name"]
        return ""
    
    def get_subtab_count(self, tab_index):
        if 0 <= tab_index < len(self.data["tabs"]):
            return len(self.data["tabs"][tab_index]["subTabs"])
        return 0
    
    def get_subtab_text_input(self, tab_index, subtab_index):
        # Retrieves the text_input from the specified subtab
        if (0 <= tab_index < len(self.data["tabs"]) and
            0 <= subtab_index < len(self.data["tabs"][tab_index]["subTabs"])):
            return self.data["tabs"][tab_index]["subTabs"][subtab_index].get("text_input", "")
        return ""

    def get_subtab_name(self, tab_index, subtab_index):
        if (0 <= tab_index < len(self.data["tabs"]) and 
            0 <= subtab_index < len(self.data["tabs"][tab_index]["subTabs"])):
            return self.data["tabs"][tab_index]["subTabs"][subtab_index]["name"]
        return ""
    
    def get_subtab_prompt(self, tab_index, subtab_index):
        if (0 <= tab_index < len(self.data["tabs"]) and 
            0 <= subtab_index < len(self.data["tabs"][tab_index]["subTabs"])):
            return self.data["tabs"][tab_index]["subTabs"][subtab_index]["prompt"]
        return ""
    
    def update_subtab_prompt(self, tab_index, subtab_index, prompt, text_input=""):
        # Updates the prompt and text_input, then saves the data
        if 0 <= tab_index < len(self.data["tabs"]) and 0 <= subtab_index < len(self.data["tabs"][tab_index]["subTabs"]):
            self.data["tabs"][tab_index]["subTabs"][subtab_index]["prompt"] = prompt
            self.data["tabs"][tab_index]["subTabs"][subtab_index]["text_input"] = text_input
            self.save_tabs()  # Save the changes
            return True
        return False


def get_window_under_mouse():
    mouse_x, mouse_y = pyautogui.position()
    window_info_list = Quartz.CGWindowListCopyWindowInfo(
        Quartz.kCGWindowListOptionOnScreenOnly, Quartz.kCGNullWindowID)
    for window in window_info_list:
        bounds = window.get('kCGWindowBounds', {})
        x = bounds.get('X', 0)
        y = bounds.get('Y', 0)
        width = bounds.get('Width', 0)
        height = bounds.get('Height', 0)
        if x <= mouse_x <= x + width and y <= mouse_y <= y + height:
            return window
    return None

# Example usage:
# window = get_window_under_mouse()
# if window:
#     print(f"Window under mouse: {window.get('kCGWindowName', 'No Title')}")


def get_window_list():
    window_list = []
    window_info_list = Quartz.CGWindowListCopyWindowInfo(
        Quartz.kCGWindowListOptionOnScreenOnly, Quartz.kCGNullWindowID)
    for window_info in window_info_list:
        window_list.append(window_info)
    return window_list

# Example usage:
# windows = get_window_list()
# for window in windows:
#     print(window.get('kCGWindowName', 'No Title'))


def on_activate():
    if not app.assistant.recorder.is_recording:
        print("🎤 Global hotkey: Start Listening")
        app.toggle_recording()
    else:
        print("🛑 Global hotkey: Stop & Process")
        app.toggle_recording()






# Configuration
load_dotenv(os.path.join(SCRIPT_DIR, ".env"))
API_KEY = os.getenv("OPENAI_API_KEY")
client = OpenAI(api_key=API_KEY)
SAMPLE_RATE = 16000
CHANNELS = 1
DTYPE = 'int16'
CHUNK = 1024
BLACKHOLE_DEVICE = "BlackHole"

# Live preview Whisper: only the last N seconds + min interval between calls.
# Sending the entire growing recording every ~2s was the main API/latency cost.
LIVE_WHISPER_TAIL_SECONDS = 8.0       # v2: reduced from 12s → less data per poll
LIVE_WHISPER_MIN_INTERVAL_SEC = 2.0   # v2: reduced from 2.5s → slightly snappier live preview

# Prompt caching: OpenAI caches prompts >1024 tokens for 5 min (50% cheaper + faster).
# We trigger it by sending large system messages (resume/JD) as the very first message,
# identical every call. The flag below enables cache_control tagging.
ENABLE_PROMPT_CACHE = True

class AudioRecorder:
    def __init__(self):
        self.frames = []
        self.is_recording = False
        self.stream = None
        self.audio_queue = queue.Queue()
        self.input_mode = "internal"  # internal = BlackHole, external = mic
        self.lock = threading.Lock()
        self.process_thread = None      # ✅ NEW

    def find_device(self):
        """
        Resolve a sounddevice input device index based on self.input_mode.

        - internal  => prefer BLACKHOLE_DEVICE
        - external  => prefer first non-BlackHole input device

        If nothing matches, return None so sounddevice uses its default.
        """
        try:
            devices = sd.query_devices()
        except Exception as e:
            print(f"⚠️ Could not query audio devices: {e}")
            return None

        target_index = None
        bh_name = BLACKHOLE_DEVICE.lower()

        if self.input_mode == "internal":
            # Prefer BlackHole for internal audio
            for idx, dev in enumerate(devices):
                try:
                    if dev.get("max_input_channels", 0) > 0 and bh_name in dev.get("name", "").lower():
                        target_index = idx
                        break
                except Exception:
                    continue
        else:
            # Prefer first *non*-BlackHole input device as "external" mic
            for idx, dev in enumerate(devices):
                try:
                    if dev.get("max_input_channels", 0) > 0 and bh_name not in dev.get("name", "").lower():
                        target_index = idx
                        break
                except Exception:
                    continue

        if target_index is None:
            print(f"⚠️ No specific device found for mode={self.input_mode!r}, falling back to default.")
            return None  # let sounddevice pick default

        try:
            print(f"🎧 Using device #{target_index}: {devices[target_index]['name']} (mode={self.input_mode})")
        except Exception:
            pass

        return target_index

    def get_snapshot(self):
        """Return a copy of all audio recorded so far, or None."""
        with self.lock:
            if not self.frames:
                return None
            return np.concatenate(self.frames).copy()

    def get_tail_snapshot(self, n_samples: int):
        """
        Return at most the last n_samples from the current recording (copy), or None.
        Used by live transcription to avoid copying the entire growing buffer every poll.
        """
        with self.lock:
            if not self.frames:
                return None
            # Walk frames from the end until we have enough samples — avoids full concat
            collected = []
            remaining = n_samples
            for frame in reversed(self.frames):
                collected.append(frame)
                remaining -= len(frame)
                if remaining <= 0:
                    break
            arr = np.concatenate(list(reversed(collected)))
            return arr[-n_samples:] if len(arr) > n_samples else arr.copy()

    def start_recording(self):
        device_id = self.find_device()
        self.frames = []
        self.is_recording = True

        def callback(indata, frames, time, status):
            if self.is_recording:
                self.audio_queue.put(indata.copy())

        self.stream = sd.InputStream(
            samplerate=SAMPLE_RATE,
            channels=CHANNELS,
            dtype=DTYPE,
            callback=callback,
            device=device_id,
            blocksize=CHUNK
        )
        self.stream.start()
        # ✅ track the thread so we can join it on stop
        self.process_thread = threading.Thread(target=self.process_audio, daemon=True)
        self.process_thread.start()

    def process_audio(self):
        while self.is_recording or not self.audio_queue.empty():
            try:
                frame = self.audio_queue.get(timeout=0.1)
            except queue.Empty:
                continue
            with self.lock:
                self.frames.append(frame)

    def stop_recording(self, filename="interviewer.wav"):
        # 🔚 stop callbacks first
        self.is_recording = False
        if self.stream:
            self.stream.stop()
            self.stream.close()

        # ✅ wait for the process_audio thread to drain the queue
        if self.process_thread and self.process_thread.is_alive():
            self.process_thread.join(timeout=1.0)

        with self.lock:
            frames_copy = list(self.frames)

        if frames_copy:
            audio_data = np.concatenate(frames_copy)
            with wave.open(filename, 'wb') as wf:
                wf.setnchannels(CHANNELS)
                wf.setsampwidth(2)
                wf.setframerate(SAMPLE_RATE)
                wf.writeframes(audio_data.tobytes())
        return filename



class ChatHistoryManager:
    def __init__(self, file_path=None):
        self.file_path = file_path or os.path.join(SCRIPT_DIR, "chats.json")
        self.sessions = []  # Each item: {"title": str, "messages": List[dict], "bookmarks": List}
        self.load()

    def save(self, force=False):
        with open(self.file_path, "w") as f:
            json.dump(self.sessions, f, indent=2)

    def load(self):
        if os.path.exists(self.file_path):
            try:
                with open(self.file_path, "r") as f:
                    self.sessions = json.load(f)
            except Exception:
                self.sessions = []

    def save_current_session(self, messages, title="AutoSave - Last Session", bookmarks=None):
        # Preserve any existing bookmarks for index-0 if none are supplied
        if bookmarks is None and self.sessions:
            bookmarks = self.sessions[0].get("bookmarks", [])
        working_session = {
            "title": title,
            "messages": messages.copy(),
            "bookmarks": bookmarks or [],
        }
        if self.sessions:
            self.sessions[0] = working_session
        else:
            self.sessions.insert(0, working_session)
        self.save()

    def add_session(self, title, messages, bookmarks=None):
        self.sessions.append({
            "title": title,
            "messages": messages,
            "bookmarks": bookmarks or [],
        })
        self.save()

    def get_titles(self):
        return [s.get("title", "Untitled") for s in self.sessions]

    def get_session(self, index):
        return self.sessions[index]["messages"] if 0 <= index < len(self.sessions) else []

    def get_session_bookmarks(self, index):
        """Return the saved bookmarks list for a session (list of [line_index, preview])."""
        if 0 <= index < len(self.sessions):
            return self.sessions[index].get("bookmarks", [])
        return []

    def update_session_bookmarks(self, index, bookmarks):
        """Overwrite the bookmarks for a session and persist immediately."""
        if 0 <= index < len(self.sessions):
            self.sessions[index]["bookmarks"] = bookmarks
            self.save()


class SQLiteChatHistoryManager(ChatHistoryManager):
    """
    SQLite-backed drop-in for ChatHistoryManager.
    Inherits the full list-based interface; overrides load/save to persist in SQLite.
    Falls back transparently to JSON if SQLite init fails.
    """

    def __init__(self, db_path=None, json_fallback_path=None):
        db_path = db_path or os.path.join(SCRIPT_DIR, "chats_v2.db")
        json_fallback_path = json_fallback_path or os.path.join(SCRIPT_DIR, "chats.json")
        self.sessions = []
        self._use_sqlite = False
        self.file_path = json_fallback_path  # used by parent save/load fallback
        try:
            self.conn = sqlite3.connect(db_path, check_same_thread=False)
            self._init_schema()
            self._use_sqlite = True
            self.load()
        except Exception:
            self._use_sqlite = False
            super().__init__(json_fallback_path)

    def _init_schema(self):
        c = self.conn.cursor()
        c.execute("""CREATE TABLE IF NOT EXISTS sessions (
            id INTEGER PRIMARY KEY AUTOINCREMENT,
            title TEXT,
            created_at TEXT
        )""")
        c.execute("""CREATE TABLE IF NOT EXISTS messages (
            id INTEGER PRIMARY KEY AUTOINCREMENT,
            session_id INTEGER,
            role TEXT,
            content TEXT,
            FOREIGN KEY(session_id) REFERENCES sessions(id) ON DELETE CASCADE
        )""")
        c.execute("""CREATE TABLE IF NOT EXISTS bookmarks (
            id INTEGER PRIMARY KEY AUTOINCREMENT,
            session_id INTEGER,
            line_index INTEGER,
            preview TEXT,
            FOREIGN KEY(session_id) REFERENCES sessions(id) ON DELETE CASCADE
        )""")
        self.conn.commit()

    def load(self):
        if not self._use_sqlite:
            return super().load()
        c = self.conn.cursor()
        c.execute("SELECT id, title, created_at FROM sessions ORDER BY id ASC")
        rows = c.fetchall()
        self.sessions = []
        for sid, title, created_at in rows:
            c2 = self.conn.cursor()
            c2.execute("SELECT role, content FROM messages WHERE session_id=? ORDER BY id", (sid,))
            msgs = [{"role": r, "content": cnt} for r, cnt in c2.fetchall()]
            c3 = self.conn.cursor()
            c3.execute("SELECT line_index, preview FROM bookmarks WHERE session_id=? ORDER BY id", (sid,))
            bks = [[li, pv] for li, pv in c3.fetchall()]
            self.sessions.append({"title": title, "messages": msgs, "bookmarks": bks, "_db_id": sid})

    def save(self, force=False):
        if not self._use_sqlite:
            return super().save(force)
        c = self.conn.cursor()
        for i, s in enumerate(self.sessions):
            db_id = s.get("_db_id")
            if db_id is None:
                c.execute("INSERT INTO sessions(title, created_at) VALUES(?,?)",
                          (s.get("title", "Chat"), datetime.datetime.utcnow().isoformat()))
                db_id = c.lastrowid
                self.sessions[i]["_db_id"] = db_id
            else:
                c.execute("UPDATE sessions SET title=? WHERE id=?", (s.get("title", "Chat"), db_id))
            c.execute("DELETE FROM messages WHERE session_id=?", (db_id,))
            for msg in s.get("messages", []):
                role = msg.get("role", "")
                content = msg.get("content", "")
                if not isinstance(content, str):
                    content = json.dumps(content)
                c.execute("INSERT INTO messages(session_id, role, content) VALUES(?,?,?)",
                          (db_id, role, content))
            c.execute("DELETE FROM bookmarks WHERE session_id=?", (db_id,))
            for bk in s.get("bookmarks", []):
                li = bk[0] if isinstance(bk, (list, tuple)) else 0
                pv = bk[1] if isinstance(bk, (list, tuple)) and len(bk) > 1 else ""
                c.execute("INSERT INTO bookmarks(session_id, line_index, preview) VALUES(?,?,?)",
                          (db_id, li, pv))
        # Remove DB rows for sessions deleted from the list
        current_ids = {s["_db_id"] for s in self.sessions if "_db_id" in s}
        c.execute("SELECT id FROM sessions")
        all_ids = {row[0] for row in c.fetchall()}
        for stale_id in all_ids - current_ids:
            c.execute("DELETE FROM sessions WHERE id=?", (stale_id,))
            c.execute("DELETE FROM messages WHERE session_id=?", (stale_id,))
            c.execute("DELETE FROM bookmarks WHERE session_id=?", (stale_id,))
        self.conn.commit()


class MockInterviewDialog(tk.Toplevel):
    """Interactive mock interview session with AI scoring."""

    def __init__(self, parent, openai_client, role="Software Engineer"):
        super().__init__(parent)
        self.title("🎯 Mock Interview")
        self.geometry("700x600")
        self.resizable(True, True)
        self.client = openai_client
        self.role = role
        self.questions = []
        self.answers = []
        self.scores = []
        self.current_q = 0
        self._build_ui()
        self._generate_questions()

    def _build_ui(self):
        top = ttk.Frame(self)
        top.pack(fill="x", padx=10, pady=8)
        ttk.Label(top, text="Mock Interview", font=("Helvetica", 14, "bold")).pack(side="left")
        self.progress_lbl = ttk.Label(top, text="Generating questions...")
        self.progress_lbl.pack(side="right")

        self.q_frame = ttk.LabelFrame(self, text="Question")
        self.q_frame.pack(fill="x", padx=10, pady=4)
        self.q_text = tk.Text(self.q_frame, height=4, wrap="word", state="disabled",
                              font=("Helvetica", 12))
        self.q_text.pack(fill="x", padx=5, pady=5)

        ans_frame = ttk.LabelFrame(self, text="Your Answer")
        ans_frame.pack(fill="both", expand=True, padx=10, pady=4)
        self.ans_text = tk.Text(ans_frame, wrap="word", font=("Helvetica", 11))
        self.ans_text.pack(fill="both", expand=True, padx=5, pady=5)

        btn_row = ttk.Frame(self)
        btn_row.pack(fill="x", padx=10, pady=8)
        ttk.Button(btn_row, text="⬅ Prev", command=self._prev_q).pack(side="left", padx=4)
        ttk.Button(btn_row, text="Submit & Score", command=self._submit_answer).pack(side="left", padx=4)
        ttk.Button(btn_row, text="Next ➡", command=self._next_q).pack(side="left", padx=4)
        ttk.Button(btn_row, text="📊 Report Card", command=self._show_report).pack(side="right", padx=4)

        self.feedback_var = tk.StringVar(value="")
        ttk.Label(self, textvariable=self.feedback_var, foreground="#2196F3",
                  wraplength=660).pack(padx=10, pady=4)

    def _generate_questions(self):
        def worker():
            prompt = (f"Generate 5 realistic interview questions for a {self.role} role. "
                      "Mix behavioral, coding, and technical questions. "
                      "Return as a numbered list, one question per line, no extra commentary.")
            try:
                resp = self.client.chat.completions.create(
                    model="gpt-4o-mini",
                    messages=[{"role": "user", "content": prompt}],
                    max_tokens=500
                )
                raw = resp.choices[0].message.content.strip()
                lines = [l.strip() for l in raw.split("\n") if l.strip() and l[0].isdigit()]
                self.questions = [re.sub(r"^\d+[\.\)]\s*", "", l) for l in lines]
                self.answers = [""] * len(self.questions)
                self.scores = [None] * len(self.questions)
                self.after(0, self._show_question)
            except Exception as e:
                self.after(0, lambda: self.feedback_var.set(f"Error generating questions: {e}"))
        threading.Thread(target=worker, daemon=True).start()

    def _show_question(self):
        if not self.questions:
            return
        idx = self.current_q
        self.progress_lbl.config(text=f"Question {idx+1} / {len(self.questions)}")
        self.q_text.config(state="normal")
        self.q_text.delete("1.0", "end")
        self.q_text.insert("1.0", self.questions[idx])
        self.q_text.config(state="disabled")
        self.ans_text.delete("1.0", "end")
        if self.answers[idx]:
            self.ans_text.insert("1.0", self.answers[idx])
        score_str = f"  [Score: {self.scores[idx]}/10]" if self.scores[idx] is not None else ""
        self.feedback_var.set(score_str)

    def _prev_q(self):
        if self.current_q > 0:
            self._save_current_answer()
            self.current_q -= 1
            self._show_question()

    def _next_q(self):
        if self.current_q < len(self.questions) - 1:
            self._save_current_answer()
            self.current_q += 1
            self._show_question()

    def _save_current_answer(self):
        if self.questions:
            self.answers[self.current_q] = self.ans_text.get("1.0", "end-1c").strip()

    def _submit_answer(self):
        self._save_current_answer()
        answer = self.answers[self.current_q]
        if not answer:
            self.feedback_var.set("Please write an answer before submitting.")
            return
        question = self.questions[self.current_q]
        self.feedback_var.set("Scoring your answer...")

        def worker():
            prompt = (f"Question: {question}\n\nCandidate answer: {answer}\n\n"
                      "Rate this answer from 1-10 and give 1-2 sentences of feedback. "
                      "Format: Score: X/10\nFeedback: <text>")
            try:
                resp = self.client.chat.completions.create(
                    model="gpt-4o-mini",
                    messages=[{"role": "user", "content": prompt}],
                    max_tokens=200
                )
                result = resp.choices[0].message.content.strip()
                m = re.search(r"Score:\s*(\d+)", result)
                score = int(m.group(1)) if m else None
                feedback_m = re.search(r"Feedback:\s*(.+)", result, re.S)
                feedback = feedback_m.group(1).strip() if feedback_m else result
                self.scores[self.current_q] = score
                self.after(0, lambda: self.feedback_var.set(
                    f"Score: {score}/10  —  {feedback}"
                ))
            except Exception as e:
                self.after(0, lambda: self.feedback_var.set(f"Scoring error: {e}"))
        threading.Thread(target=worker, daemon=True).start()

    def _show_report(self):
        answered = sum(1 for a in self.answers if a)
        scored = [s for s in self.scores if s is not None]
        avg = sum(scored) / len(scored) if scored else 0
        report = tk.Toplevel(self)
        report.title("📊 Interview Report Card")
        report.geometry("500x400")
        txt = tk.Text(report, wrap="word", font=("Helvetica", 11), padx=10, pady=10)
        txt.pack(fill="both", expand=True)
        txt.insert("end", f"Mock Interview Report\n{'='*40}\n\n")
        txt.insert("end", f"Questions answered: {answered}/{len(self.questions)}\n")
        txt.insert("end", f"Average score: {avg:.1f}/10\n\n")
        for i, (q, a, s) in enumerate(zip(self.questions, self.answers, self.scores)):
            score_str = f"{s}/10" if s is not None else "Not scored"
            txt.insert("end", f"Q{i+1} [{score_str}]: {q}\n")
            if a:
                txt.insert("end", f"  Your answer: {a[:200]}{'...' if len(a)>200 else ''}\n\n")
        txt.config(state="disabled")
        ttk.Button(report, text="Close", command=report.destroy).pack(pady=8)


class ChatGPTAssistant:
    def __init__(self,app):
        self.app = app
        self.recorder = AudioRecorder()
        self.streaming = False
        self.current_response = ""
        self.messages = [{"role": "system", "content": "You are a helpful interview assistant. Provide detailed technical answers and ask follow-up questions when appropriate."}]
        self.lock = threading.Lock()
        self.last_scroll_position = 0
        self.font_size = 12
        self.stream_thread = None 
        
        # ====== MODEL SETTINGS ======
        self.current_model = "gpt-4o"          # Default model (can be changed via UI)
        self.max_retries = 3                   # Number of retries for failed API calls
        
        # ====== OPTIMIZATION SETTINGS ======
        # These control how we send context to GPT while keeping FULL history locally
        self.optimization_mode = False         # Default: Full context; True = Speed/cheaper
        self.max_rounds_for_model = 4          # Kept for backwards compat (budget-based now)
        self.summary_message = None            # Rolling summary of older conversation
        self.summary_threshold_rounds = 4      # Start summarising after 4 rounds
        self.image_detail_level = "low"        # "low" = 85 tokens, "high" = 765+ tokens
        self._summary_in_progress = False      # Prevent concurrent summarisation
        self._pending_summary_thread = None    # Background summary thread
        self._system_msg_truncated = False     # Track if system msgs were truncated
        self._last_summary_rounds = 0          # Round count at last summary generation

        # v2 additions
        self.transcription_mode = "api"        # "api" | "local"
        self.auto_copy_enabled = False
        self._local_whisper_model = None       # lazy-loaded FasterWhisperModel

    def max_output_tokens_for_request(self) -> int:
        """Cap completion length by UI answer mode — faster replies and lower cost."""
        mode = getattr(self.app, "answer_mode", "default") if self.app else "default"
        return {
            "quick": 400,
            "default": 1000,
            "detailed": 2200,
            "code": 1800,
        }.get(mode, 1000)

    def get_model_for_question(self) -> str:
        """
        Smart model routing — use fast/cheap gpt-4o-mini for simple questions,
        reserve gpt-4o for complex ones. Saves ~300-500ms latency + 10x cost.

        behavioral / general  → gpt-4o-mini  (fast, cheap, great for STAR answers)
        coding / system_design → gpt-4o      (needs deep reasoning)
        technical              → gpt-4o-mini  (usually concept explanations)
        manual override        → self.current_model (user chose via UI button)
        """
        if not getattr(self.app, '_auto_route_model', True):
            return self.current_model  # user pinned a specific model
        qtype = getattr(self.app, '_last_question_type', 'general')
        if qtype in ('coding', 'system_design'):
            return 'gpt-4o'
        return 'gpt-4o-mini'  # behavioral, technical, general → fast model
    
    def _maybe_summarize_history(self):
        """
        Trigger a rolling background summary of the conversation history.

        ALWAYS runs regardless of optimization_mode — summarisation is free
        (gpt-4o-mini, ~800 tokens) and critical for keeping context on long chats.

        Rules:
        • Threshold: 4 completed Q&A rounds (8 user+assistant messages)
        • Refresh: every 3 new rounds after the last summary
        • NEVER blocks — always runs in a daemon thread
        • Called AFTER each completed response, never before
        """
        if self._summary_in_progress:
            return  # Already running in background

        user_assistant = [
            m for m in self.messages
            if isinstance(m, dict) and m.get("role") in ("user", "assistant")
        ]
        rounds = len(user_assistant) // 2

        START_THRESHOLD = 4   # Start summarising after 4 complete rounds
        REFRESH_EVERY   = 3   # Re-summarise every 3 new rounds

        if rounds < START_THRESHOLD:
            return  # Chat too short — no need yet

        # If we already have a recent summary, skip until enough new rounds have accumulated
        if self.summary_message:
            # Store the round count at last summary so we know when to refresh
            last_summarised = getattr(self, '_last_summary_rounds', 0)
            if rounds < last_summarised + REFRESH_EVERY:
                return  # Summary still fresh

        # Kick off background summarisation — NEVER blocks the answer
        self._summary_in_progress = True
        self._last_summary_rounds = rounds
        self._pending_summary_thread = threading.Thread(
            target=self._run_background_summary,
            args=(list(user_assistant),),   # snapshot — safe to use off-thread
            daemon=True
        )
        self._pending_summary_thread.start()
        print(f"📝 Background summary started ({rounds} rounds, non-blocking)")

    def _run_background_summary(self, user_assistant_msgs: list):
        """
        Runs summarization in background thread - DOES NOT BLOCK your response.
        Only updates summary_message when complete.
        """
        try:
            # 1) Build a plain-text transcript to summarize
            transcript_lines = []
            for m in user_assistant_msgs:
                role = m.get("role")
                role_label = "User" if role == "user" else "Assistant"

                content = m.get("content", "")
                if isinstance(content, list):
                    # multimodal: extract only text chunks
                    text_parts = []
                    for c in content:
                        if not isinstance(c, dict):
                            continue
                        if c.get("type") == "text":
                            text_parts.append(c.get("text", ""))
                    content = "\n".join(text_parts)
                else:
                    content = str(content)

                transcript_lines.append(f"{role_label}: {content}")

            transcript = "\n".join(transcript_lines)
            
            # Limit transcript length to avoid token limits on summary call
            if len(transcript) > 15000:
                transcript = transcript[:15000] + "\n... [truncated for summary]"

            # 2) Call a smaller/faster model to summarize
            resp = client.chat.completions.create(
                model="gpt-4o-mini",
                messages=[
                    {
                        "role": "system",
                        "content": (
                            "You are summarizing an ongoing job interview conversation. "
                            "Extract key points discussed: technical topics, candidate responses, "
                            "questions asked, skills mentioned, and any important context. "
                            "Be concise but preserve critical interview details."
                        ),
                    },
                    {"role": "user", "content": transcript},
                ],
                max_tokens=800  # Keep summary concise
            )
            summary_text = resp.choices[0].message.content.strip()

            # 3) Store as synthetic system message (thread-safe update)
            self.summary_message = {
                "role": "system",
                "content": (
                    "[Rolling Interview Summary — covers all previous Q&A]\n"
                    f"{summary_text}"
                ),
            }
            print(f"✅ Background summary complete: {len(summary_text)} chars — "
                  f"context preserved for next question")

        except Exception as e:
            print(f"❌ Background summary error: {e}")
        finally:
            self._summary_in_progress = False


    
    def _build_messages_for_model(self):
        """
        Build the message list sent to the API.

        Goals:
          1. ALWAYS preserve full interview context via rolling summary (never lose it)
          2. ALWAYS stay under a hard token budget so responses stay fast
          3. ALWAYS include your resume/JD system prompts (truncated if huge)
          4. Works identically regardless of optimization_mode toggle

        Token budget strategy:
          - Hard cap: TOKEN_BUDGET_TARGET tokens sent to API
          - System prompts eat up to SYSTEM_TOKEN_CAP tokens each
          - Summary gets ~800 tokens
          - Remaining budget filled greedily with most-recent Q&A pairs
          - "Optimized" mode = tighter budget (faster/cheaper)
          - "Full" mode = wider budget (more raw context, still capped)
        """
        # ── Split messages by role ──────────────────────────────────────────
        system_msgs = []
        other_msgs  = []
        for m in self.messages:
            if not isinstance(m, dict):
                continue
            role = m.get("role")
            if role == "system":
                system_msgs.append(m)
            elif role in ("user", "assistant"):
                other_msgs.append(m)

        total_qa = len(other_msgs)

        # ── Token budgets ───────────────────────────────────────────────────
        # Opt mode OFF = quality mode: wider window, still capped
        # Opt mode ON  = speed mode: tighter window
        if self.optimization_mode:
            TOKEN_BUDGET_TARGET = 8_000   # tight — fast responses
            SYSTEM_TOKEN_CAP    = 4_000   # chars per system message
        else:
            TOKEN_BUDGET_TARGET = 14_000  # generous — more raw context
            SYSTEM_TOKEN_CAP    = 8_000

        # ── 1. Build system section (resume, JD, base prompt) ──────────────
        optimized_system = []
        for sys_msg in system_msgs:
            content = sys_msg.get("content", "")
            if isinstance(content, str) and len(content) > SYSTEM_TOKEN_CAP:
                content = content[:SYSTEM_TOKEN_CAP] + "\n… [truncated — see summary for rest]"
                print(f"✂️ System msg truncated to {SYSTEM_TOKEN_CAP} chars")
            optimized_system.append({"role": "system", "content": content})

        # ── 2. Inject rolling summary (always, if available) ───────────────
        if self.summary_message:
            optimized_system.append(self.summary_message)

        # ── 3. Fill recent Q&A up to token budget ──────────────────────────
        # Start with the budget left after system tokens
        sys_tokens = estimate_tokens_for_messages(optimized_system)
        remaining  = TOKEN_BUDGET_TARGET - sys_tokens

        # Greedily pick messages from most-recent backward until budget full
        selected_qa = []
        for msg in reversed(other_msgs):
            msg_tokens = estimate_tokens_for_messages([msg])
            if remaining - msg_tokens < 0 and len(selected_qa) >= 4:
                break  # budget full — but always keep at least 2 rounds (4 msgs)
            # Optimise images in this message
            opt_msg = self._optimize_message_images(msg)
            # Strip images from messages that are not in the last 4 (2 rounds)
            if len(selected_qa) >= 4:
                opt_msg = self._strip_images_keep_text(opt_msg)
            # Truncate very long individual messages in tight mode
            if self.optimization_mode and len(selected_qa) >= 4:
                opt_msg = self._truncate_long_message(opt_msg, max_chars=2000)
            selected_qa.insert(0, opt_msg)
            remaining -= msg_tokens

        sent     = len(selected_qa)
        skipped  = total_qa - sent
        print(f"🔧 Context window: {sent}/{total_qa} Q&A msgs sent "
              f"({'opt' if self.optimization_mode else 'full'} mode, "
              f"~{TOKEN_BUDGET_TARGET - remaining:,} tokens, "
              f"{skipped} older msgs covered by summary)")

        # ── 4. Append meta-instructions ─────────────────────────────────────
        final_messages = optimized_system + selected_qa

        if self.app and hasattr(self.app, 'get_answer_mode_instruction'):
            mode_instruction = self.app.get_answer_mode_instruction()
            if mode_instruction:
                final_messages.append({"role": "system", "content": mode_instruction})

        qtype = getattr(self.app, '_last_question_type', 'general') if self.app else 'general'
        type_instruction = _QUESTION_TYPE_INSTRUCTIONS.get(qtype, "")
        if type_instruction:
            final_messages.append({"role": "system", "content": type_instruction})

        return final_messages
    
    def _truncate_long_message(self, msg: dict, max_chars: int = 2000) -> dict:
        """Truncate very long messages while preserving structure."""
        content = msg.get("content")
        
        if isinstance(content, str) and len(content) > max_chars:
            return {"role": msg.get("role"), "content": content[:max_chars] + "... [truncated]"}
        
        if isinstance(content, list):
            truncated_content = []
            for item in content:
                if isinstance(item, dict) and item.get("type") == "text":
                    text = item.get("text", "")
                    if len(text) > max_chars:
                        truncated_content.append({"type": "text", "text": text[:max_chars] + "... [truncated]"})
                    else:
                        truncated_content.append(item)
                else:
                    truncated_content.append(item)
            return {"role": msg.get("role"), "content": truncated_content}
        
        return msg
    
    def _strip_images_keep_text(self, msg: dict) -> dict:
        """
        Remove images from a message but keep all text content.
        Used for aggressive optimization of older messages.
        """
        content = msg.get("content")
        if not isinstance(content, list):
            return msg  # No images, return as-is
        
        # Extract only text items
        text_parts = []
        for item in content:
            if isinstance(item, dict) and item.get("type") == "text":
                text_parts.append(item.get("text", ""))
            elif isinstance(item, str):
                text_parts.append(item)
        
        if not text_parts:
            # If no text, add placeholder
            text_parts = ["[Image was here]"]
        
        combined_text = " ".join(text_parts)
        return {"role": msg.get("role"), "content": combined_text}
    
    def _optimize_message_images(self, msg: dict) -> dict:
        """
        Optimize images in a message by adding detail:low parameter.
        This reduces token usage from 765+ to 85 tokens per image.
        Original message is NOT modified - returns a new optimized copy.
        """
        if not self.optimization_mode:
            return msg
            
        content = msg.get("content")
        if not isinstance(content, list):
            return msg  # No images, return as-is
        
        # Create a copy with optimized images
        optimized_content = []
        for item in content:
            if isinstance(item, dict) and item.get("type") == "image_url":
                # Add detail parameter for optimization
                image_url_data = item.get("image_url", {})
                optimized_item = {
                    "type": "image_url",
                    "image_url": {
                        "url": image_url_data.get("url", ""),
                        "detail": self.image_detail_level  # "low" = 85 tokens
                    }
                }
                optimized_content.append(optimized_item)
            else:
                optimized_content.append(item)
        
        return {"role": msg.get("role"), "content": optimized_content}


            
    def cancel_streaming(self):
        self.streaming = False
        if self.stream_thread and self.stream_thread.is_alive():
            try:
                self.stream_thread.join(timeout=1)
            except:
                pass  # Avoid crash if join fails



    def load_document(self, file_path: str) -> tuple:
        """
        Load a single file of any type into the assistant's system context.
        Uses extract_text_from_file() for robust, textract-free extraction.
        Returns (success: bool, message: str).
        """
        base = os.path.basename(file_path)
        try:
            text = extract_text_from_file(file_path)
            if not text.strip():
                return False, f"⚠️ {base} appears empty or has no readable text."
            self.messages.append({
                "role": "system",
                "content": (
                    f"Attached document '{base}':\n"
                    f"{text[:50000]}"          # guard against huge files filling context
                )
            })
            size_hint = f"{len(text):,} chars"
            return True, f"📄 {base} loaded ({size_hint})."
        except RuntimeError as e:
            return False, f"❌ {str(e)}"
        except Exception as e:
            return False, f"❌ {base}: {str(e)}"

    def load_resume(self, file_path: str) -> tuple:
        """Legacy shim kept for hotkey / backward-compat. Delegates to load_document()."""
        return self.load_document(file_path)


    def _transcribe_local(self, filename: str) -> str:
        """Transcribe using local faster-whisper model (no API call)."""
        if not _FASTER_WHISPER_AVAILABLE:
            raise RuntimeError("faster-whisper not installed")
        if self._local_whisper_model is None:
            self._local_whisper_model = FasterWhisperModel("base", device="cpu", compute_type="int8")
        segments, _ = self._local_whisper_model.transcribe(filename, language="en")
        return " ".join(seg.text for seg in segments).strip()

    def transcribe_audio(self, filename, prompt: str | None = None):
        """Transcribe audio — routes to local or API based on transcription_mode."""
        if self.transcription_mode == "local" and _FASTER_WHISPER_AVAILABLE:
            try:
                return self._transcribe_local(filename)
            except Exception as e:
                print(f"⚠️ Local transcription failed, falling back to API: {e}")

        last_error = None
        for retry in range(self.max_retries):
            try:
                with open(filename, "rb") as audio_file:
                    transcription = client.audio.transcriptions.create(
                        model="whisper-1",
                        file=audio_file,
                        prompt=prompt or ""
                    )
                return transcription.text
            except Exception as e:
                last_error = e
                if retry < self.max_retries - 1:
                    wait_time = (retry + 1) * 2
                    print(f"⚠️ Transcription failed (attempt {retry+1}/{self.max_retries}): {e}")
                    time.sleep(wait_time)
                else:
                    return f"❌ Transcription error after {self.max_retries} retries: {str(last_error)}"

        return f"❌ Transcription error: {str(last_error)}"

    def diagnose_performance(self) -> dict:
        """
        Analyze current chat for performance issues.
        Returns diagnostic info about token usage and latency sources.
        """
        diag = {
            "total_messages": len(self.messages),
            "system_messages": 0,
            "user_messages": 0,
            "assistant_messages": 0,
            "images_count": 0,
            "estimated_total_tokens": 0,
            "estimated_system_tokens": 0,
            "estimated_conversation_tokens": 0,
            "estimated_image_tokens": 0,
            "optimization_mode": self.optimization_mode,
            "would_send_messages": 0,
            "issues": [],
            "recommendations": []
        }
        
        for msg in self.messages:
            if not isinstance(msg, dict):
                continue
            role = msg.get("role", "")
            content = msg.get("content", "")
            
            # Count by role
            if role == "system":
                diag["system_messages"] += 1
                tokens = len(str(content)) // 4
                diag["estimated_system_tokens"] += tokens
            elif role == "user":
                diag["user_messages"] += 1
            elif role == "assistant":
                diag["assistant_messages"] += 1
            
            # Count images
            if isinstance(content, list):
                for item in content:
                    if isinstance(item, dict) and item.get("type") == "image_url":
                        diag["images_count"] += 1
                        # Image tokens: 85 (low) to 765+ (high)
                        if self.optimization_mode:
                            diag["estimated_image_tokens"] += 85
                        else:
                            diag["estimated_image_tokens"] += 765
            
            # Estimate text tokens
            if isinstance(content, str):
                diag["estimated_conversation_tokens"] += len(content) // 4
            elif isinstance(content, list):
                for item in content:
                    if isinstance(item, dict) and item.get("type") == "text":
                        diag["estimated_conversation_tokens"] += len(item.get("text", "")) // 4
        
        # Single build — same payload Whisper/GPT optimization uses
        model_msgs = self._build_messages_for_model()
        diag["would_send_messages"] = len(model_msgs)
        
        # Calculate total estimated tokens
        diag["estimated_total_tokens"] = (
            diag["estimated_system_tokens"] + 
            diag["estimated_conversation_tokens"] + 
            diag["estimated_image_tokens"]
        )
        
        diag["will_send_tokens"] = estimate_tokens_for_messages(model_msgs, self.optimization_mode)
        diag["will_send_messages"] = len(model_msgs)
        
        # Identify issues
        if diag["estimated_system_tokens"] > 8000:
            diag["issues"].append(f"⚠️ System messages are very large ({diag['estimated_system_tokens']} tokens)")
            diag["recommendations"].append("System messages will be auto-truncated")
        
        if diag["images_count"] > 5:
            diag["issues"].append(f"⚠️ Many images in chat ({diag['images_count']})")
            diag["recommendations"].append("Old images will be stripped automatically")
        
        if diag["total_messages"] > 20 and not self.summary_message:
            diag["issues"].append(f"⚠️ Long conversation ({diag['total_messages']} msgs) without summary")
            diag["recommendations"].append("Click 'Force Summarize' to generate summary now")
        
        if diag["will_send_tokens"] > 15000:
            diag["issues"].append(f"🔴 High token count ({diag['will_send_tokens']:,} tokens will be sent)")
            diag["recommendations"].append("Consider starting a new chat for faster responses")
        
        return diag
    
    def print_performance_report(self):
        """Print a detailed performance report to console."""
        diag = self.diagnose_performance()
        
        print("\n" + "="*60)
        print("📊 PERFORMANCE DIAGNOSTIC REPORT")
        print("="*60)
        print(f"Total Messages: {diag['total_messages']}")
        print(f"  - System: {diag['system_messages']}")
        print(f"  - User: {diag['user_messages']}")
        print(f"  - Assistant: {diag['assistant_messages']}")
        print(f"  - Images: {diag['images_count']}")
        print()
        print(f"Estimated Tokens:")
        print(f"  - System: ~{diag['estimated_system_tokens']:,}")
        print(f"  - Conversation: ~{diag['estimated_conversation_tokens']:,}")
        print(f"  - Images: ~{diag['estimated_image_tokens']:,}")
        print(f"  - TOTAL: ~{diag['estimated_total_tokens']:,}")
        print()
        print(f"Optimization Mode: {'ON ⚡' if diag['optimization_mode'] else 'OFF 🐢'}")
        print(f"Messages sent to API: {diag['would_send_messages']}/{diag['total_messages']}")
        print(f"Has Summary: {'Yes ✅' if self.summary_message else 'No ❌'}")
        print()
        
        if diag["issues"]:
            print("⚠️ ISSUES FOUND:")
            for issue in diag["issues"]:
                print(f"  {issue}")
            print()
            print("💡 RECOMMENDATIONS:")
            for rec in diag["recommendations"]:
                print(f"  • {rec}")
        else:
            print("✅ No performance issues detected")
        
        print("="*60 + "\n")
        return diag

    def stream_gpt_response(self, text_widget, status_label, button, on_complete=None):
        self.cancel_streaming()  # 🔴 Cancel any ongoing output

        
        def run_stream():
            start_time = time.time()
            # NOTE: _maybe_summarize_history() is called in the finally block AFTER
            # streaming completes, so it never delays the answer. The summary produced
            # will be ready for the NEXT question.

            # Lock only for the brief message-list mutations; release before API call
            # so the lock is never held during network I/O (was previously held for seconds).
            with self.lock:
                self.current_response = ""
                self.streaming = True
                placeholder = {"role": "assistant", "content": ""}
                self.messages.append(placeholder)
                build_start = time.time()
                model_messages = self._build_messages_for_model()
                build_time = time.time() - build_start
                total_msgs = len(self.messages)

            try:
                sent_msgs = len(model_messages)
                estimated_tokens = estimate_tokens_for_messages(model_messages, self.optimization_mode)
                print(f"📊 Performance: {sent_msgs}/{total_msgs} msgs, ~{estimated_tokens:,} tokens, model: {self.current_model}, build: {build_time*1000:.0f}ms")

                # Smart model routing: pick fastest model for this question type
                chosen_model = self.get_model_for_question()
                print(f"🤖 Model: {chosen_model} (qtype={getattr(self.app, '_last_question_type', 'general')})")

                api_start = time.time()
                stream = None
                last_error = None

                for retry in range(self.max_retries):
                    try:
                        stream = client.chat.completions.create(
                            model=chosen_model,
                            messages=model_messages,
                            stream=True,
                            max_tokens=self.max_output_tokens_for_request(),
                        )
                        break
                    except Exception as e:
                        last_error = e
                        if retry < self.max_retries - 1:
                            # Short sleeps only — long waits make the app feel frozen
                            wait_time = 0.5 * (retry + 1)   # 0.5s, 1.0s (was 2s, 4s)
                            print(f"⚠️ API call failed (attempt {retry+1}/{self.max_retries}): {e}")
                            status_label.config(text=f"⚠️ Retrying ({retry+1}/{self.max_retries})…")
                            time.sleep(wait_time)
                        else:
                            raise last_error

                if stream is None:
                    raise Exception(f"Failed after {self.max_retries} retries: {last_error}")

                first_token_time = None
                buffer = ""
                last_update = time.time()

                def _insert_answer_header():
                    first_visible = text_widget.index("@0,0")
                    text_widget.config(state=tk.NORMAL)
                    text_widget.insert(tk.END, "------------------\nANSWER: ")
                    text_widget.config(state=tk.DISABLED)
                    text_widget.see(first_visible)
                text_widget.after(0, _insert_answer_header)

                output_chars = 0
                for chunk in stream:
                    if not self.streaming:
                        break
                    delta = chunk.choices[0].delta.content if chunk.choices[0].delta else ""
                    if delta:
                        if first_token_time is None:
                            first_token_time = time.time()
                            print(f"⏱️ Time to first token: {(first_token_time - api_start)*1000:.0f}ms")

                        buffer += delta
                        output_chars += len(delta)
                        self.current_response += delta
                        placeholder["content"] = self.current_response

                        if time.time() - last_update > 0.05 or len(buffer) > 20:
                            self.update_text_widget(text_widget, buffer)
                            buffer = ""
                            last_update = time.time()

                if buffer:
                    self.update_text_widget(text_widget, buffer)

                total_time = time.time() - start_time
                print(f"⏱️ Total response time: {total_time:.1f}s | 📈 ~{estimated_tokens:,} in, ~{output_chars // 4:,} out tokens")

            except Exception as e:
                placeholder["content"] = f"❌ GPT Error: {str(e)}"
                self.update_text_widget(text_widget, f"\n{placeholder['content']}\n")

            finally:
                self.streaming = False
                button.config(state=tk.NORMAL)
                status_label.config(text="✅ Ready")
                # Post-process the answer:
                #   150ms → math formatter  (LaTeX → readable Unicode, boxes)
                #   350ms → code renderer   (``` blocks → embedded no-wrap widgets)
                # Both run non-blocking via text_widget.after() so they never
                # stall the UI thread.
                _has_math = '\\' in self.current_response or '$' in self.current_response
                _has_code = '```' in self.current_response
                self._reformat_math_in_widget(text_widget)
                if _has_code:
                    self._render_code_blocks_in_widget(text_widget)
                if self.auto_copy_enabled and self.current_response and not _has_math:
                    try:
                        pyperclip.copy(self.current_response)
                    except Exception:
                        pass
                if self.app:
                    self.app.chat_manager.save_current_session(self.messages)
                    if on_complete is not None:
                        self.app.after(600, on_complete)
                # Summarise after the response — never before — so the next
                # question benefits from an up-to-date summary with zero delay.
                self._maybe_summarize_history()



        self.stream_thread = threading.Thread(target=run_stream, daemon=True)
        self.stream_thread.start()





    def update_text_widget(self, text_widget, new_text_part: str):
        # Called from background streaming thread — schedule on main thread via after().
        #
        # Key insight: yview_moveto(fraction) uses a PROPORTIONAL position.
        # As text is appended and the document grows longer, the same fraction points
        # to a lower and lower line — so every insert silently scrolls the viewport down.
        #
        # Correct fix: save "@0,0" — the absolute text index (e.g. "42.0") of the
        # character visible at the top-left pixel. Line numbers are absolute; line 42
        # stays line 42 no matter how many lines are appended after it.
        # After the insert, see(saved_index) brings that exact line back into view.
        def _do_insert():
            # Absolute index of the top-left visible character
            first_visible = text_widget.index("@0,0")
            text_widget.config(state=tk.NORMAL)
            text_widget.insert(tk.END, new_text_part)
            text_widget.config(state=tk.DISABLED)
            # Restore that exact line to the top of the viewport
            text_widget.see(first_visible)

        text_widget.after(0, _do_insert)


    def _reformat_math_in_widget(self, text_widget):
        """After streaming finishes, find the last ANSWER block in the response box
        and replace raw LaTeX with human-readable Unicode / boxed formulas."""
        raw_response = self.current_response
        if not raw_response:
            return
        # Quick check: only bother if there's actual math notation present
        if '\\' not in raw_response and '$' not in raw_response:
            return

        def _do_reformat():
            try:
                full = text_widget.get("1.0", tk.END)
                marker = "ANSWER: "
                last_pos = full.rfind(marker)
                if last_pos == -1:
                    return
                answer_start_char = last_pos + len(marker)
                # Convert character offset → Tkinter line.col index
                prefix = full[:answer_start_char]
                line_no = prefix.count('\n') + 1
                col_no = len(prefix.split('\n')[-1])
                tk_idx = f"{line_no}.{col_no}"

                raw = text_widget.get(tk_idx, tk.END).rstrip('\n')
                formatted = format_math_for_display(raw)
                if formatted == raw:
                    return  # nothing changed, skip widget update

                first_visible = text_widget.index("@0,0")
                text_widget.config(state=tk.NORMAL)
                text_widget.delete(tk_idx, tk.END)
                text_widget.insert(tk.END, formatted)
                text_widget.config(state=tk.DISABLED)
                text_widget.see(first_visible)

                # Update the stored response to the formatted version so
                # auto-copy and bookmarks get the readable text
                self.current_response = formatted
                if self.auto_copy_enabled:
                    try:
                        import pyperclip as _pc
                        _pc.copy(formatted)
                    except Exception:
                        pass
            except Exception as exc:
                print(f"[math reformat] error: {exc}")

        # Small delay so the final streaming buffer flush has settled
        text_widget.after(150, _do_reformat)


    def highlight_code(self, text_widget):
        """Legacy — code blocks are now rendered as embedded widgets by
        _render_code_blocks_in_widget(). Kept as no-op so old call sites
        don't crash."""
        pass

    # ------------------------------------------------------------------
    # Code-block renderer — replaces ``` fences with embedded Text widgets
    # ------------------------------------------------------------------
    def _render_code_blocks_in_widget(self, text_widget):
        """
        After streaming finishes, scan the last answer for ``` code blocks
        and replace each one with an embedded scrollable Text widget that has:

          • wrap=NONE  → ASCII diagrams never break across lines
          • Horizontal scrollbar → scroll wide diagrams without distortion
          • Menlo/Courier monospace font → every character same width
          • Dark VS Code-style chrome (header bar + frame)

        This runs 350 ms after streaming so it chains after the math formatter.
        """
        raw_response = self.current_response
        if not raw_response or '```' not in raw_response:
            return  # nothing to render

        # Preferred monospace fonts for ASCII art (macOS-first)
        CODE_FONTS = [('Menlo', 11), ('Courier New', 11), ('Courier', 11)]

        def _pick_font():
            import tkinter.font as tkfont
            available = tkfont.families()
            for name, size in CODE_FONTS:
                if name in available:
                    return (name, size)
            return ('Courier', 11)

        def _char_offset_to_tk(full_text: str, char_offset: int) -> str:
            """Convert an absolute character offset in full_text → 'line.col' index."""
            prefix = full_text[:char_offset]
            line_no = prefix.count('\n') + 1
            col_no  = len(prefix) - prefix.rfind('\n') - 1
            return f"{line_no}.{col_no}"

        def _do():
            try:
                full = text_widget.get("1.0", tk.END)

                # Locate the start of the last answer
                marker     = "ANSWER: "
                last_pos   = full.rfind(marker)
                if last_pos == -1:
                    return
                ans_start = last_pos + len(marker)

                # Find all ``` blocks inside that answer
                pattern = re.compile(r'```(\w*)\n?(.*?)```', re.DOTALL)
                blocks  = list(pattern.finditer(full, pos=ans_start))
                if not blocks:
                    return

                code_font    = _pick_font()
                first_visible = text_widget.index("@0,0")
                text_widget.config(state=tk.NORMAL)

                # Process REVERSE so earlier block's tk-indices stay valid
                for match in reversed(blocks):
                    lang         = (match.group(1) or "").strip()
                    code_content = match.group(2).rstrip('\n')

                    abs_start = match.start()
                    abs_end   = match.end()

                    tk_start = _char_offset_to_tk(full, abs_start)
                    tk_end   = _char_offset_to_tk(full, abs_end)

                    # ── Delete the raw ``` block ──────────────────────────
                    text_widget.delete(tk_start, tk_end)

                    # ── Build the embedded code widget ────────────────────
                    lines   = code_content.split('\n')
                    max_len = max((len(l) for l in lines), default=40)
                    h = min(max(len(lines), 2), 30)          # 2-30 lines
                    w = min(max(max_len + 4, 44), 100)       # 44-100 chars

                    # Outer container
                    outer = tk.Frame(
                        text_widget,
                        bg='#0d1117',
                        highlightthickness=1,
                        highlightbackground='#30363d'
                    )

                    # Header bar  ── "PYTHON" / "TEXT" / "CODE"
                    lang_label = lang.upper() if lang else "CODE"
                    hdr = tk.Label(
                        outer,
                        text=f"  {lang_label}  ",
                        bg='#161b22', fg='#8b949e',
                        font=(code_font[0], 9, 'bold'),
                        anchor='w', padx=6, pady=4
                    )
                    hdr.pack(fill='x', side='top')

                    tk.Frame(outer, bg='#30363d', height=1).pack(fill='x', side='top')

                    # Scrollbar frame
                    inner = tk.Frame(outer, bg='#0d1117')
                    inner.pack(fill='both', expand=True, side='top')

                    xscroll = tk.Scrollbar(inner, orient='horizontal', bg='#161b22')
                    yscroll = tk.Scrollbar(inner, orient='vertical',   bg='#161b22')

                    ct = tk.Text(
                        inner,
                        wrap=tk.NONE,           # ← KEY: never wrap lines
                        font=code_font,
                        bg='#0d1117', fg='#e6edf3',
                        selectbackground='#264f78',
                        insertbackground='white',
                        width=w, height=h,
                        bd=0, padx=12, pady=8,
                        xscrollcommand=xscroll.set,
                        yscrollcommand=yscroll.set,
                        cursor='arrow',
                    )
                    xscroll.config(command=ct.xview)
                    yscroll.config(command=ct.yview)

                    # Only show y-scrollbar if content taller than view
                    if len(lines) > h:
                        yscroll.pack(side='right', fill='y')
                    xscroll.pack(side='bottom', fill='x')
                    ct.pack(side='left', fill='both', expand=True)

                    ct.insert('1.0', code_content)
                    ct.config(state='disabled')

                    # ── Insert into main widget ───────────────────────────
                    # Order of insertions at tk_start (each pushes prior ones right):
                    #   Step 1: insert trailing \n  → result at tk_start: \n
                    #   Step 2: window_create       → [window] \n
                    #   Step 3: insert leading  \n  → \n [window] \n
                    text_widget.insert(tk_start, '\n')
                    text_widget.window_create(tk_start, window=outer, padx=2, pady=4)
                    text_widget.insert(tk_start, '\n')

                text_widget.config(state=tk.DISABLED)
                text_widget.see(first_visible)

            except Exception as exc:
                print(f"[code render] {exc}")
                import traceback; traceback.print_exc()
                try:
                    text_widget.config(state=tk.DISABLED)
                except Exception:
                    pass

        # 350 ms — after the math formatter (150 ms) has already run
        text_widget.after(350, _do)


class Application(tk.Tk):
    def __init__(self):
        super().__init__()
        self.title("INTERVIEW ASSISTANT")

        # --- Load UI prefs first
        self.ui_prefs = UIPreferences.load()

        # Use saved geometry if present (falls back to your hardcoded one)
        self.geometry(self.ui_prefs.get("geometry", "643x967+-644+25"))

        self.toggle_lock = threading.Lock()
        self._last_persisted_hash = None

        self.is_processing_audio = False
        self.live_transcription_running = False
        self.latest_live_question = ""   # last incremental text from Whisper
        self.live_question_index = None  # index of the "Live Question" line in the Text widget
        
        # Answer Quality Mode: "default", "quick", "detailed", "code"
        self.answer_mode = "default"  # Default mode - normal GPT behavior
        
        # Multi-Model Support: Available models with descriptions
        self.available_models = {
            "gpt-4o": "🧠 GPT-4o (Best quality, slower, higher cost)",
            "gpt-4o-mini": "⚡ GPT-4o-mini (Fast, cheaper, good for simple Q&A)",
            "gpt-4-turbo": "🚀 GPT-4-Turbo (Balance of speed & quality)"
        }
        self.current_model = "gpt-4o"  # Default model
        
        # Connection status
        self.api_connected = False
        
        # Audio level tracking
        self.current_audio_level = 0
        
        # UI Mode: "modern" or "classic"
        self.ui_mode = self.ui_prefs.get("ui_mode", "modern")

        # v2 feature flags
        self.stealth_mode = False
        self.vad_mode = False
        self._vad_running = False
        self._last_question_type = "general"
        self._auto_route_model = True   # True = smart routing; False = always use user-picked model

        self.assistant = ChatGPTAssistant(app=self)
        self.prompt_manager = PromptManager()
        self.chat_manager = SQLiteChatHistoryManager()
        

        # If user saved a preferred font size, use it before building widgets
        if "response_font_size" in self.ui_prefs:
            self.assistant.font_size = int(self.ui_prefs["response_font_size"])

        self.setup_ui()
        self.load_chat_tabs()

        # Auto-load autosave session (unchanged)
        # Auto-load autosave session (unchanged)
        # Auto-load autosave session (safer)
        if self.chat_manager.sessions and self.chat_manager.sessions[0].get("title") == "AutoSave - Last Session":
            msgs = self.chat_manager.sessions[0].get("messages", [])
            if isinstance(msgs, list):
                self.assistant.messages = msgs
                self.display_chat_history()


            self.status.config(text="🕑 Resumed from last auto-save session")

        # Bind paste to input_entry directly (not bind_all) to prevent double-paste
        self.input_entry.bind("<Command-v>", self.handle_paste)
        self.input_entry.bind("<Control-v>", self.handle_paste)  # For non-Mac keyboards
        self.sidebar_visible = True
        self.current_tab = -1
        self.current_subtab = -1
        self.always_on_top = False

        # F1 already prints geometry
        self.bind("<F1>", lambda e: print("Window geometry:", self.geometry()))
        # 💾 New: F2 save, F3 apply
        self.bind("<F2>", lambda e: self.save_ui_prefs())
        self.bind("<F3>", lambda e: self.apply_ui_prefs())
        
        # 🚀 Quick Setup shortcut: Cmd+Shift+S (or Ctrl+Shift+S on Windows)
        self.bind("<Command-Shift-s>", lambda e: self.open_quick_setup())
        self.bind("<Control-Shift-s>", lambda e: self.open_quick_setup())
        
        # 🔖 Bookmark shortcut: Cmd+B or F4
        self.bind("<Command-b>", lambda e: self.add_bookmark_at_cursor())
        self.bind("<Control-b>", lambda e: self.add_bookmark_at_cursor())
        self.bind("<F4>", lambda e: self.add_bookmark_at_cursor())
        # Go to next bookmark one-by-one (cycle): F5
        self.bind("<F5>", lambda e: self.go_to_next_bookmark())

        # 📌 Default interview: Cmd+Shift+I feeds all default subtab instructions at once
        self.bind("<Command-Shift-i>", lambda e: self.apply_default_interview_instructions())
        self.bind("<Control-Shift-i>", lambda e: self.apply_default_interview_instructions())

        # Chat scroll: PgUp = top, PgDn = end, Up/Down = paragraph (when focus not in input)
        self.bind("<Prior>", lambda e: self._scroll_chat_to_top())
        self.bind("<Next>", lambda e: self._scroll_chat_to_bottom())
        self.bind("<Up>", lambda e: self._on_up_key(e))
        self.bind("<Down>", lambda e: self._on_down_key(e))

        # Apply sash (split) after widgets exist
        self.after(0, self.apply_ui_prefs)

        # Load tabs after UI setup
        self.load_tabs()
        # Ensure we start within limits
        self.after(0, lambda: self.auto_prune_chats(max_chats=10))

    def update_live_question_in_ui(self, text: str):
        """
        Safely update the 'Live Question: ...' line in the response_box.
        This is called from a background thread via .after().
        """
        if not self.live_transcription_running:
            return

        try:
            self.response_box.config(state=tk.NORMAL)
            
            # Find "Live Question:" and replace everything after it on that line
            pos = self.response_box.search("Live Question:", "1.0", tk.END)
            if pos:
                # Get line number and delete from "Live Question:" to end of line
                line_num = pos.split('.')[0]
                self.response_box.delete(pos, f"{line_num}.end")
                # Insert the updated text
                self.response_box.insert(pos, f"Live Question: {text}")
            
            self.response_box.config(state=tk.DISABLED)
            # Only follow live transcription updates if already at the bottom
            if self.response_box.yview()[1] >= 0.99:
                self.response_box.see(tk.END)
        except Exception as e:
            print(f"Live UI update error: {e}")



    def live_transcription_loop(self):
        """
        Runs in a background thread while recording.
        Sends only the last LIVE_WHISPER_TAIL_SECONDS of audio to Whisper on a paced
        interval — final transcription on stop still uses the full WAV.

        (Previously every poll sent *all* recorded audio, which ballooned uploads/API time.)
        """
        last_text = ""
        tail_samples = max(int(SAMPLE_RATE * LIVE_WHISPER_TAIL_SECONDS), SAMPLE_RATE // 2)
        last_api_time = 0.0

        while self.live_transcription_running and self.assistant.recorder.is_recording:
            now = time.time()
            if now - last_api_time < LIVE_WHISPER_MIN_INTERVAL_SEC:
                time.sleep(0.12)
                continue

            aud = self.assistant.recorder.get_tail_snapshot(tail_samples)
            if aud is None:
                time.sleep(0.25)
                continue

            if len(aud) < SAMPLE_RATE // 2:
                time.sleep(0.2)
                continue

            whisper_prompt = last_text[:200] if last_text else ""
            text = ""
            temp_path = None
            try:
                fd, temp_path = tempfile.mkstemp(suffix=".wav")
                os.close(fd)
                with wave.open(temp_path, 'wb') as wf:
                    wf.setnchannels(CHANNELS)
                    wf.setsampwidth(2)
                    wf.setframerate(SAMPLE_RATE)
                    wf.writeframes(aud.tobytes())

                result = self.assistant.transcribe_audio(temp_path, prompt=whisper_prompt)
                text = (result or "").strip()
                if text.startswith("❌"):
                    text = ""
                last_api_time = time.time()
            except Exception as e:
                print(f"❌ Live transcription error: {e}")
                last_api_time = time.time()
                time.sleep(0.8)
            finally:
                if temp_path:
                    try:
                        os.remove(temp_path)
                    except Exception:
                        pass

            if text and text != last_text:
                last_text = text
                self.latest_live_question = text
                self.after(0, lambda t=text: self.update_live_question_in_ui(t))

            time.sleep(0.08)




    def auto_prune_chats(self, max_chats=10):
        """
        If there are more than `max_chats` real chats (excluding AutoSave),
        automatically delete all but the most recent real chat.
        Always keep "AutoSave - Last Session" if present.
        """
        AUTO_TITLE = "AutoSave - Last Session"

        # Split out real chats vs AutoSave
        real_indices = [i for i, s in enumerate(self.chat_manager.sessions)
                        if s.get("title") != AUTO_TITLE]
        real_count = len(real_indices)

        if real_count <= max_chats:
            return  # nothing to do

        # Keep only the MOST RECENT real chat (the last appended),
        # plus AutoSave if it exists
        keep_real_idx = real_indices[-1]  # most recent real chat by your append order
        new_sessions = []
        kept_title = None

        for i, s in enumerate(self.chat_manager.sessions):
            title = s.get("title", "Untitled")
            if i == keep_real_idx or title == AUTO_TITLE:
                new_sessions.append(s)
                if i == keep_real_idx:
                    kept_title = title

        removed_count = len(self.chat_manager.sessions) - len(new_sessions)
        if removed_count > 0:
            self.chat_manager.sessions = new_sessions
            self.chat_manager.save()
            self.load_chat_tabs()

            # Reselect the kept real chat if it exists, otherwise select AutoSave
            reselect_index = 0
            for i, s in enumerate(self.chat_manager.sessions):
                if s.get("title", "") == kept_title:
                    reselect_index = i
                    break
            kept_item_id = f"chat_{reselect_index}"
            if self.chat_tabs.exists(kept_item_id):
                self.chat_tabs.selection_set(kept_item_id)
                self.chat_tabs.see(kept_item_id)

            self.status.config(
                text=f"🧹 Auto-pruned {removed_count} chat(s). Kept recent chat{f' “{kept_title}”' if kept_title else ''} and “{AUTO_TITLE}”."
            )


    def _get_tree_open_state(self, tree: ttk.Treeview):
        """Return a list of item IDs that are expanded (open=True) in the given Treeview."""
        open_iids = []

        def walk(iid):
            try:
                if tree.item(iid, 'open'):
                    open_iids.append(iid)
            except Exception:
                pass
            for child in tree.get_children(iid):
                walk(child)

        # top-level nodes
        for root in tree.get_children(''):
            walk(root)

        return open_iids


    def _apply_tree_open_state(self, tree: ttk.Treeview, open_iids):
        """Expand items whose IDs are in open_iids (ignore any that don't exist)."""
        if not open_iids:
            return
        def do_apply():
            for iid in open_iids:
                if tree.exists(iid):
                    tree.item(iid, open=True)
        # Apply now and once more shortly after, in case the tree was just rebuilt
        self.after(0, do_apply)
        self.after(120, do_apply)

    def _generate_session_title(self):
        """Create a nice title for the current chat, reusing your start_new_chat logic."""
        timestamp = time.strftime("%Y-%m-%d %H:%M:%S")

        # Look for attached documents in system messages (supports both old and new format)
        doc_names = []
        for msg in self.assistant.messages:
            if not isinstance(msg, dict):
                continue
            if msg.get("role") != "system":
                continue
            content = msg.get("content", "")
            if not isinstance(content, str):
                continue
            # New format: "Attached document 'filename':\n..."
            m = re.match(r"Attached document '(.+?)':", content)
            if m:
                doc_names.append(os.path.splitext(os.path.basename(m.group(1)))[0])
                continue
            # Legacy format: "Use this resume content to contextualize answers (from file: name)"
            m2 = re.search(r'from file:\s*(.+?)\)', content)
            if m2:
                doc_names.append(os.path.splitext(os.path.basename(m2.group(1).strip()))[0])

        if doc_names:
            label = ", ".join(doc_names[:3])          # show up to 3 filenames
            if len(doc_names) > 3:
                label += f" +{len(doc_names) - 3} more"
            return f"{label} - {timestamp}"
        return timestamp


    def _persist_working_chat_if_needed(self):
        """
        If current working chat has user content and is not already saved as a real session
        (i.e., beyond 'AutoSave - Last Session'), save it now as a new session.
        """
        msgs = self.assistant.messages or []
        # Must have at least one user message to be meaningful
        if not any(m.get("role") == "user" for m in msgs):
            return

        # Fast O(1) hash: message count + fingerprint of last message content.
        # Avoids serialising the entire history on every chat-switch.
        try:
            last_msg = msgs[-1] if msgs else {}
            last_content = last_msg.get("content", "")
            if isinstance(last_content, list):
                last_content = str(last_content)
            cur_hash = (len(msgs), hash(last_content[:300]))
            if self._last_persisted_hash is not None and self._last_persisted_hash == cur_hash:
                return
        except Exception:
            cur_hash = None

        # Avoid duplicates: if an identical non-AutoSave session already exists, skip
        for s in self.chat_manager.sessions:
            if s.get("title") != "AutoSave - Last Session" and s.get("messages") == msgs:
                # already saved
                self._last_persisted_hash = cur_hash
                return

        # Save as a proper session with a generated title
        title = self._generate_session_title()
        self.chat_manager.add_session(title, msgs.copy())
        self.chat_manager.save()
        self.load_chat_tabs()
        # 🔁 Auto-prune when over the limit
        self.auto_prune_chats(max_chats=10)


        self.status.config(text=f"💾 Saved current chat as: {title}")
        self._last_persisted_hash = cur_hash


    def load_tabs(self):
        self.tab_tree.delete(*self.tab_tree.get_children())

        for i in range(self.prompt_manager.get_tab_count()):
            tab_id = self.tab_tree.insert("", "end", text=self.prompt_manager.get_tab_name(i), iid=f"tab_{i}")
            for j in range(self.prompt_manager.get_subtab_count(i)):
                subtab_id = self.tab_tree.insert(tab_id, "end", text=self.prompt_manager.get_subtab_name(i, j), iid=f"sub_{i}_{j}")

        # # Automatically select first subtab if available
        # if self.prompt_manager.get_tab_count() > 0 and self.prompt_manager.get_subtab_count(0) > 0:
        #     self.tab_tree.selection_set(f"sub_0_0")
        #     self.after(100, lambda: self.on_tab_select(None))  # Trigger selection logic

    def save_ui_prefs(self):
        """Save current window geometry, split position, font size, and tabs/subtabs open state."""
        try:
            sash = self.paned.sashpos(0)
        except Exception:
            sash = None
        
        # Save sidebar internal split (tabs vs chats)
        # sidebar_paned is tk.PanedWindow — use sash_coord(), not sashpos()
        try:
            sidebar_sash = self.sidebar_paned.sash_coord(0)[1]  # (x, y) → take y
        except Exception:
            sidebar_sash = None

        # Merge with existing prefs so we keep default_interview_subtabs, ui_mode, etc.
        prefs = UIPreferences.load()
        prefs["geometry"] = self.geometry()
        prefs["paned_sash"] = sash
        prefs["sidebar_sash"] = sidebar_sash
        prefs["response_font_size"] = int(self.assistant.font_size)
        prefs["tab_tree_open"] = self._get_tree_open_state(self.tab_tree)
        UIPreferences.save(prefs)
        self.status.config(text="💾 Saved UI defaults (geometry, splits, font, dropdowns).")
        print("Saved UI Prefs:", prefs)


    def apply_ui_prefs(self, *_):
        """Apply saved defaults (geometry, split, font size, tabs/subtabs open state)."""
        prefs = self.ui_prefs = UIPreferences.load()

        # Geometry
        if "geometry" in prefs:
            try:
                self.geometry(prefs["geometry"])
            except Exception as e:
                print("Geometry apply error:", e)

        # Font
        if "response_font_size" in prefs:
            try:
                self.assistant.font_size = int(prefs["response_font_size"])
                self.response_box.config(font=('Consolas', self.assistant.font_size))
            except Exception as e:
                print("Font apply error:", e)

        # Main split (sidebar vs main content)
        if "paned_sash" in prefs and prefs["paned_sash"] is not None:
            try:
                self.paned.sashpos(0, int(prefs["paned_sash"]))
            except Exception as e:
                print("Sash apply error, retrying...", e)
                self.after(50, lambda: self.paned.sashpos(0, int(prefs["paned_sash"])))

        # Sidebar internal split (tabs vs chats)
        # Uses tk.PanedWindow.sash_place(index, x, y) — clamp so chat section
        # always gets at least 150px (can never be hidden again).
        if "sidebar_sash" in prefs and prefs["sidebar_sash"] is not None:
            def apply_sidebar_sash():
                try:
                    sidebar_h = self.sidebar_paned.winfo_height()
                    saved_y = int(prefs["sidebar_sash"])
                    # Clamp: leave at least 150px for chat section at the bottom
                    max_y = max(120, sidebar_h - 150)
                    sash_y = max(120, min(saved_y, max_y))
                    self.sidebar_paned.sash_place(0, 0, sash_y)
                except Exception as e:
                    print("Sidebar sash apply error:", e)
            self.after(200, apply_sidebar_sash)  # give time for widget to render

        # NEW: tabs/subtabs expanded state
        if "tab_tree_open" in prefs:
            self._apply_tree_open_state(self.tab_tree, prefs["tab_tree_open"])

        self.status.config(text="✅ Applied UI defaults.")


        
    def on_chat_tab_select(self, event):
        selected = self.chat_tabs.selection()
        if not selected:
            return

        # ✅ Ensure the working chat is saved before switching away
        self._persist_working_chat_if_needed()

        tab_id = selected[0]
        if tab_id.startswith("chat_"):
            index = int(tab_id.split("_")[1])
            self._current_chat_index = index          # track which session is active
            self.assistant.messages = self.chat_manager.get_session(index)
            self.display_chat_history()               # restores bookmarks for this session
            self.status.config(text=f"📂 Loaded chat: {self.chat_manager.get_titles()[index]}")

    def add_new_tab(self):
        name = simpledialog.askstring("New Tab", "Enter tab name:")
        if name:
            tab_index = self.prompt_manager.add_tab(name)
            self.tab_tree.insert("", "end", text=name, iid=f"tab_{tab_index}")
            self.add_subtab_btn.config(state=tk.NORMAL)

        
    def toggle_input_mode(self):
        if self.assistant.recorder.input_mode == "internal":
            self.assistant.recorder.input_mode = "external"
            self.toggle_input_btn.config(text="🎧 Mic")
            self.status.config(text="🎧 Switched to External Microphone")
        else:
            self.assistant.recorder.input_mode = "internal"
            self.toggle_input_btn.config(text="🔈 BlackHole")
            self.status.config(text="🔈 Switched to Internal Audio (BlackHole)")

    # ========== CONNECTION STATUS ==========
    def check_api_connection(self):
        """
        Connection pre-warm: send a tiny real chat completion on startup.
        This establishes the HTTP/2 connection & TLS session so the FIRST real
        question gets the same latency as subsequent ones (saves ~200-400ms).
        """
        def _warm():
            try:
                # Send a 1-token completion — establishes TCP+TLS connection in the pool
                client.chat.completions.create(
                    model="gpt-4o-mini",
                    messages=[{"role": "user", "content": "hi"}],
                    max_tokens=1,
                    stream=False,
                )
                self.api_connected = True
                self.after(0, lambda: self.connection_label.config(text="🟢 Connected"))
                print("✅ Connection pre-warmed — first question will be fast")
            except Exception as e:
                self.api_connected = False
                self.after(0, lambda: self.connection_label.config(text="🔴 Offline"))
                print(f"API pre-warm failed: {e}")

        threading.Thread(target=_warm, daemon=True).start()
        # Re-check every 90 seconds (less aggressive — connection stays warm)
        self.after(90000, self.check_api_connection)

    # ========== MULTI-MODEL SUPPORT ==========
    def toggle_model(self):
        """Cycle through available models with description."""
        models = list(self.available_models.keys())
        model_labels = {
            "gpt-4o": "🧠 4o",
            "gpt-4o-mini": "⚡ Mini",
            "gpt-4-turbo": "🚀 Turbo"
        }
        
        current_idx = models.index(self.current_model)
        next_idx = (current_idx + 1) % len(models)
        self.current_model = models[next_idx]
        
        # Update button and show description
        self.model_btn.config(text=model_labels[self.current_model])
        self.status.config(text=f"🔄 Model: {self.available_models[self.current_model]}")
        
        # Update the assistant's model
        self.assistant.current_model = self.current_model

    # ========== UI MODE TOGGLE ==========
    def toggle_ui_mode(self):
        """Toggle between modern and classic UI mode (requires restart)."""
        if self.ui_mode == "modern":
            self.ui_mode = "classic"
            msg = "🎨 Classic UI mode selected. Restart app to apply."
        else:
            self.ui_mode = "modern"
            msg = "🎨 Modern UI mode selected. Restart app to apply."
        
        # Save preference
        self.ui_prefs["ui_mode"] = self.ui_mode
        UIPreferences.save(self.ui_prefs)
        self.status.config(text=msg)
        
        # Ask if user wants to restart now
        if messagebox.askyesno("UI Mode Changed", f"{msg}\n\nRestart now?"):
            self._restart_app()
    
    def _restart_app(self):
        """Restart the application."""
        import subprocess
        try:
            python = sys.executable
            script = os.path.abspath(sys.argv[0])
            subprocess.Popen([python, script])
            self.destroy()
            os._exit(0)
        except Exception as e:
            self.status.config(text=f"❌ Restart failed: {e}")

    # ========== AUDIO LEVEL INDICATOR ==========
    def update_audio_level(self):
        """Update the audio level indicator during recording."""
        if not self.assistant.recorder.is_recording:
            self.audio_level_bar['value'] = 0
            self.audio_level_label.config(text="--")
            return
        
        # Get current audio level from recorder (only last 0.1s needed for RMS)
        snapshot = self.assistant.recorder.get_tail_snapshot(1600)
        if snapshot is not None and len(snapshot) > 0:
            rms = np.sqrt(np.mean(snapshot**2))
            # Convert to percentage (normalize to reasonable range)
            level = min(100, int(rms / 300 * 100))
            self.current_audio_level = level
            self.audio_level_bar['value'] = level
            self.audio_level_label.config(text=f"{level}%")
        
        # Continue updating while recording
        if self.assistant.recorder.is_recording:
            self.after(100, self.update_audio_level)

    def toggle_answer_mode(self):
        """Cycle through answer quality modes: Default → Quick → Detailed → Code → Default"""
        modes = ["default", "quick", "detailed", "code"]
        mode_labels = {
            "default": "🔘 Default",
            "quick": "⚡ Quick",
            "detailed": "📝 Detailed", 
            "code": "💻 Code"
        }
        mode_descriptions = {
            "default": "Normal GPT behavior",
            "quick": "Short, concise answers",
            "detailed": "Comprehensive explanations",
            "code": "Focus on code examples"
        }
        
        # Cycle to next mode
        current_idx = modes.index(self.answer_mode)
        next_idx = (current_idx + 1) % len(modes)
        self.answer_mode = modes[next_idx]
        
        # Update button text
        self.answer_mode_btn.config(text=mode_labels[self.answer_mode])
        self.status.config(text=f"🎯 Answer Mode: {mode_descriptions[self.answer_mode]}")

    def get_answer_mode_instruction(self):
        """Return instruction text based on current answer mode"""
        instructions = {
            "default": "",  # No special instruction - normal GPT behavior
            "quick": "\n\n[INSTRUCTION: Provide SHORT, CONCISE answers like a real human expert would in a conversation. Be brief and natural. IMPORTANT: Sound like a real person - avoid AI phrases like 'Certainly!', 'I'd be happy to', 'As an AI'. Don't be overly formal or structured. Just answer naturally like a knowledgeable colleague would. Maximum 2-3 sentences.]",
            "detailed": "\n\n[INSTRUCTION: Provide COMPREHENSIVE explanations like a senior developer explaining to a teammate. IMPORTANT: Sound completely HUMAN - avoid typical AI patterns like 'Certainly!', 'Great question!', 'I'd be happy to help', 'As an AI'. Don't use excessive bullet points or overly structured formats. Speak naturally, use casual professional tone, include real-world context. Explain like you're a human expert having a conversation, not a chatbot giving a formal response.]",
            "code": "\n\n[INSTRUCTION: Focus on CODE EXAMPLES. Provide working code snippets with brief explanations. Prioritize practical, copy-paste ready code.]"
        }
        return instructions.get(self.answer_mode, "")

    def display_chat_history(self, max_rounds=20):
        self.response_box.config(state=tk.NORMAL)
        self.response_box.delete(1.0, tk.END)

        # Normalize / skip malformed messages
        ua = []
        for m in self.assistant.messages:
            if not isinstance(m, dict):
                continue
            role = m.get("role")
            if role in ("user", "assistant"):
                ua.append(m)

        # keep last max_rounds * 2 msgs
        keep = max_rounds * 2
        recent_ua = ua[-keep:] if len(ua) > keep else ua

        for msg in recent_ua:
            role = msg.get("role")
            content = msg.get("content", "")

            # Normalize content to text
            if isinstance(content, list):
                text = "\n".join(
                    c.get("text", "[non-text]")
                    if c.get("type") == "text"
                    else "[Image]"
                    for c in content
                )
            else:
                text = str(content)

            if role == "user":
                self.response_box.insert(
                    tk.END,
                    f"\n\n---------------------------------------------------------------------\nQUESTION: {text.strip()}\n"
                )
            elif role == "assistant":
                self.response_box.insert(
                    tk.END,
                    f"------------------\nANSWER: {text.strip()}\n"
                )

        self.response_box.config(state=tk.DISABLED)
        self.response_box.see(tk.END)

        # Re-apply any saved bookmarks for this session
        self._restore_bookmarks()


    def setup_ui(self):
        # Create paned window for sidebar and main content
        self.paned = ttk.PanedWindow(self, orient=tk.HORIZONTAL)
        self.paned.pack(fill="both", expand=True)

        # Create sidebar frame
        self.sidebar = ttk.Frame(self.paned, width=200)
        self.paned.add(self.sidebar, weight=0)
        
        # Create toggle button at top
        self.toggle_btn = ttk.Button(self.sidebar, text="☰", width=2, command=self.toggle_sidebar)
        self.toggle_btn.pack(pady=5, fill="x")

        # ====== RESIZABLE SIDEBAR SECTIONS ======
        # Use tk.PanedWindow (not ttk) because it supports per-pane minsize.
        # ttk.PanedWindow has no minsize — either section can collapse to 0 and
        # become impossible to recover by dragging.
        self.sidebar_paned = tk.PanedWindow(
            self.sidebar, orient=tk.VERTICAL,
            sashrelief=tk.RIDGE, sashwidth=6, sashpad=2
        )
        self.sidebar_paned.pack(fill="both", expand=True, padx=5, pady=5)

        # ----- TOP SECTION: Tabs/Subtabs -----
        self.tab_section = ttk.Frame(self.sidebar_paned)
        # minsize=120 → Prompts section can never collapse below 120px
        self.sidebar_paned.add(self.tab_section, minsize=120, stretch="always")
        
        ttk.Label(self.tab_section, text="📋 Prompts & Subtabs").pack(anchor="w")
        
        # Create tab treeview with scrollbar
        self.tab_frame = ttk.Frame(self.tab_section)
        self.tab_frame.pack(fill="both", expand=True)
        
        self.tab_tree = ttk.Treeview(self.tab_frame, show="tree", selectmode="browse")
        self.tab_tree.pack(fill="both", expand=True, side="left")
        self.tab_tree.bind("<<TreeviewSelect>>", self.on_tab_select)
        
        # Add drag-and-drop for tabs/subtabs
        self._setup_drag_drop(self.tab_tree, "tabs")
        
        # Add scrollbar to tab_tree
        tab_scrollbar = ttk.Scrollbar(self.tab_frame, orient="vertical", command=self.tab_tree.yview)
        tab_scrollbar.pack(side="right", fill="y")
        self.tab_tree.configure(yscrollcommand=tab_scrollbar.set)

        # Saved profiles (same panel as subtabs; double-click to apply)
        ttk.Label(self.tab_section, text="📋 Saved profiles").pack(anchor="w", pady=(8, 2))
        profile_list_frame = ttk.Frame(self.tab_section)
        profile_list_frame.pack(fill="x", pady=(0, 5))
        self.profile_listbox = tk.Listbox(
            profile_list_frame, height=4, font=("Arial", 9),
            selectmode=tk.SINGLE, activestyle="dotbox", highlightthickness=0
        )
        self.profile_listbox.pack(side="left", fill="both", expand=True)
        profile_scroll = ttk.Scrollbar(profile_list_frame, orient="vertical", command=self.profile_listbox.yview)
        profile_scroll.pack(side="right", fill="y")
        self.profile_listbox.configure(yscrollcommand=profile_scroll.set)
        self.profile_listbox.bind("<Double-Button-1>", self._on_profile_double_click)
        self.profile_listbox.bind("<Button-3>", self._on_profile_right_click)
        self.profile_listbox.bind("<Control-Button-1>", self._on_profile_right_click)
        profile_btn_frame = ttk.Frame(self.tab_section)
        profile_btn_frame.pack(fill="x")
        ttk.Button(profile_btn_frame, text="Edit order", width=10, command=self._edit_selected_profile_order).pack(side="left", padx=2)
        self._refresh_profile_list()

        # ----- BOTTOM SECTION: Chat History -----
        self.chat_section = ttk.Frame(self.sidebar_paned)
        # minsize=150 → Past Chats section can never collapse below 150px
        self.sidebar_paned.add(self.chat_section, minsize=150, stretch="always")
        
        ttk.Label(self.chat_section, text="💬 Past Chats").pack(anchor="w")
        
        # Create chat history treeview with scrollbar
        self.chat_frame = ttk.Frame(self.chat_section)
        self.chat_frame.pack(fill="both", expand=True)
        
        self.chat_tabs = ttk.Treeview(self.chat_frame, show="tree", selectmode="browse")
        self.chat_tabs.pack(fill="both", expand=True, side="left")
        self.chat_tabs.bind("<<TreeviewSelect>>", self.on_chat_tab_select)
        
        # Add drag-and-drop for chats
        self._setup_drag_drop(self.chat_tabs, "chats")
        
        # Add scrollbar to chat_tabs
        chat_scrollbar = ttk.Scrollbar(self.chat_frame, orient="vertical", command=self.chat_tabs.yview)
        chat_scrollbar.pack(side="right", fill="y")
        self.chat_tabs.configure(yscrollcommand=chat_scrollbar.set)

        # Create buttons frame (Default first so it's always visible in small windows)
        btn_frame = ttk.Frame(self.sidebar)
        btn_frame.pack(fill="x", padx=5, pady=5)

        # Quick Setup / Default first: multi-select subtabs + Set as default (hotkey Cmd+Shift+I)
        self.quick_setup_btn = ttk.Button(btn_frame, text="📌 Default", command=self.open_quick_setup, width=10)
        self.quick_setup_btn.pack(side="left", padx=2)

        self.add_tab_btn = ttk.Button(btn_frame, text="+ Tab", command=self.add_new_tab)
        self.add_tab_btn.pack(side="left", fill="x", expand=True, padx=2)

        self.add_subtab_btn = ttk.Button(btn_frame, text="+ Sub", command=self.add_new_subtab, state=tk.DISABLED)
        self.add_subtab_btn.pack(side="left", fill="x", expand=True, padx=2)
        
        self.delete_chat_btn = ttk.Button(self.sidebar, text="🗑 Delete Chat", command=self.delete_chat)
        self.delete_chat_btn.pack(side="left", padx=4)

        self.rename_chat_btn = ttk.Button(self.sidebar, text="✏️ Rename Chat", command=self.rename_chat)
        self.rename_chat_btn.pack(side="left", padx=4)


        # Create main content frame (existing UI)
        self.main_frame = ttk.Frame(self.paned)
        self.paned.add(self.main_frame, weight=1)

        # ====== HEADER BAR ======
        header_frame = ttk.Frame(self.main_frame)
        header_frame.pack(fill="x", padx=10, pady=(5, 2))
        
        # Connection status (left)
        self.connection_label = ttk.Label(header_frame, text="🔴", font=('Arial', 10))
        self.connection_label.pack(side="left", padx=(0, 5))
        
        # Main status label (left-center)
        self.status = ttk.Label(header_frame, text="🔊 Ready", font=('Arial', 10))
        self.status.pack(side="left", padx=5)
        
        # Audio level indicator (right)
        self.audio_level_frame = ttk.Frame(header_frame)
        self.audio_level_frame.pack(side="right", padx=5)
        ttk.Label(self.audio_level_frame, text="🎙", font=('Arial', 9)).pack(side="left")
        self.audio_level_bar = ttk.Progressbar(self.audio_level_frame, length=60, mode='determinate', maximum=100)
        self.audio_level_bar.pack(side="left", padx=2)
        self.audio_level_label = ttk.Label(self.audio_level_frame, text="--", font=('Arial', 8), width=4)
        self.audio_level_label.pack(side="left")
        
        # Check API connection on startup
        self.after(500, self.check_api_connection)

        text_frame = ttk.Frame(self.main_frame)
        text_frame.pack(fill="both", expand=True, padx=10)

        # ====== BOOKMARK/POINTER PANEL (like debug breakpoints) ======
        # Pack FIRST (side=right) so it reserves space before response_box expands
        self.bookmark_frame = ttk.Frame(text_frame, width=28)
        self.bookmark_frame.pack(side="right", fill="y", padx=(2, 0))
        self.bookmark_frame.pack_propagate(False)
        
        # Bookmark header with tooltip
        bookmark_header = ttk.Label(self.bookmark_frame, text="📍", font=('Arial', 9))
        bookmark_header.pack(pady=1)
        
        # Bookmark listbox (shows pointers)
        self.bookmark_listbox = tk.Listbox(
            self.bookmark_frame, 
            bg='#2d2d30', 
            fg='#ffd700',  # Gold color for pointers
            selectbackground='#4a4a00',
            selectforeground='#ffff00',
            font=('Arial', 8),
            width=3,
            highlightthickness=0,
            borderwidth=0
        )
        self.bookmark_listbox.pack(fill="both", expand=True)
        self.bookmark_listbox.bind("<<ListboxSelect>>", self._on_bookmark_click)
        self.bookmark_listbox.bind("<Double-Button-1>", self._on_bookmark_delete)
        
        # Store bookmarks: [(line_index, question_preview), ...]
        self.bookmarks = []
        self._current_bookmark_index = -1  # for "go to next bookmark" hotkey cycling
        self._current_chat_index = 0       # tracks which session is loaded (0 = AutoSave)

        # Scrollbar - pack second (side=right)
        scrollbar = ttk.Scrollbar(text_frame)
        scrollbar.pack(side="right", fill="y")

        # Response box - pack last (side=left, expand) takes remaining space
        self.response_box = tk.Text(text_frame, wrap=tk.WORD, font=('Consolas', self.assistant.font_size),
                                    bg='#343541', fg='white', insertbackground='white', selectbackground='#4E4E4E',
                                    highlightthickness=0)
        self.response_box.pack(side="left", fill="both", expand=True)
        self.response_box.insert(tk.END, "🤖 Start a new conversation or ask your first question...")
        self.response_box.config(state=tk.DISABLED, yscrollcommand=scrollbar.set)
        scrollbar.config(command=self.response_box.yview)
        
        self.response_box.tag_configure('code', foreground='#4EC9B0')
        self.response_box.tag_configure('bookmark_highlight', background='#4a4a00', foreground='#ffff00')
        
        # Right-click context menu for bookmarking
        self.response_box.bind("<Button-2>", self._show_bookmark_menu)  # Middle click on Mac
        self.response_box.bind("<Button-3>", self._show_bookmark_menu)  # Right click
        self.response_box.bind("<Control-Button-1>", self._show_bookmark_menu)  # Ctrl+click on Mac

        # Chat scroll keys (PgUp/PgDn/Up/Down) are bound on the main window above

        # ====== MODERN CONTROL PANEL - 2 ROWS ======
        control_container = ttk.Frame(self.main_frame)
        control_container.pack(fill="x", padx=10, pady=5)
        
        # ----- ROW 1: MAIN ACTIONS (Recording & Chat) -----
        row1 = ttk.Frame(control_container)
        row1.pack(fill="x", pady=(0, 3))
        
        # LEFT: Primary actions
        self.record_btn = ttk.Button(row1, text="🎤 Listen", command=self.toggle_recording, width=10)
        self.record_btn.pack(side="left", padx=2)
        
        self.stop_btn = ttk.Button(row1, text="⏹ Stop", command=self.stop_output, state=tk.DISABLED, width=8)
        self.stop_btn.pack(side="left", padx=2)
        
        # Separator
        ttk.Separator(row1, orient="vertical").pack(side="left", fill="y", padx=8)
        
        self.new_chat_btn = ttk.Button(row1, text="🆕 New Chat", command=self.start_new_chat, width=10)
        self.new_chat_btn.pack(side="left", padx=2)
        
        self.upload_btn = ttk.Button(row1, text="📁 Resume/JD", command=self.upload_resume, width=10)
        self.upload_btn.pack(side="left", padx=2)
        
        # RIGHT: Window controls
        self.topmost_btn = ttk.Button(row1, text="📌", command=self.toggle_always_on_top, width=3)
        self.topmost_btn.pack(side="right", padx=2)
        
        # UI Mode toggle (Classic/Modern)
        self.ui_mode_btn = ttk.Button(row1, text="🎨", command=self.toggle_ui_mode, width=3)
        self.ui_mode_btn.pack(side="right", padx=2)
        
        # ----- ROW 2: SETTINGS & OPTIONS -----
        row2 = ttk.Frame(control_container)
        row2.pack(fill="x", pady=(0, 0))
        
        # Model selector
        self.model_btn = ttk.Button(row2, text="🧠 4o", command=self.toggle_model, width=7)
        self.model_btn.pack(side="left", padx=2)
        
        # Answer mode
        self.answer_mode_btn = ttk.Button(row2, text="🔘 Default", command=self.toggle_answer_mode, width=10)
        self.answer_mode_btn.pack(side="left", padx=2)
        
        # Fast/Full mode (default: Full)
        self.optimize_btn = ttk.Button(row2, text="🐢 Full", command=self.toggle_optimization_mode, width=7)
        self.optimize_btn.pack(side="left", padx=2)
        
        # Separator
        ttk.Separator(row2, orient="vertical").pack(side="left", fill="y", padx=8)
        
        # Audio input toggle (compact)
        self.toggle_input_btn = ttk.Button(row2, text="🔈 BlackHole", command=self.toggle_input_mode, width=10)
        self.toggle_input_btn.pack(side="left", padx=2)
        
        # Separator
        ttk.Separator(row2, orient="vertical").pack(side="left", fill="y", padx=8)
        
        # Bookmarks
        self.bookmark_btn = ttk.Button(row2, text="🔖", command=self.add_bookmark_at_cursor, width=3)
        self.bookmark_btn.pack(side="left", padx=2)
        
        self.clear_bookmarks_btn = ttk.Button(row2, text="🗑", command=self.clear_all_bookmarks, width=3)
        self.clear_bookmarks_btn.pack(side="left", padx=2)
        
        # Report button
        self.diag_btn = ttk.Button(row2, text="📊", command=self.show_performance_dialog, width=3)
        self.diag_btn.pack(side="left", padx=2)

        # RIGHT: Font controls
        font_frame = ttk.Frame(row2)
        font_frame.pack(side="right", padx=2)
        ttk.Button(font_frame, text="A+", command=self.increase_font, width=3).pack(side="left")
        ttk.Button(font_frame, text="A-", command=self.decrease_font, width=3).pack(side="left")

        # ====== ROW 3 — v2 competitive features ======
        row3 = ttk.Frame(control_container)
        row3.pack(fill="x", pady=(2, 0))

        self.stealth_btn = ttk.Button(row3, text="🕵️ Stealth", command=self.toggle_stealth_mode, width=10)
        self.stealth_btn.pack(side="left", padx=2)

        self.vad_btn = ttk.Button(row3, text="🎙 VAD", command=self.toggle_vad_mode, width=8)
        self.vad_btn.pack(side="left", padx=2)

        self.local_asr_btn = ttk.Button(row3, text="🔊 Local ASR", command=self.toggle_local_asr, width=10)
        self.local_asr_btn.pack(side="left", padx=2)

        self.auto_copy_btn = ttk.Button(row3, text="📋 Auto-Copy", command=self.toggle_auto_copy, width=10)
        self.auto_copy_btn.pack(side="left", padx=2)

        self.mock_btn = ttk.Button(row3, text="🎯 Mock Interview", command=self.open_mock_interview, width=13)
        self.mock_btn.pack(side="left", padx=2)

        self.route_btn = ttk.Button(row3, text="⚡ Smart Route", command=self.toggle_smart_route, width=12)
        self.route_btn.pack(side="left", padx=2)

        # ====== INPUT BAR (bottom) ======
        input_frame = ttk.Frame(self.main_frame)
        input_frame.pack(side="bottom", fill="x", padx=10, pady=8)

        self.input_entry = tk.Text(input_frame, font=('Arial', 13), height=2,
                                   wrap="word", relief="solid", borderwidth=1)
        self.input_entry.pack(side="left", fill="x", expand=True, padx=(0, 8))

        # Enter / Numpad-Enter → send message; return "break" suppresses the
        # default newline that tk.Text would otherwise insert.
        self.input_entry.bind("<Return>", self._on_input_enter)
        self.input_entry.bind("<KP_Enter>", self._on_input_enter)   # Number pad Enter
        # Shift+Enter → insert a real newline without sending
        self.input_entry.bind("<Shift-Return>", self._on_input_shift_enter)

        self.submit_btn = ttk.Button(input_frame, text="Send ➡️", width=8, command=self.submit_text_question)
        self.submit_btn.pack(side="right")

    def load_chat_tabs(self):
        print("Chat Tabs:", self.chat_tabs)  # Debugging line
        self.chat_tabs.delete(*self.chat_tabs.get_children())
        for i, title in enumerate(self.chat_manager.get_titles()):
            self.chat_tabs.insert("", "end", iid=f"chat_{i}", text=title)
            
    def delete_chat(self):
        """
        Delete ALL chats except:
            - the currently selected chat
            - the "AutoSave - Last Session" chat
        """
        AUTO_TITLE = "AutoSave - Last Session"

        # Require a selection to know which one to keep (the 'current chat')
        selected = self.chat_tabs.selection()
        if not selected:
            messagebox.showwarning("No Chat Selected", "Please select the chat you want to KEEP.")
            return

        tab_id = selected[0]
        if not tab_id.startswith("chat_"):
            messagebox.showwarning("Invalid Selection", "Please select a valid chat.")
            return

        keep_index = int(tab_id.split("_")[1])
        titles = self.chat_manager.get_titles()
        if keep_index < 0 or keep_index >= len(titles):
            messagebox.showwarning("Invalid Selection", "Selected chat index is out of range.")
            return

        keep_title = titles[keep_index]

        # Build the new chat list, keeping only the selected chat and the AutoSave session (if present)
        original_count = len(self.chat_manager.sessions)
        new_sessions = []
        for i, s in enumerate(self.chat_manager.sessions):
            title = s.get("title", "Untitled")
            if i == keep_index or title == AUTO_TITLE:
                new_sessions.append(s)

        removed_count = original_count - len(new_sessions)
        if removed_count <= 0:
            messagebox.showinfo("Nothing to Delete", "There are no other chats to delete.")
            return

        confirm = messagebox.askyesno(
            "Delete Chats",
            f"This will permanently delete {removed_count} chat(s), keeping only:\n\n"
            f"• {keep_title}\n"
            f"• {AUTO_TITLE} (if it exists)\n\n"
            "Are you sure?"
        )
        if not confirm:
            return

        # Commit changes
        self.chat_manager.sessions = new_sessions
        self.chat_manager.save()

        # Refresh UI list
        self.load_chat_tabs()

        # Reselect the kept chat (find it again by title)
        reselect_index = 0
        for i, s in enumerate(self.chat_manager.sessions):
            if s.get("title", "") == keep_title:
                reselect_index = i
                break
        kept_item_id = f"chat_{reselect_index}"
        if self.chat_tabs.exists(kept_item_id):
            self.chat_tabs.selection_set(kept_item_id)
            self.chat_tabs.see(kept_item_id)

        self.status.config(text=f"🧹 Deleted {removed_count} chat(s). Kept: “{keep_title}” and “{AUTO_TITLE}”.")
        print(f"Deleted {removed_count} chat(s). Kept: {keep_title} (index {reselect_index}) and '{AUTO_TITLE}'.")


    def rename_chat(self):
        selected = self.chat_tabs.selection()
        if not selected:
            messagebox.showwarning("No Chat Selected", "Please select a chat first")
            return

        # Get the chat index from the selected tab
        tab_id = selected[0]
        if tab_id.startswith("chat_"):
            index = int(tab_id.split("_")[1])
            old_title = self.chat_manager.get_titles()[index]

            # Ask for a new name
            new_name = simpledialog.askstring("Rename Chat", f"Enter a new name for the chat '{old_title}':")
            if new_name:
                # Update the title in the chat history manager
                self.chat_manager.sessions[index]["title"] = new_name
                self.chat_manager.save()

                # Update the chat tab title in the UI
                self.chat_tabs.item(tab_id, text=new_name)

                self.status.config(text=f"🔄 Renamed chat to: {new_name}")
                print(f"Renamed chat to: {new_name}")



    def toggle_sidebar(self):
        self.sidebar_visible = not self.sidebar_visible

        if self.sidebar_visible:
            self.paned.forget(self.mini_sidebar)  # Remove mini toggle
            self.paned.insert(0, self.sidebar)    # Restore full sidebar
        else:
            self.paned.forget(self.sidebar)       # Remove full sidebar

            # Add a minimal sidebar with just the ☰ button
            self.mini_sidebar = ttk.Frame(self.paned, width=30)
            toggle_only_btn = ttk.Button(self.mini_sidebar, text="☰", command=self.toggle_sidebar, width=2)
            toggle_only_btn.pack(padx=2, pady=5, fill="x")
            self.paned.insert(0, self.mini_sidebar)



    def add_new_subtab(self):
    # Get selected tab
        selected = self.tab_tree.selection()
        if not selected:
            messagebox.showwarning("No Tab Selected", "Please select a tab first")
            return
                
        tab_id = selected[0]
        if not tab_id.startswith("tab_"):
            # Find parent tab
            parent = self.tab_tree.parent(tab_id)
            if parent:
                tab_id = parent
        
        if tab_id.startswith("tab_"):
            tab_index = int(tab_id.split("_")[1])
            name = simpledialog.askstring("New Subtask", "Enter subtab name:")
            if name:
                prompt = simpledialog.askstring("Prompt", "Enter prompt text:")
                
                # Use the prompt text as both the prompt and text_input
                text_input = prompt  # Directly use the prompt text as the input for this subtab

                subtab_index = self.prompt_manager.add_subtab(tab_index, name, prompt or "", text_input or "")
                self.tab_tree.insert(tab_id, "end", text=name, iid=f"sub_{tab_index}_{subtab_index}")

    # ============================================================================
    # QUICK SETUP - Multi-select subtabs and send in ONE message
    # ============================================================================
    
    def open_quick_setup(self):
        """Open Quick Setup dialog with multi-select subtabs and profiles."""
        dialog = tk.Toplevel(self)
        dialog.title("🚀 Quick Setup - Interview Initialization")
        dialog.geometry("500x600")
        dialog.transient(self)
        dialog.grab_set()
        
        # Center on parent
        dialog.update_idletasks()
        x = self.winfo_x() + (self.winfo_width() // 2) - (500 // 2)
        y = self.winfo_y() + (self.winfo_height() // 2) - (600 // 2)
        dialog.geometry(f"+{x}+{y}")
        
        # ---- PROFILES SECTION ----
        profile_frame = ttk.LabelFrame(dialog, text="📋 Saved Profiles (One-Click Apply)")
        profile_frame.pack(fill="x", padx=10, pady=5)
        
        profiles = self._load_setup_profiles()
        self._profile_vars = {}
        
        profile_btn_frame = ttk.Frame(profile_frame)
        profile_btn_frame.pack(fill="x", padx=5, pady=5)
        
        if profiles:
            for i, (name, subtab_ids) in enumerate(profiles.items()):
                btn = ttk.Button(
                    profile_btn_frame, 
                    text=f"▶ {name}", 
                    command=lambda n=name, s=subtab_ids, d=dialog: self._apply_profile(n, s, d)
                )
                btn.pack(side="left", padx=2, pady=2)
                if i >= 4:  # Limit visible profiles
                    break
        else:
            ttk.Label(profile_btn_frame, text="No saved profiles yet", foreground="gray").pack()
        
        # ---- SUBTABS SELECTION ----
        select_frame = ttk.LabelFrame(dialog, text="☑️ Select Prompts to Apply (multi-select)")
        select_frame.pack(fill="both", expand=True, padx=10, pady=5)
        
        # Scrollable canvas for checkboxes
        canvas = tk.Canvas(select_frame)
        scrollbar = ttk.Scrollbar(select_frame, orient="vertical", command=canvas.yview)
        scrollable_frame = ttk.Frame(canvas)
        
        scrollable_frame.bind(
            "<Configure>",
            lambda e: canvas.configure(scrollregion=canvas.bbox("all"))
        )
        
        canvas.create_window((0, 0), window=scrollable_frame, anchor="nw")
        canvas.configure(yscrollcommand=scrollbar.set)
        
        canvas.pack(side="left", fill="both", expand=True)
        scrollbar.pack(side="right", fill="y")
        
        # Create checkboxes for each tab/subtab
        self._setup_checkboxes = {}  # {subtab_id: (var, text)}
        
        for tab_idx in range(self.prompt_manager.get_tab_count()):
            tab_name = self.prompt_manager.get_tab_name(tab_idx)
            
            # Tab header
            tab_label = ttk.Label(scrollable_frame, text=f"📁 {tab_name}", font=('Arial', 11, 'bold'))
            tab_label.pack(anchor="w", pady=(10, 2), padx=5)
            
            # Subtabs
            for sub_idx in range(self.prompt_manager.get_subtab_count(tab_idx)):
                sub_name = self.prompt_manager.get_subtab_name(tab_idx, sub_idx)
                sub_text = self.prompt_manager.get_subtab_text_input(tab_idx, sub_idx) or \
                           self.prompt_manager.get_subtab_prompt(tab_idx, sub_idx) or sub_name
                
                var = tk.BooleanVar(value=False)
                subtab_id = f"sub_{tab_idx}_{sub_idx}"
                self._setup_checkboxes[subtab_id] = (var, sub_text, sub_name)
                
                cb = ttk.Checkbutton(
                    scrollable_frame, 
                    text=f"  {sub_name}", 
                    variable=var
                )
                cb.pack(anchor="w", padx=20)
        
        # Pre-check default interview subtabs if set (so user sees and can edit current default)
        default_ids = self.ui_prefs.get("default_interview_subtabs") or []
        for sid in default_ids:
            if sid in self._setup_checkboxes:
                self._setup_checkboxes[sid][0].set(True)
        
        # ---- ACTION BUTTONS ----
        action_frame = ttk.Frame(dialog)
        action_frame.pack(fill="x", padx=10, pady=10)
        
        # Select/Deselect All
        ttk.Button(
            action_frame, 
            text="☑ Select All", 
            command=lambda: self._toggle_all_checkboxes(True)
        ).pack(side="left", padx=2)
        
        ttk.Button(
            action_frame, 
            text="☐ Deselect All", 
            command=lambda: self._toggle_all_checkboxes(False)
        ).pack(side="left", padx=2)
        
        ttk.Separator(action_frame, orient="vertical").pack(side="left", fill="y", padx=10)
        
        # Set as default interview (saves current selection as default; used by hotkey until you change it)
        ttk.Button(
            action_frame,
            text="📌 Set as default",
            command=lambda: self._save_default_interview_subtabs(dialog)
        ).pack(side="left", padx=2)
        if default_ids:
            ttk.Button(
                action_frame,
                text="✏️ Edit default order",
                command=self._open_edit_default_order_dialog
            ).pack(side="left", padx=2)
        
        # Save as Profile
        ttk.Button(
            action_frame, 
            text="💾 Save Profile", 
            command=lambda: self._save_current_as_profile(dialog)
        ).pack(side="left", padx=2)
        
        # Apply button (main action)
        apply_btn = ttk.Button(
            action_frame, 
            text="🚀 APPLY SELECTED", 
            command=lambda: self._apply_quick_setup(dialog),
            style="Accent.TButton"
        )
        apply_btn.pack(side="right", padx=2)
        
        # Status label (show current default count if set)
        default_count = len([x for x in (self.ui_prefs.get("default_interview_subtabs") or []) if x in self._setup_checkboxes])
        status_text = f"Default interview: {default_count} prompts. Select subtabs → Set as default → use hotkey Cmd+Shift+I in interview." if default_count else "Select prompts and click Apply to send in ONE message"
        self._setup_status = ttk.Label(dialog, text=status_text)
        self._setup_status.pack(pady=5)
        
        # Keyboard shortcut
        dialog.bind("<Return>", lambda e: self._apply_quick_setup(dialog))
        dialog.bind("<Escape>", lambda e: dialog.destroy())
    
    def _toggle_all_checkboxes(self, select: bool):
        """Select or deselect all checkboxes."""
        for subtab_id, (var, _, _) in self._setup_checkboxes.items():
            var.set(select)
    
    def _load_setup_profiles(self) -> dict:
        """Load saved setup profiles from file."""
        profile_path = os.path.join(os.path.dirname(__file__), "setup_profiles.json")
        try:
            if os.path.exists(profile_path):
                with open(profile_path, "r") as f:
                    return json.load(f)
        except Exception as e:
            print(f"Error loading profiles: {e}")
        return {}
    
    def _save_setup_profiles(self, profiles: dict):
        """Save setup profiles to file."""
        profile_path = os.path.join(os.path.dirname(__file__), "setup_profiles.json")
        try:
            with open(profile_path, "w") as f:
                json.dump(profiles, f, indent=2)
        except Exception as e:
            print(f"Error saving profiles: {e}")
    
    def _save_default_interview_subtabs(self, dialog):
        """Save currently selected subtabs as the default interview set. Hotkey Cmd+Shift+I will use this until you change it."""
        selected = [sid for sid, (var, _, _) in self._setup_checkboxes.items() if var.get()]
        if not selected:
            self._setup_status.config(text="⚠ Select at least one prompt, then click Set as default.")
            return
        self.ui_prefs["default_interview_subtabs"] = selected
        UIPreferences.save(self.ui_prefs)
        self._setup_status.config(text=f"✅ Default interview set: {len(selected)} prompts. Use Cmd+Shift+I in interview to feed them.")
        messagebox.showinfo("Default set", f"Default interview set to {len(selected)} prompts.\n\nIn interview: attach resume, paste JD, then press Cmd+Shift+I to feed all instructions at once.")

    def _save_current_as_profile(self, dialog):
        """Save currently selected subtabs as a profile."""
        selected = [sid for sid, (var, _, _) in self._setup_checkboxes.items() if var.get()]
        
        if not selected:
            messagebox.showwarning("No Selection", "Please select at least one prompt to save as profile")
            return
        
        name = simpledialog.askstring("Save Profile", "Enter profile name:", parent=dialog)
        if name:
            profiles = self._load_setup_profiles()
            profiles[name] = selected
            self._save_setup_profiles(profiles)
            self._refresh_profile_list()
            self._setup_status.config(text=f"✅ Profile '{name}' saved with {len(selected)} prompts!")
            messagebox.showinfo("Saved", f"Profile '{name}' saved!\nIt appears in the sidebar under Saved profiles — double-click to apply.")
    
    def _refresh_profile_list(self):
        """Reload the Saved profiles listbox from setup_profiles.json."""
        if not hasattr(self, "profile_listbox"):
            return
        self.profile_listbox.delete(0, tk.END)
        profiles = self._load_setup_profiles()
        for name in sorted(profiles.keys()):
            self.profile_listbox.insert(tk.END, name)

    def _on_profile_double_click(self, event):
        """Double-click on a saved profile: apply it one-by-one."""
        sel = self.profile_listbox.curselection()
        if not sel:
            return
        name = self.profile_listbox.get(sel[0])
        profiles = self._load_setup_profiles()
        if name not in profiles:
            self.status.config(text=f"Profile '{name}' not found.")
            return
        self.apply_profile_by_ids(profiles[name], profile_name=name)

    def _on_profile_right_click(self, event):
        """Right-click: show menu Apply / Edit order."""
        sel = self.profile_listbox.curselection()
        if not sel:
            return
        name = self.profile_listbox.get(sel[0])
        menu = tk.Menu(self, tearoff=0)
        menu.add_command(label="Apply (one-by-one)", command=lambda: self._apply_profile_by_name(name))
        menu.add_command(label="Edit order", command=lambda: self._open_edit_profile_order_dialog(name))
        menu.tk_popup(event.x_root, event.y_root)

    def _apply_profile_by_name(self, name):
        profiles = self._load_setup_profiles()
        if name in profiles:
            self.apply_profile_by_ids(profiles[name], profile_name=name)

    def _edit_selected_profile_order(self):
        """Edit order of the selected profile (from Edit order button)."""
        sel = self.profile_listbox.curselection()
        if not sel:
            self.status.config(text="Select a profile first, then click Edit order.")
            return
        name = self.profile_listbox.get(sel[0])
        self._open_edit_profile_order_dialog(name)

    def _open_edit_profile_order_dialog(self, profile_name: str):
        """Dialog to reorder subtabs in a profile: listbox with Move Up / Move Down / Save."""
        profiles = self._load_setup_profiles()
        if profile_name not in profiles:
            messagebox.showwarning("Not found", f"Profile '{profile_name}' not found.")
            return
        subtab_ids = list(profiles[profile_name])
        if not subtab_ids:
            messagebox.showinfo("Empty", "Profile has no prompts.")
            return
        dialog = tk.Toplevel(self)
        dialog.title(f"Edit order: {profile_name}")
        dialog.geometry("320x340")
        dialog.transient(self)
        ttk.Label(dialog, text=f"Reorder prompts (top = first sent):").pack(anchor="w", padx=10, pady=5)
        list_frame = ttk.Frame(dialog)
        list_frame.pack(fill="both", expand=True, padx=10, pady=5)
        order_listbox = tk.Listbox(list_frame, height=12, font=("Arial", 10), selectmode=tk.SINGLE)
        order_listbox.pack(side="left", fill="both", expand=True)
        scroll = ttk.Scrollbar(list_frame, orient="vertical", command=order_listbox.yview)
        scroll.pack(side="right", fill="y")
        order_listbox.configure(yscrollcommand=scroll.set)
        names_and_ids = []
        for sid in subtab_ids:
            _, names = self._get_combined_prompt_for_subtab_ids([sid])
            name_display = names[0] if names else sid
            names_and_ids.append((name_display, sid))
            order_listbox.insert(tk.END, name_display)
        def refresh_listbox():
            order_listbox.delete(0, tk.END)
            for (nm, _) in names_and_ids:
                order_listbox.insert(tk.END, nm)
        def move_up():
            sel = order_listbox.curselection()
            if not sel or sel[0] == 0:
                return
            i = sel[0]
            names_and_ids[i], names_and_ids[i - 1] = names_and_ids[i - 1], names_and_ids[i]
            refresh_listbox()
            order_listbox.selection_set(i - 1)
        def move_down():
            sel = order_listbox.curselection()
            if not sel or sel[0] >= len(names_and_ids) - 1:
                return
            i = sel[0]
            names_and_ids[i], names_and_ids[i + 1] = names_and_ids[i + 1], names_and_ids[i]
            refresh_listbox()
            order_listbox.selection_set(i + 1)
        def save_order():
            new_ids = [x[1] for x in names_and_ids]
            profiles = self._load_setup_profiles()
            profiles[profile_name] = new_ids
            self._save_setup_profiles(profiles)
            self._refresh_profile_list()
            dialog.destroy()
            self.status.config(text=f"✅ Order saved for '{profile_name}'.")
        btn_frame = ttk.Frame(dialog)
        btn_frame.pack(fill="x", padx=10, pady=10)
        ttk.Button(btn_frame, text="↑ Move up", command=move_up).pack(side="left", padx=2)
        ttk.Button(btn_frame, text="↓ Move down", command=move_down).pack(side="left", padx=2)
        ttk.Button(btn_frame, text="Save order", command=save_order).pack(side="right", padx=2)
        dialog.bind("<Escape>", lambda e: dialog.destroy())

    def _open_edit_default_order_dialog(self):
        """Dialog to reorder default interview subtabs (hotkey Cmd+Shift+I)."""
        default_ids = self.ui_prefs.get("default_interview_subtabs") or []
        if not default_ids:
            messagebox.showinfo("No default", "Set a default interview first (select subtabs → Set as default).")
            return
        dialog = tk.Toplevel(self)
        dialog.title("Edit default interview order")
        dialog.geometry("320x340")
        dialog.transient(self)
        ttk.Label(dialog, text="Reorder prompts (top = first sent). Hotkey: Cmd+Shift+I").pack(anchor="w", padx=10, pady=5)
        list_frame = ttk.Frame(dialog)
        list_frame.pack(fill="both", expand=True, padx=10, pady=5)
        order_listbox = tk.Listbox(list_frame, height=12, font=("Arial", 10), selectmode=tk.SINGLE)
        order_listbox.pack(side="left", fill="both", expand=True)
        scroll = ttk.Scrollbar(list_frame, orient="vertical", command=order_listbox.yview)
        scroll.pack(side="right", fill="y")
        order_listbox.configure(yscrollcommand=scroll.set)
        names_and_ids = []
        for sid in default_ids:
            _, names = self._get_combined_prompt_for_subtab_ids([sid])
            name_display = names[0] if names else sid
            names_and_ids.append((name_display, sid))
            order_listbox.insert(tk.END, name_display)
        def refresh_listbox():
            order_listbox.delete(0, tk.END)
            for (nm, _) in names_and_ids:
                order_listbox.insert(tk.END, nm)
        def move_up():
            sel = order_listbox.curselection()
            if not sel or sel[0] == 0:
                return
            i = sel[0]
            names_and_ids[i], names_and_ids[i - 1] = names_and_ids[i - 1], names_and_ids[i]
            refresh_listbox()
            order_listbox.selection_set(i - 1)
        def move_down():
            sel = order_listbox.curselection()
            if not sel or sel[0] >= len(names_and_ids) - 1:
                return
            i = sel[0]
            names_and_ids[i], names_and_ids[i + 1] = names_and_ids[i + 1], names_and_ids[i]
            refresh_listbox()
            order_listbox.selection_set(i + 1)
        def save_order():
            self.ui_prefs["default_interview_subtabs"] = [x[1] for x in names_and_ids]
            UIPreferences.save(self.ui_prefs)
            dialog.destroy()
            self.status.config(text="✅ Default interview order saved.")
        btn_frame = ttk.Frame(dialog)
        btn_frame.pack(fill="x", padx=10, pady=10)
        ttk.Button(btn_frame, text="↑ Move up", command=move_up).pack(side="left", padx=2)
        ttk.Button(btn_frame, text="↓ Move down", command=move_down).pack(side="left", padx=2)
        ttk.Button(btn_frame, text="Save order", command=save_order).pack(side="right", padx=2)
        dialog.bind("<Escape>", lambda e: dialog.destroy())

    def apply_profile_by_ids(self, subtab_ids: list, profile_name: str = None):
        """Apply a list of subtab IDs one-by-one: send each prompt to GPT, get answer, then next. Order = subtab_ids order. 'Intro' subtab is always run first if present."""
        valid_ids = []
        for sid in subtab_ids:
            combined, names = self._get_combined_prompt_for_subtab_ids([sid])
            if combined:
                valid_ids.append(sid)
        if not valid_ids:
            self.status.config(text="No valid prompts in this profile.")
            return
        # Put "Intro" first so you can copy it quickly (rest keep original order)
        intro_ids = []
        other_ids = []
        for sid in valid_ids:
            _, names = self._get_combined_prompt_for_subtab_ids([sid])
            name = (names[0] or "").strip().lower() if names else ""
            if name == "intro":
                intro_ids.append(sid)
            else:
                other_ids.append(sid)
        valid_ids = intro_ids + other_ids
        self._profile_queue = list(valid_ids)
        self._profile_name = profile_name or "Profile"
        self._profile_first = True
        self.response_box.config(state=tk.NORMAL)
        self.response_box.insert(
            tk.END,
            f"\n\n---------------------------------------------------------------------\n"
            f"🚀 {self._profile_name}: {len(valid_ids)} prompts (one-by-one)\n"
        )
        self.response_box.config(state=tk.DISABLED)
        if self.response_box.yview()[1] >= 0.99:
            self.response_box.see(tk.END)
        self._send_next_profile_prompt()

    def _send_next_profile_prompt(self):
        """Send the next prompt in _profile_queue; called after each stream completes."""
        if not getattr(self, "_profile_queue", None):
            return
        if not self._profile_queue:
            self.status.config(text=f"✅ {getattr(self, '_profile_name', 'Profile')}: all done!")
            return
        sid = self._profile_queue.pop(0)
        combined_prompt, selected_names = self._get_combined_prompt_for_subtab_ids([sid])
        if not combined_prompt:
            self.after(0, self._on_profile_stream_complete)
            return
        if getattr(self, "_profile_first", False):
            self._profile_first = False
            current = self.input_entry.get("1.0", tk.END).strip()
            if current:
                combined_prompt = f"{current}\n\n{combined_prompt}"
            self.input_entry.delete("1.0", tk.END)
            if hasattr(self, "pending_attachments") and self.pending_attachments:
                content = [{"type": "text", "text": combined_prompt}]
                content.extend(self.pending_attachments)
                del self.pending_attachments
                self.assistant.messages.append({"role": "user", "content": content})
            else:
                self.assistant.messages.append({"role": "user", "content": combined_prompt})
        else:
            self.assistant.messages.append({"role": "user", "content": combined_prompt})
        name = selected_names[0] if selected_names else sid
        self.response_box.config(state=tk.NORMAL)
        self.response_box.insert(tk.END, f"\n\n---- QUESTION ({name}) ----\n")
        self.response_box.config(state=tk.DISABLED)
        if self.response_box.yview()[1] >= 0.99:
            self.response_box.see(tk.END)
        self.chat_manager.save_current_session(self.assistant.messages)
        self.assistant.cancel_streaming()
        remaining = len(self._profile_queue)
        self.status.config(text=f"📌 {self._profile_name}: {name} ({(remaining + 1)} left)")
        self.assistant.stream_gpt_response(
            self.response_box, self.status, self.record_btn,
            on_complete=self._on_profile_stream_complete
        )

    def _on_profile_stream_complete(self):
        """When one-by-one stream finishes, wait for streaming=False then send next (so full answer is shown)."""
        def maybe_send_next():
            if self.assistant.streaming:
                self.after(400, maybe_send_next)
                return
            if getattr(self, "_profile_queue", None) and self._profile_queue:
                self.after(300, self._send_next_profile_prompt)
            else:
                self.status.config(text=f"✅ {getattr(self, '_profile_name', 'Profile')}: all done!")
        maybe_send_next()

    def _apply_profile(self, profile_name: str, subtab_ids: list, dialog):
        """Apply a saved profile - select checkboxes and apply."""
        # First, deselect all
        self._toggle_all_checkboxes(False)
        
        # Select the ones in the profile
        for subtab_id in subtab_ids:
            if subtab_id in self._setup_checkboxes:
                self._setup_checkboxes[subtab_id][0].set(True)
        
        # Auto-apply
        self._apply_quick_setup(dialog)
    
    def _apply_quick_setup(self, dialog):
        """Apply selected prompts one-by-one (ordered as in tree)."""
        selected_ids = [sid for sid, (var, _, _) in self._setup_checkboxes.items() if var.get()]
        if not selected_ids:
            messagebox.showwarning("No Selection", "Please select at least one prompt")
            return
        dialog.destroy()
        self.apply_profile_by_ids(selected_ids, profile_name="Quick Setup")

    def _get_combined_prompt_for_subtab_ids(self, subtab_ids):
        """Build combined text and names for a list of subtab IDs (e.g. ['sub_0_0','sub_0_1']). Skips invalid IDs."""
        selected_texts = []
        selected_names = []
        for sid in subtab_ids:
            if not isinstance(sid, str) or not sid.startswith("sub_"):
                continue
            parts = sid.split("_")
            if len(parts) != 3:
                continue
            try:
                tab_idx = int(parts[1])
                sub_idx = int(parts[2])
            except ValueError:
                continue
            if tab_idx < 0 or tab_idx >= self.prompt_manager.get_tab_count():
                continue
            if sub_idx < 0 or sub_idx >= self.prompt_manager.get_subtab_count(tab_idx):
                continue
            text = self.prompt_manager.get_subtab_text_input(tab_idx, sub_idx) or \
                   self.prompt_manager.get_subtab_prompt(tab_idx, sub_idx) or ""
            name = self.prompt_manager.get_subtab_name(tab_idx, sub_idx) or sid
            selected_texts.append((text or name).strip())
            selected_names.append(name)
        combined = "\n\n---\n\n".join(selected_texts) if selected_texts else ""
        return combined, selected_names

    def apply_default_interview_instructions(self):
        """Feed default interview subtabs one-by-one. Use after attaching resume and pasting JD; hotkey Cmd+Shift+I."""
        self.ui_prefs = UIPreferences.load()
        default_ids = self.ui_prefs.get("default_interview_subtabs") or []
        if not default_ids:
            self.status.config(text="📌 No default interview set. Open Quick Setup (🚀) → select subtabs → Set as default.")
            return
        self.apply_profile_by_ids(default_ids, profile_name="Default interview")

    def on_tab_select(self, event):
        # Reentrancy guard (TreeviewSelect can fire more than once)
        if getattr(self, "_subtab_sending", False):
            return

        selected = self.tab_tree.selection()
        if not selected:
            self.current_tab = -1
            self.current_subtab = -1
            return

        item_id = selected[0]

        # Tab (top level) clicked: just remember selection; do not send
        if item_id.startswith("tab_"):
            self.current_tab = int(item_id.split("_")[1])
            self.current_subtab = -1
            self.add_subtab_btn.config(state=tk.NORMAL)
            return

        # Only handle subtab clicks from here on
        if not item_id.startswith("sub_"):
            return

        parts = item_id.split("_")
        self.current_tab = int(parts[1])
        self.current_subtab = int(parts[2])
        self.add_subtab_btn.config(state=tk.NORMAL)

        # Prefer text_input, then prompt, then subtab name
        sub_name   = self.prompt_manager.get_subtab_name(self.current_tab, self.current_subtab) or "Request"
        prompt     = self.prompt_manager.get_subtab_prompt(self.current_tab, self.current_subtab) or ""
        text_input = self.prompt_manager.get_subtab_text_input(self.current_tab, self.current_subtab) or ""
        sub_text   = (text_input or prompt or sub_name).strip()

        # If nothing found, still send the subtab name as a fallback
        if not sub_text:
            sub_text = sub_name

        # Append sub_text to whatever is already in the entry (do NOT wipe)
        current = self.input_entry.get("1.0", tk.END).strip()
        prefix = "" if (not current or current.endswith((" ", "\n", "\t"))) else " "
        self.input_entry.insert(tk.END, f"{prefix}{sub_text}")

        # Auto-send the combined message (existing text + subchat text + any queued images)
        try:
            self._subtab_sending = True  # guard
            # submit_text_question() will:
            #  - read the entry,
            #  - include self.pending_attachments (if any),
            #  - clear entry & attachments,
            #  - stream the response.
            self.submit_text_question()

            # UX hint
            if hasattr(self, "pending_attachments"):
                # (submit_text_question clears pending_attachments after sending)
                pass
            self.status.config(text="🚀 Sub chat sent with your current text and any attached image(s).")
        finally:
            self._subtab_sending = False




        
    def handle_paste(self, event=None):
        # 1) Try text first
        try:
            text = self.clipboard_get()
            if isinstance(text, str) and text.strip():
                # ✅ Append text at caret, do not clear previous input or attachments
                self.input_entry.insert(tk.INSERT, text)
                # keep any queued attachments
                return "break"
        except tk.TclError:
            # no text in clipboard; try image next
            pass
        except Exception as e:
            print(f"❌ Paste (text) failed: {e}")
            self.status.config(text=f"❌ Paste error: {e}")
            # fall through to image

        # 2) Try image
        try:
            image = ImageGrab.grabclipboard()
            if isinstance(image, Image.Image):
                # Compress image for faster transmission and lower token cost
                b64_image = compress_image_png(image, max_size=1280)
                original_size = image.size[0] * image.size[1] * 3 // 1024  # Rough KB estimate
                compressed_size = len(b64_image) // 1024
                print(f"📎 Image compressed: ~{original_size}KB → {compressed_size}KB")

                # ✅ Keep a growing list of attachments
                if not hasattr(self, 'pending_attachments'):
                    self.pending_attachments = []

                self.pending_attachments.append({
                    "type": "image_url",
                    "image_url": {"url": f"data:image/png;base64,{b64_image}"}
                })

                # ✅ Insert a placeholder at caret (don't wipe existing text)
                idx = len(self.pending_attachments)
                self.input_entry.insert(tk.INSERT, f" [📎 Image {idx}] ")

                self.status.config(text=f"📎 {idx} image(s) attached ({compressed_size}KB). Paste more or Enter to send.")
                return "break"
        except Exception as e:
            print(f"❌ Paste (image) failed: {e}")
            self.status.config(text=f"❌ Paste error: {e}")

        # Let native paste happen if neither text nor image detected
        return None





    # def setup_ui(self):
    #     self.status = ttk.Label(self, text="🔊 Ready", style='TLabel')
    #     self.status.pack(pady=5, anchor="w", padx=10)
        


    #     text_frame = ttk.Frame(self)
    #     text_frame.pack(fill="both", expand=True, padx=10)

    #     self.response_box = tk.Text(
    #         text_frame, wrap=tk.WORD, font=('Consolas', self.assistant.font_size),
    #         bg='#343541', fg='white', insertbackground='white',
    #         selectbackground='#4E4E4E', highlightthickness=0
    #     )
    #     self.response_box.pack(side="left", fill="both", expand=True)
    #     self.response_box.insert(tk.END, "🤖 Start a new conversation or ask your first question...")
    #     self.response_box.config(state=tk.DISABLED)
    #     self.response_box.tag_configure('code', foreground='#4EC9B0')

    #     scrollbar = ttk.Scrollbar(text_frame, command=self.response_box.yview)
    #     scrollbar.pack(side="right", fill="y")
    #     self.response_box.config(yscrollcommand=scrollbar.set)

    #     # Button bar above chat entry
    #     control_frame = ttk.Frame(self)
    #     control_frame.pack(fill="x", padx=10, pady=5)

    #     self.record_btn = ttk.Button(control_frame, text="🎤 Listen", command=self.toggle_recording)
    #     self.record_btn.pack(side="left", padx=4)
        
    #     self.toggle_input_btn = ttk.Button(control_frame, text="🔈 Internal Audio (BlackHole)", command=self.toggle_input_mode)
    #     self.toggle_input_btn.pack(side="left", padx=4)
        
        

    #     self.new_chat_btn = ttk.Button(control_frame, text="🆕 New", command=self.start_new_chat)
    #     self.new_chat_btn.pack(side="left", padx=4)

    #     self.stop_btn = ttk.Button(control_frame, text="⏹ Stop", command=self.stop_output, state=tk.DISABLED)
    #     self.stop_btn.pack(side="left", padx=4)

    #     self.upload_btn = ttk.Button(control_frame, text="📁 Resume", command=self.upload_resume)
    #     self.upload_btn.pack(side="left", padx=4)

    #     font_controls = ttk.Frame(control_frame)
    #     font_controls.pack(side="left", padx=10)
    #     ttk.Button(font_controls, text="A+", command=self.increase_font).pack(side="left")
    #     ttk.Button(font_controls, text="A-", command=self.decrease_font).pack(side="left")
        
        


    #     self.topmost_btn = ttk.Button(control_frame, text="📌 Pin", command=self.toggle_always_on_top)
    #     self.topmost_btn.pack(side="right", padx=4)

    #     # Chat input bar at bottom
    #     input_frame = ttk.Frame(self)
    #     input_frame.pack(side="bottom", fill="x", padx=10, pady=5)

    #     self.input_entry = ttk.Entry(input_frame, font=('Arial', 14), width=80)
    #     self.input_entry.pack(side="left", fill="x", expand=True, padx=(5, 10))

    #     self.input_entry.bind("<Return>", lambda event: self.submit_text_question())

    #     self.submit_btn = ttk.Button(input_frame, text="➡️", width=4, command=self.submit_text_question)
    #     self.submit_btn.pack(side="right")
        
    #     self.load_chat_tabs()

    def capture_and_submit_screenshot(self):
        self.status.config(text="📸 Capturing screen...")
        print("📸 Got the screen capture")

        screenshot = pyautogui.screenshot()
        
        # Use compressed image for faster transmission
        b64_image = compress_image_png(screenshot, max_size=1280)
        print(f"📸 Screenshot compressed: {len(b64_image) // 1024}KB")

        # Build multimodal message (text + image)
        content = [
            {"type": "text", "text": "Please analyze this screenshot."},
            {"type": "image_url", "image_url": {"url": f"data:image/png;base64,{b64_image}"}}
        ]

        # UI preview
        self.response_box.config(state=tk.NORMAL)
        self.response_box.insert(
            tk.END,
            "\n\n---------------------------------------------------------------------\nQUESTION: [Screenshot attached]\n"
        )
        self.response_box.config(state=tk.DISABLED)
        self.response_box.see(tk.END)

        # Add to chat history & stream
        self.assistant.messages.append({"role": "user", "content": content})
        self.chat_manager.save_current_session(self.assistant.messages)

        self.status.config(text="🧠 Analyzing screenshot...")
        threading.Thread(
            target=self.assistant.stream_gpt_response,
            args=(self.response_box, self.status, self.record_btn),
            daemon=True
        ).start()


    def _on_input_enter(self, event=None):
        """Enter / KP_Enter → send the message. Return 'break' to suppress the
        newline that tk.Text would insert by default."""
        self.submit_text_question()
        return "break"

    def _on_input_shift_enter(self, event=None):
        """Shift+Enter → insert a real newline so the user can write multi-line
        questions without accidentally sending. Return 'break' to stop the event
        reaching the default <Return> binding."""
        self.input_entry.insert(tk.INSERT, "\n")
        return "break"

    def submit_text_question(self):
        question = self.input_entry.get("1.0", tk.END).strip()
        self.input_entry.delete("1.0", tk.END)

        # If nothing to send at all
        if not question and not hasattr(self, 'pending_attachments'):
            return

        # Screenshot shortcut remains
        if question == "--":
            self.capture_and_submit_screenshot()
            return

        content = []

        # ✅ Always include text if present
        if question:
            content.append({"type": "text", "text": question})

        # ✅ Always include any queued images
        if hasattr(self, 'pending_attachments'):
            content.extend(self.pending_attachments)
            del self.pending_attachments  # clear only after sending

        # Flatten for UI display
        flat_text = "\n".join(
            c["text"] if c["type"] == "text" else "[Image]" for c in content
        )

        self.response_box.config(state=tk.NORMAL)
        self.response_box.insert(
            tk.END,
            f"\n\n---------------------------------------------------------------------\nQUESTION: {flat_text.strip()}\n"
        )
        self.response_box.config(state=tk.DISABLED)
        self.response_box.see(tk.END)

        # Send to GPT
        if any(c["type"] == "image_url" for c in content):
            self.assistant.messages.append({"role": "user", "content": content})
        else:
            self.assistant.messages.append({"role": "user", "content": flat_text})

        # Mark as dirty so auto-persist can save if you switch chats
        self._last_persisted_hash = None

        self._last_question_type = classify_question(flat_text)
        self.chat_manager.save_current_session(self.assistant.messages)
        self.assistant.cancel_streaming()
        self.status.config(text=f"[{self._last_question_type.upper()}] Thinking...")
        self.assistant.stream_gpt_response(self.response_box, self.status, self.record_btn)

    # ------------------------------------------------------------------ #
    #  Native macOS picker: choose file or folder  (single dialog)      #
    # ------------------------------------------------------------------ #
    def _pick_files_or_folders_native(self) -> list:
        """
        Open a single native macOS NSOpenPanel (same panel Finder uses) that
        lets the user select ANY mix of files and folders in one go.

        Strategy (most reliable → least reliable):
          1. AppKit.NSOpenPanel  – already available because the app imports
             Quartz (both are part of PyObjC). No extra permissions required.
          2. osascript fallback  – in case PyObjC import somehow fails.
          3. tkinter filedialog  – files-only last resort.

        Returns a flat list of individual file paths (folders are expanded
        recursively). Returns [] on cancel or error.
        """

        raw_paths = []

        # ── 1. AppKit NSOpenPanel (primary) ──────────────────────────────
        try:
            import AppKit  # part of PyObjC, same framework as Quartz

            panel = AppKit.NSOpenPanel.openPanel()
            panel.setCanChooseFiles_(True)
            panel.setCanChooseDirectories_(True)
            panel.setAllowsMultipleSelection_(True)
            panel.setTitle_("Attach Files or Folders")
            panel.setPrompt_("Attach")
            panel.setMessage_("Select files or folders  •  ⌘-click to select multiple")

            # runModal() blocks on the main thread — fine here because this
            # method is always called from a tkinter button callback (main thread).
            ok = panel.runModal()
            if ok == 1:   # NSModalResponseOK
                raw_paths = [str(url.path()) for url in panel.URLs()]
            else:
                return []   # user cancelled

        except Exception:
            # ── 2. osascript fallback ─────────────────────────────────────
            try:
                apple_script = (
                    'set theItems to choose file or folder '
                    'with multiple selections allowed\n'
                    'set output to ""\n'
                    'repeat with anItem in theItems\n'
                    '    set output to output & POSIX path of anItem & "\\n"\n'
                    'end repeat\n'
                    'return output'
                )
                res = subprocess.run(
                    ['osascript', '-e', apple_script],
                    capture_output=True, text=True, timeout=120
                )
                if res.returncode == 0:
                    raw_paths = [p.strip() for p in res.stdout.splitlines() if p.strip()]
                else:
                    # ── 3. tkinter filedialog (files only, always works) ──
                    chosen = filedialog.askopenfilenames(
                        title="Select Files  (⌘-click for multiple)",
                        filetypes=[("All Files", "*.*")]
                    )
                    raw_paths = list(chosen)
            except Exception:
                raw_paths = []

        # ── Expand any selected folders recursively ───────────────────────
        all_files = []
        for path in raw_paths:
            if os.path.isdir(path):
                self.status.config(text=f"🔍 Scanning {os.path.basename(path)}…")
                self.update_idletasks()
                all_files.extend(collect_files_from_folder(path))
            elif os.path.isfile(path):
                all_files.append(path)
        return all_files

    def upload_resume(self):
        """
        Attach files or folders into the assistant context via a single native
        macOS Finder dialog. The user can select any mix of files and folders
        — just like a regular Open panel. Folders are walked recursively.
        Supports all common file types: .py, .pdf, .docx, .xlsx, .pptx, .txt, …
        """
        self.status.config(text="📂 Opening file picker…")
        self.update_idletasks()

        file_paths = self._pick_files_or_folders_native()

        if not file_paths:
            self.status.config(text="⚠️ No files selected.")
            return

        # ── Process every collected file ──────────────────────────────────
        self.status.config(text=f"⏳ Loading {len(file_paths)} file(s)…")
        self.update_idletasks()

        loaded_msgs, failed_msgs = [], []
        loaded_names = []                       # track filenames for the status bar
        for fp in file_paths:
            success, msg = self.assistant.load_document(fp)
            if success:
                loaded_msgs.append(msg)
                loaded_names.append(os.path.basename(fp))
            else:
                failed_msgs.append(msg)

        # ── Status bar — show actual filename(s), not just a count ───────
        parts = []
        if loaded_names:
            if len(loaded_names) == 1:
                # Single file: show full name
                parts.append(f"📎 {loaded_names[0]} attached")
            else:
                # Multiple files: show first name + how many more
                parts.append(f"📎 {loaded_names[0]} + {len(loaded_names) - 1} more attached")
        if failed_msgs:
            parts.append(f"⚠️ {len(failed_msgs)} skipped/failed")
        self.status.config(text="  |  ".join(parts) if parts else "Nothing loaded.")

        # Only pop a warning for genuine failures (not cleanly-skipped binaries)
        if failed_msgs:
            real_failures = [m for m in failed_msgs if not m.startswith("⚠️ Skipped")]
            if real_failures:
                messagebox.showwarning(
                    "Some files could not be loaded",
                    "\n".join(real_failures[:20])
                )

    def toggle_recording(self):
        with self.toggle_lock:
            if self.is_processing_audio:
                # Instead of blocking, we now force a reset if the user explicitly tries to record
                print("⚡️ Interrupting ongoing processing for new recording.")
                self.stop_output()                # Trigger stop
                self.is_processing_audio = False  # Reset
                self.assistant.streaming = False
                self.assistant.current_response = ""

            if not self.assistant.recorder.is_recording:
                self.assistant.streaming = False

                # Show listening + create a Live Question line
                self.response_box.config(state=tk.NORMAL)
                self.response_box.insert(tk.END, "\n\n🎙 Listening to your question...\n")
                # Remember where the live question line starts
                self.live_question_index = self.response_box.index(tk.END)
                self.response_box.insert(tk.END, "Live Question: ")
                self.response_box.config(state=tk.DISABLED)
                self.response_box.see(tk.END)

                self.assistant.recorder.start_recording()
                self.status.config(text="🎙 Listening to interviewer...")
                self.record_btn.config(text="🛑 Stop & Process")
                self.stop_btn.config(state=tk.DISABLED)

                # 🔴 Start live / incremental transcription
                self.live_transcription_running = True
                threading.Thread(target=self.live_transcription_loop, daemon=True).start()
                
                # 🎙 Start audio level indicator
                self.after(100, self.update_audio_level)

            else:
                self.is_processing_audio = True  # Set flag to True when processing audio
                self.record_btn.config(state=tk.DISABLED)
                self.stop_btn.config(state=tk.NORMAL)
                self.status.config(text="💭 Processing question...")
                threading.Thread(target=self.process_recording, daemon=True).start()



    def process_recording(self):
        # Stop live transcription immediately
        self.live_transcription_running = False
        
        try:
            # All Tkinter calls from this background thread must go through self.after().
            self.after(0, lambda: self.status.config(text="⏳ Getting complete question..."))

            filename = self.assistant.recorder.stop_recording()

            hint = (self.latest_live_question or "").strip()
            question = self.assistant.transcribe_audio(
                filename, prompt=hint[:450] if hint else None
            )

            if not question or question.startswith("❌"):
                msg = question if question else "⚠️ No speech detected"
                self.after(0, lambda m=msg: self.status.config(text=m))
                return

            question = question.strip()

            def _update_ui_with_question(q):
                # Remove the "Listening..." placeholder block
                content = self.response_box.get("1.0", tk.END)
                listening_idx = content.rfind("🎙 Listening to your question...")
                self.response_box.config(state=tk.NORMAL)
                if listening_idx != -1:
                    self.response_box.delete(f"1.0+{listening_idx}c", tk.END)
                self.response_box.insert(tk.END, f"\n\n---------------------------------------------------------------------\nQUESTION: {q}\n")
                self.response_box.config(state=tk.DISABLED)
                self.response_box.see(tk.END)

                self._last_question_type = classify_question(q)
                self.assistant.messages.append({"role": "user", "content": q})
                self.chat_manager.save_current_session(self.assistant.messages)
                self.status.config(text=f"[{self._last_question_type.upper()}] Thinking...")
                self.assistant.cancel_streaming()
                self.assistant.stream_gpt_response(self.response_box, self.status, self.record_btn)
                self.chat_manager.save_current_session(self.assistant.messages)

            self.after(0, lambda q=question: _update_ui_with_question(q))

        finally:
            self.is_processing_audio = False






    def start_new_chat(self):
        # Save current session if not empty
        if any(isinstance(m, dict) and m.get("role") == "user" for m in self.assistant.messages):
            # _generate_session_title() already handles both the new
            # "Attached document 'filename':" format and the legacy resume format.
            session_title = self._generate_session_title()

            self.chat_manager.add_session(session_title, self.assistant.messages.copy())
            self.chat_manager.save()
            self.load_chat_tabs()
            # 🔁 Auto-prune when over the limit
            self.auto_prune_chats(max_chats=10)

        # Start fresh session
        self.assistant.messages = [{
            "role": "system",
            "content": "You are a helpful interview assistant. "
                       "Provide detailed technical answers and ask follow-up questions when appropriate."
        }]
        self.response_box.config(state=tk.NORMAL)
        self.response_box.delete(1.0, tk.END)
        self.response_box.insert(tk.END, "🤖 New conversation started...")
        self.response_box.config(state=tk.DISABLED)
        self.status.config(text="🆕 New chat started")




    def stop_output(self):
        self.assistant.streaming = False
        self.stop_btn.config(state=tk.DISABLED)
        self.status.config(text="⏹ Output stopped")

    def increase_font(self):
        self.assistant.font_size = min(24, self.assistant.font_size + 1)
        self.response_box.config(font=('Consolas', self.assistant.font_size))

    def decrease_font(self):
        self.assistant.font_size = max(8, self.assistant.font_size - 1)
        self.response_box.config(font=('Consolas', self.assistant.font_size))

    def toggle_always_on_top(self):
        self.always_on_top = not self.always_on_top
        self.attributes("-topmost", self.always_on_top)
        self.topmost_btn.config(text="📌 Unpin Window" if self.always_on_top else "📌 Pin Window")

    # ── v2: Stealth Mode ──────────────────────────────────────────────────
    def toggle_stealth_mode(self):
        self.stealth_mode = not self.stealth_mode
        if _APPKIT_AVAILABLE:
            sharing = 4 if self.stealth_mode else 0  # 4 = NSWindowSharingNone
            try:
                for win in _AppKit.NSApp.windows():
                    win.setSharingType_(sharing)
            except Exception:
                pass
        self.stealth_btn.config(text="🕵️ Stealth ON" if self.stealth_mode else "🕵️ Stealth")
        self.status.config(text="Stealth ON — hidden from screen share" if self.stealth_mode
                           else "Stealth mode OFF")

    # ── v2: VAD Auto-Record ───────────────────────────────────────────────
    def toggle_vad_mode(self):
        self.vad_mode = not self.vad_mode
        if self.vad_mode:
            if not self._vad_running:
                self._vad_running = True
                threading.Thread(
                    target=self._vad_monitor_loop, daemon=True, name="VAD-Monitor"
                ).start()
            self.vad_btn.config(text="🎙 VAD ON")
            dev = self.assistant.recorder.find_device()
            dev_name = "default"
            try:
                dev_name = sd.query_devices(dev)["name"] if dev is not None else "default"
            except Exception:
                pass
            self.status.config(text=f"🎙 VAD active — monitoring: {dev_name}")
        else:
            self._vad_running = False   # ← signals the loop to exit cleanly
            self.vad_mode    = False
            self.vad_btn.config(text="🎙 VAD")
            self.status.config(text="🎙 VAD disabled")

    def _vad_monitor_loop(self):
        """
        Persistent VAD loop — opens ONE sd.InputStream on the correct device
        (BlackHole or external mic, matching what AudioRecorder uses) and
        monitors RMS energy to auto-start/stop recording.

        State machine
        ─────────────
        IDLE     → wait for ONSET_CHUNKS consecutive loud chunks (debounce)
                   → call toggle_recording()  [start]
        RECORDING → wait for SILENCE_SECS of quiet after MIN_RECORD_SECS
                   → call toggle_recording()  [stop]
        """
        # ── Tuning constants ───────────────────────────────────────────────
        SPEECH_RMS      = 0.012   # RMS threshold — tune if needed
        SILENCE_SECS    = 2.2     # silence before auto-stop
        MIN_RECORD_SECS = 1.5     # don't auto-stop before this many seconds
        ONSET_CHUNKS    = 3       # consecutive loud chunks to confirm speech start
        # ──────────────────────────────────────────────────────────────────

        # Resolve the SAME device the main recorder uses
        device_id = self.assistant.recorder.find_device()

        # Query device's native channel count (BlackHole 2ch → 2 channels)
        try:
            dev_info   = sd.query_devices(device_id) if device_id is not None \
                         else sd.query_devices(sd.default.device[0])
            n_channels = max(1, min(int(dev_info.get("max_input_channels", 1)), 8))
            native_sr  = int(dev_info.get("default_samplerate", 48000))
            dev_name   = dev_info.get("name", str(device_id))
        except Exception:
            n_channels, native_sr, dev_name = 1, 48000, str(device_id)

        print(f"🎙 VAD: opening stream on [{device_id}] {dev_name!r} "
              f"({n_channels}ch @ {native_sr} Hz)")

        rms_q   = queue.Queue(maxsize=200)
        BLOCK   = max(native_sr // 10, 512)   # ~100 ms per chunk

        def _cb(indata, frames, t, status):
            # Mix to mono regardless of how many channels the device has
            mono = indata.mean(axis=1) if indata.ndim > 1 else indata[:, 0]
            rms  = float(np.sqrt(np.mean(mono.astype(np.float32) ** 2)))
            try:
                rms_q.put_nowait(rms)
            except queue.Full:
                pass  # drop if consumer is too slow

        stream = None
        try:
            stream = sd.InputStream(
                device=device_id,
                samplerate=native_sr,
                channels=n_channels,
                dtype="float32",
                blocksize=BLOCK,
                callback=_cb,
            )
            stream.start()

            # ── State ──────────────────────────────────────────────────
            onset_count   = 0
            speaking      = False
            silence_start = None
            rec_start     = None

            while self._vad_running and self.vad_mode:
                try:
                    rms = rms_q.get(timeout=0.4)
                except queue.Empty:
                    continue

                is_rec  = self.assistant.recorder.is_recording   # ← FIXED (was .recording)
                is_busy = self.is_processing_audio

                # ── IDLE: looking for speech onset ──────────────────
                if not speaking:
                    if rms > SPEECH_RMS:
                        onset_count += 1
                        if onset_count >= ONSET_CHUNKS and not is_rec and not is_busy:
                            speaking      = True
                            silence_start = None
                            rec_start     = time.time()
                            self.after(0, self.toggle_recording)
                            self.after(0, lambda: self.status.config(
                                text=f"🎙 VAD: speech detected (RMS={rms:.4f}) — recording…"))
                            print(f"🎙 VAD start  rms={rms:.4f}  onset={onset_count}")
                    else:
                        onset_count = max(0, onset_count - 1)  # decay (not hard-reset)

                # ── SPEAKING: looking for end of speech ─────────────
                else:
                    if not is_rec:
                        # Recording stopped externally (user pressed Stop)
                        speaking = False; onset_count = 0; silence_start = None
                        continue

                    elapsed = time.time() - (rec_start or time.time())

                    if rms > SPEECH_RMS:
                        silence_start = None          # reset silence timer on any sound
                        onset_count   = ONSET_CHUNKS  # keep onset saturated
                    else:
                        if elapsed >= MIN_RECORD_SECS:
                            if silence_start is None:
                                silence_start = time.time()
                            elif time.time() - silence_start >= SILENCE_SECS:
                                speaking      = False
                                silence_start = None
                                onset_count   = 0
                                rec_start     = None
                                if self.assistant.recorder.is_recording:
                                    self.after(0, self.toggle_recording)
                                    self.after(0, lambda: self.status.config(
                                        text="🎙 VAD: silence detected — processing…"))
                                    print(f"🎙 VAD stop   silence={SILENCE_SECS}s")

        except Exception as exc:
            print(f"⚠️ VAD stream error: {exc}")
            self.after(0, lambda: self.status.config(text=f"⚠️ VAD error: {exc}"))
        finally:
            if stream is not None:
                try:
                    stream.stop()
                    stream.close()
                except Exception:
                    pass
            self._vad_running = False
            print("🎙 VAD loop exited")

    # ── v2: Local ASR ─────────────────────────────────────────────────────
    def toggle_local_asr(self):
        if not _FASTER_WHISPER_AVAILABLE:
            messagebox.showinfo("Local ASR",
                                "faster-whisper not installed.\nRun: pip install faster-whisper")
            return
        is_local = self.assistant.transcription_mode != "local"
        self.assistant.transcription_mode = "local" if is_local else "api"
        self.local_asr_btn.config(text="🔊 Local ASR ON" if is_local else "🔊 Local ASR")
        self.status.config(text=f"Transcription: {'local (faster-whisper)' if is_local else 'OpenAI Whisper API'}")

    # ── v2: Auto-Copy ─────────────────────────────────────────────────────
    def toggle_auto_copy(self):
        self.assistant.auto_copy_enabled = not self.assistant.auto_copy_enabled
        is_on = self.assistant.auto_copy_enabled
        self.auto_copy_btn.config(text="📋 Auto-Copy ON" if is_on else "📋 Auto-Copy")
        self.status.config(text=f"Auto-copy {'enabled — answers copied to clipboard' if is_on else 'disabled'}")

    # ── v2: Mock Interview ────────────────────────────────────────────────
    def open_mock_interview(self):
        role = "Software Engineer"
        try:
            role = self.tab_var.get() or role
        except Exception:
            pass
        MockInterviewDialog(self, client, role=role)

    # ── v2: Smart Model Routing ───────────────────────────────────────────
    def toggle_smart_route(self):
        """
        Toggle auto model routing ON/OFF.
        ON  → behavioral/general → gpt-4o-mini (2-3x faster, 10x cheaper)
              coding/system_design → gpt-4o (best quality)
        OFF → always use whatever model the 🧠 button shows.
        """
        self._auto_route_model = not self._auto_route_model
        is_on = self._auto_route_model
        self.route_btn.config(text="⚡ Smart Route ON" if is_on else "⚡ Smart Route")
        self.status.config(
            text="⚡ Smart routing ON — mini for behavioral, 4o for coding/design"
            if is_on else "⚡ Smart routing OFF — using fixed model"
        )

    def show_performance_dialog(self):
        """Show performance diagnostic dialog."""
        diag = self.assistant.diagnose_performance()
        self.assistant.print_performance_report()  # Also print to console
        
        dialog = tk.Toplevel(self)
        dialog.title("📊 Performance Diagnostics")
        dialog.geometry("450x500")
        dialog.transient(self)
        
        # Center on parent
        dialog.update_idletasks()
        x = self.winfo_x() + (self.winfo_width() // 2) - (450 // 2)
        y = self.winfo_y() + (self.winfo_height() // 2) - (500 // 2)
        dialog.geometry(f"+{x}+{y}")
        
        # Header
        ttk.Label(dialog, text="📊 Performance Analysis", font=('Arial', 14, 'bold')).pack(pady=10)
        
        # Stats frame
        stats_frame = ttk.LabelFrame(dialog, text="Current Chat Statistics")
        stats_frame.pack(fill="x", padx=15, pady=5)
        
        stats = [
            (f"Total Messages:", f"{diag['total_messages']}"),
            (f"  System:", f"{diag['system_messages']}"),
            (f"  User:", f"{diag['user_messages']}"),
            (f"  Assistant:", f"{diag['assistant_messages']}"),
            (f"  Images:", f"{diag['images_count']}"),
        ]
        
        for label, value in stats:
            row = ttk.Frame(stats_frame)
            row.pack(fill="x", padx=10, pady=2)
            ttk.Label(row, text=label).pack(side="left")
            ttk.Label(row, text=value, font=('Arial', 10, 'bold')).pack(side="right")
        
        # Token estimate frame
        token_frame = ttk.LabelFrame(dialog, text="Token Usage (Full Chat vs Optimized)")
        token_frame.pack(fill="x", padx=15, pady=5)
        
        token_stats = [
            ("Full Chat Tokens:", f"~{diag['estimated_total_tokens']:,}", diag['estimated_total_tokens'] > 30000),
            ("→ WILL SEND:", f"~{diag.get('will_send_tokens', 0):,} tokens", diag.get('will_send_tokens', 0) > 15000),
        ]
        
        for label, value, is_warning in token_stats:
            row = ttk.Frame(token_frame)
            row.pack(fill="x", padx=10, pady=2)
            ttk.Label(row, text=label).pack(side="left")
            color = 'red' if is_warning else 'green' if '→' in label else 'black'
            ttk.Label(row, text=value, font=('Arial', 10, 'bold'), foreground=color).pack(side="right")
        
        # Show reduction percentage
        if diag['estimated_total_tokens'] > 0:
            reduction = (1 - diag.get('will_send_tokens', 0) / diag['estimated_total_tokens']) * 100
            reduction_row = ttk.Frame(token_frame)
            reduction_row.pack(fill="x", padx=10, pady=2)
            ttk.Label(reduction_row, text="Token Reduction:").pack(side="left")
            ttk.Label(reduction_row, text=f"{reduction:.0f}% saved!", 
                      font=('Arial', 10, 'bold'), foreground='green').pack(side="right")
        
        # Optimization status
        opt_frame = ttk.LabelFrame(dialog, text="Optimization Status")
        opt_frame.pack(fill="x", padx=15, pady=5)
        
        opt_row = ttk.Frame(opt_frame)
        opt_row.pack(fill="x", padx=10, pady=5)
        opt_status = "⚡ ON (Fast Mode)" if diag['optimization_mode'] else "🐢 OFF (Full Context)"
        ttk.Label(opt_row, text="Mode:").pack(side="left")
        ttk.Label(opt_row, text=opt_status, font=('Arial', 10, 'bold')).pack(side="right")
        
        sent_row = ttk.Frame(opt_frame)
        sent_row.pack(fill="x", padx=10, pady=2)
        ttk.Label(sent_row, text="Messages sent to API:").pack(side="left")
        ttk.Label(sent_row, text=f"{diag['would_send_messages']}/{diag['total_messages']}", 
                  font=('Arial', 10, 'bold')).pack(side="right")
        
        has_summary = "Yes ✅" if self.assistant.summary_message else "No ❌"
        sum_row = ttk.Frame(opt_frame)
        sum_row.pack(fill="x", padx=10, pady=2)
        ttk.Label(sum_row, text="Has Summary:").pack(side="left")
        ttk.Label(sum_row, text=has_summary, font=('Arial', 10, 'bold')).pack(side="right")
        
        # Issues and recommendations
        if diag["issues"]:
            issues_frame = ttk.LabelFrame(dialog, text="⚠️ Issues Found")
            issues_frame.pack(fill="x", padx=15, pady=5)
            
            for issue in diag["issues"]:
                ttk.Label(issues_frame, text=issue, wraplength=400).pack(anchor="w", padx=10, pady=2)
            
            ttk.Label(issues_frame, text="\n💡 Recommendations:", font=('Arial', 10, 'bold')).pack(anchor="w", padx=10)
            for rec in diag["recommendations"]:
                ttk.Label(issues_frame, text=f"• {rec}", wraplength=400).pack(anchor="w", padx=20, pady=1)
        else:
            ttk.Label(dialog, text="✅ No performance issues detected!", 
                      font=('Arial', 11, 'bold'), foreground='green').pack(pady=10)
        
        # Action buttons
        btn_frame = ttk.Frame(dialog)
        btn_frame.pack(pady=15)
        
        def force_summarize():
            self.assistant._maybe_summarize_history()
            self.status.config(text="🔄 Triggering background summarization...")
            dialog.destroy()
        
        ttk.Button(btn_frame, text="🔄 Force Summarize", command=force_summarize).pack(side="left", padx=5)
        ttk.Button(btn_frame, text="🆕 New Chat", command=lambda: [self.start_new_chat(), dialog.destroy()]).pack(side="left", padx=5)
        ttk.Button(btn_frame, text="Close", command=dialog.destroy).pack(side="left", padx=5)
    
    # ============================================================================
    # CHAT SCROLL - PgUp/PgDn and Up/Down paragraph navigation
    # ============================================================================

    def _scroll_chat_to_top(self):
        """Scroll chat to top (PgUp hotkey)."""
        try:
            self.response_box.see("1.0")
            return "break"
        except Exception:
            return None

    def _scroll_chat_to_bottom(self):
        """Scroll chat to end (PgDn hotkey)."""
        try:
            self.response_box.see(tk.END)
            return "break"
        except Exception:
            return None

    def _scroll_chat_paragraph_up(self):
        """Scroll chat up by one paragraph (double-newline boundary)."""
        try:
            visible = self.response_box.index("@0,0")
            # Find previous "\n\n" (paragraph start) or go to 1.0
            prev = self.response_box.search("\n\n", visible, backwards=True, stopindex="1.0", regexp=False)
            if prev:
                self.response_box.see(prev)
            else:
                self.response_box.see("1.0")
            return "break"
        except Exception:
            return None

    def _scroll_chat_paragraph_down(self):
        """Scroll chat down by one paragraph (double-newline boundary)."""
        try:
            # Start from bottom of visible area to get "next" paragraph
            visible_bottom = self.response_box.index("@0,%d" % self.response_box.winfo_height())
            next_para = self.response_box.search("\n\n", visible_bottom, forwards=True, stopindex=tk.END, regexp=False)
            if next_para:
                self.response_box.see(next_para)
            else:
                self.response_box.see(tk.END)
            return "break"
        except Exception:
            return None

    def _on_up_key(self, event):
        """Up arrow: scroll chat one paragraph up only when focus is not in input (so typing is unaffected)."""
        try:
            if self.focus_get() == self.input_entry:
                return None  # Let Entry handle (cursor movement)
            return self._scroll_chat_paragraph_up()
        except Exception:
            return None

    def _on_down_key(self, event):
        """Down arrow: scroll chat one paragraph down only when focus is not in input."""
        try:
            if self.focus_get() == self.input_entry:
                return None
            return self._scroll_chat_paragraph_down()
        except Exception:
            return None

    # ============================================================================
    # BOOKMARK/POINTER SYSTEM - Quick navigation to questions (like debug breakpoints)
    # ============================================================================
    
    def add_bookmark_at_cursor(self):
        """
        Add a bookmark at the nearest QUESTION in the response box.
        Searches backward from current view to find 'QUESTION:' line.
        """
        try:
            # Get current visible position
            visible_start = self.response_box.index("@0,0")
            
            # Search backward for "QUESTION:" from current view
            question_pos = self.response_box.search(
                "QUESTION:", 
                visible_start, 
                backwards=True, 
                stopindex="1.0"
            )
            
            # If not found backwards, search forward
            if not question_pos:
                question_pos = self.response_box.search(
                    "QUESTION:", 
                    visible_start, 
                    forwards=True, 
                    stopindex=tk.END
                )
            
            if not question_pos:
                # No question found, bookmark current visible line
                question_pos = visible_start
                self._add_bookmark(question_pos, "📍 Manual mark")
            else:
                # Get the question text (rest of that line)
                line_end = f"{question_pos.split('.')[0]}.end"
                question_text = self.response_box.get(question_pos, line_end).strip()
                
                # Truncate for display
                if len(question_text) > 40:
                    question_text = question_text[:37] + "..."
                
                self._add_bookmark(question_pos, question_text)
                
        except Exception as e:
            print(f"Bookmark error: {e}")
            self.status.config(text=f"❌ Bookmark error: {e}")
    
    def add_bookmark_at_position(self, line_index: str, preview: str = ""):
        """Add a bookmark at a specific position (called programmatically)."""
        self._add_bookmark(line_index, preview or f"Q at line {line_index.split('.')[0]}")
    
    # ── Bookmark persistence helpers ──────────────────────────────────── #

    def _save_current_bookmarks(self):
        """Persist the current in-memory bookmarks to chats.json immediately."""
        self.chat_manager.update_session_bookmarks(
            self._current_chat_index,
            [[idx, preview] for idx, preview in self.bookmarks]
        )

    def _restore_bookmarks(self):
        """
        Re-apply saved bookmarks after display_chat_history() rebuilds the text.
        Called at the end of display_chat_history() and on app startup.
        """
        # Wipe any stale in-memory state first
        self.bookmarks.clear()
        self.bookmark_listbox.delete(0, tk.END)
        self._current_bookmark_index = -1

        saved = self.chat_manager.get_session_bookmarks(self._current_chat_index)
        for entry in saved:
            if len(entry) < 2:
                continue
            line_index, preview = str(entry[0]), str(entry[1])
            self.bookmarks.append((line_index, preview))
            self.bookmark_listbox.insert(tk.END, f"Q{len(self.bookmarks)}")
            self._highlight_bookmark(line_index)

    # ─────────────────────────────────────────────────────────────────── #

    def _add_bookmark(self, line_index: str, preview: str):
        """Internal method to add a bookmark."""
        # Check if already bookmarked (same line)
        line_num = line_index.split('.')[0]
        for existing_idx, _ in self.bookmarks:
            if existing_idx.split('.')[0] == line_num:
                self.status.config(text="⚠️ This line is already bookmarked")
                return
        
        # Add to bookmarks list
        bookmark_num = len(self.bookmarks) + 1
        self.bookmarks.append((line_index, preview))
        
        # Add to listbox (show number)
        self.bookmark_listbox.insert(tk.END, f"Q{bookmark_num}")
        
        # Highlight the bookmarked line in response box
        self._highlight_bookmark(line_index)
        
        self.status.config(text=f"🔖 Bookmark #{bookmark_num} added: {preview[:30]}...")
        print(f"📍 Bookmark #{bookmark_num} at {line_index}: {preview}")
        self._save_current_bookmarks()   # persist immediately
    
    def _highlight_bookmark(self, line_index: str):
        """Highlight a bookmarked line in the response box."""
        try:
            line_num = line_index.split('.')[0]
            line_start = f"{line_num}.0"
            line_end = f"{line_num}.end"
            
            self.response_box.config(state=tk.NORMAL)
            self.response_box.tag_add('bookmark_highlight', line_start, line_end)
            self.response_box.config(state=tk.DISABLED)
        except Exception as e:
            print(f"Highlight error: {e}")
    
    def go_to_next_bookmark(self):
        """Cycle to the next bookmark (hotkey F5). Jumps and flashes; wraps from last to first."""
        if not self.bookmarks:
            self.status.config(text="ℹ️ No bookmarks — add some with 🔖 or F4")
            return "break"
        self._current_bookmark_index = (self._current_bookmark_index + 1) % len(self.bookmarks)
        idx = self._current_bookmark_index
        line_index, preview = self.bookmarks[idx]
        self.response_box.see(line_index)
        self._flash_bookmark(line_index)
        self.bookmark_listbox.selection_clear(0, tk.END)
        self.bookmark_listbox.selection_set(idx)
        self.bookmark_listbox.see(idx)
        self.status.config(text=f"📍 Bookmark {idx + 1}/{len(self.bookmarks)}: {preview[:40]}...")
        return "break"

    def _on_bookmark_click(self, event=None):
        """Jump to the selected bookmark."""
        selection = self.bookmark_listbox.curselection()
        if not selection:
            return
        
        idx = selection[0]
        if idx < len(self.bookmarks):
            self._current_bookmark_index = idx  # keep hotkey cycle in sync
            line_index, preview = self.bookmarks[idx]
            
            # Jump to that line
            self.response_box.see(line_index)
            
            # Flash highlight the line briefly
            self._flash_bookmark(line_index)
            
            self.status.config(text=f"📍 Jumped to: {preview[:40]}...")
    
    def _flash_bookmark(self, line_index: str):
        """Flash a bookmark line to draw attention."""
        try:
            line_num = line_index.split('.')[0]
            line_start = f"{line_num}.0"
            line_end = f"{line_num}.end"
            
            # Create flash tag
            self.response_box.tag_configure('bookmark_flash', background='#ffff00', foreground='#000000')
            
            self.response_box.config(state=tk.NORMAL)
            self.response_box.tag_add('bookmark_flash', line_start, line_end)
            self.response_box.config(state=tk.DISABLED)
            
            # Remove flash after 500ms
            self.after(500, lambda: self._remove_flash(line_start, line_end))
        except Exception as e:
            print(f"Flash error: {e}")
    
    def _remove_flash(self, start, end):
        """Remove the flash highlight."""
        try:
            self.response_box.config(state=tk.NORMAL)
            self.response_box.tag_remove('bookmark_flash', start, end)
            self.response_box.config(state=tk.DISABLED)
        except:
            pass
    
    def _on_bookmark_delete(self, event=None):
        """Delete a bookmark on double-click."""
        selection = self.bookmark_listbox.curselection()
        if not selection:
            return
        
        idx = selection[0]
        if idx < len(self.bookmarks):
            line_index, preview = self.bookmarks[idx]
            
            # Remove highlight
            try:
                line_num = line_index.split('.')[0]
                self.response_box.config(state=tk.NORMAL)
                self.response_box.tag_remove('bookmark_highlight', f"{line_num}.0", f"{line_num}.end")
                self.response_box.config(state=tk.DISABLED)
            except:
                pass
            
            # Remove from list
            del self.bookmarks[idx]
            self.bookmark_listbox.delete(idx)
            
            # Renumber remaining bookmarks
            self._renumber_bookmarks()
            # Keep hotkey cycle index in range
            n = len(self.bookmarks)
            if n == 0:
                self._current_bookmark_index = -1
            elif self._current_bookmark_index >= n:
                self._current_bookmark_index = n - 1
            elif self._current_bookmark_index >= idx:
                self._current_bookmark_index = max(0, self._current_bookmark_index - 1)
            
            self.status.config(text=f"🗑 Bookmark removed: {preview[:30]}...")
            self._save_current_bookmarks()   # persist immediately

    def _renumber_bookmarks(self):
        """Renumber bookmarks after deletion."""
        self.bookmark_listbox.delete(0, tk.END)
        for i, (_, _) in enumerate(self.bookmarks):
            self.bookmark_listbox.insert(tk.END, f"Q{i+1}")
    
    def clear_all_bookmarks(self):
        """Clear all bookmarks."""
        if not self.bookmarks:
            self.status.config(text="ℹ️ No bookmarks to clear")
            return
        
        # Remove all highlights
        try:
            self.response_box.config(state=tk.NORMAL)
            self.response_box.tag_remove('bookmark_highlight', "1.0", tk.END)
            self.response_box.config(state=tk.DISABLED)
        except:
            pass
        
        # Clear list
        self.bookmarks.clear()
        self.bookmark_listbox.delete(0, tk.END)
        self._current_bookmark_index = -1
        self._save_current_bookmarks()   # persist the empty list

        self.status.config(text="🗑 All bookmarks cleared")
    
    def auto_bookmark_question(self, question_text: str):
        """
        Automatically add a bookmark when a new question is submitted.
        Called from submit_text_question or process_recording.
        """
        try:
            # Find the latest "QUESTION:" in the response box
            end_pos = self.response_box.index(tk.END)
            question_pos = self.response_box.search(
                "QUESTION:", 
                end_pos, 
                backwards=True, 
                stopindex="1.0"
            )
            
            if question_pos:
                preview = question_text[:40] + "..." if len(question_text) > 40 else question_text
                self._add_bookmark(question_pos, f"QUESTION: {preview}")
        except Exception as e:
            print(f"Auto-bookmark error: {e}")
    
    def _show_bookmark_menu(self, event):
        """Show context menu for bookmarking at click position."""
        # Get click position in text widget
        click_index = self.response_box.index(f"@{event.x},{event.y}")
        line_num = click_index.split('.')[0]
        
        # Get line content
        line_start = f"{line_num}.0"
        line_end = f"{line_num}.end"
        line_text = self.response_box.get(line_start, line_end).strip()
        
        # Create popup menu
        menu = tk.Menu(self, tearoff=0)
        
        # Check if this line is already bookmarked
        is_bookmarked = any(idx.split('.')[0] == line_num for idx, _ in self.bookmarks)
        
        if is_bookmarked:
            menu.add_command(
                label="🗑 Remove Bookmark", 
                command=lambda: self._remove_bookmark_at_line(line_num)
            )
        else:
            preview = line_text[:30] + "..." if len(line_text) > 30 else line_text
            menu.add_command(
                label=f"🔖 Bookmark Here", 
                command=lambda: self._add_bookmark(line_start, preview or f"Line {line_num}")
            )
        
        menu.add_separator()
        
        # Find nearest question
        nearest_q = self.response_box.search("QUESTION:", click_index, backwards=True, stopindex="1.0")
        if nearest_q:
            q_line = nearest_q.split('.')[0]
            q_text = self.response_box.get(nearest_q, f"{q_line}.end").strip()[:25]
            menu.add_command(
                label=f"🔖 Mark Question: {q_text}...", 
                command=lambda: self._add_bookmark(nearest_q, q_text)
            )
        
        menu.add_separator()
        menu.add_command(label="📋 Show All Questions", command=self._show_all_questions_dialog)

        # ── Copy question + images for use in other AI tools ─────────────
        if nearest_q:
            menu.add_separator()
            menu.add_command(
                label="📤 Copy Question + Images  (for Claude / ChatGPT)",
                command=lambda: self._copy_question_with_images(nearest_q)
            )

        # Show menu at click position
        try:
            menu.tk_popup(event.x_root, event.y_root)
        finally:
            menu.grab_release()
    
    def _remove_bookmark_at_line(self, line_num: str):
        """Remove bookmark at a specific line."""
        for i, (idx, _) in enumerate(self.bookmarks):
            if idx.split('.')[0] == line_num:
                # Remove highlight
                try:
                    self.response_box.config(state=tk.NORMAL)
                    self.response_box.tag_remove('bookmark_highlight', f"{line_num}.0", f"{line_num}.end")
                    self.response_box.config(state=tk.DISABLED)
                except:
                    pass
                
                del self.bookmarks[i]
                self._renumber_bookmarks()
                self._save_current_bookmarks()   # persist immediately
                self.status.config(text=f"🗑 Bookmark removed")
                return
    
    # ------------------------------------------------------------------ #
    #  Copy question + images  (for re-use in Claude / ChatGPT / etc.)  #
    # ------------------------------------------------------------------ #

    def _copy_question_with_images(self, question_pos: str):
        """
        Copy the CLEAN question text (no placeholders) to clipboard and save
        attached images to ~/Desktop as numbered PNGs. Then show an action
        dialog so the user can drag the images into Cursor / Claude.ai.

        question_pos – text-widget index of the 'QUESTION:' line (e.g. '14.0').
        """
        import base64, re as _re

        # ── 1. Count which QUESTION this is (1-based) in the response box ─
        #   We walk from the top and count every "QUESTION:" occurrence up to
        #   and including question_pos to determine the N-th user turn.
        q_count = 0
        search_from = "1.0"
        while True:
            found = self.response_box.search("QUESTION:", search_from, stopindex=tk.END)
            if not found:
                break
            q_count += 1
            # Stop once we reach or pass the target line
            if self.response_box.compare(found, ">=", question_pos):
                break
            search_from = f"{found} +1c"

        # ── 2. Get the N-th user message from self.assistant.messages ──────
        #   This gives us the raw, placeholder-free text + real image data.
        user_msgs = [
            m for m in self.assistant.messages
            if isinstance(m, dict) and m.get("role") == "user"
        ]
        # q_count is 1-based; clamp to available messages
        msg_index = max(0, min(q_count - 1, len(user_msgs) - 1))
        matched_msg = user_msgs[msg_index] if user_msgs else None

        # ── 3. Extract CLEAN text directly from the message ───────────────
        #   Avoids any "[📎 Image X]" or "[Image]" placeholder that the UI
        #   injects into the response box display.
        if matched_msg:
            content = matched_msg.get("content", "")
            if isinstance(content, str):
                clean_text = content.strip()
            elif isinstance(content, list):
                clean_text = "\n".join(
                    part.get("text", "").strip()
                    for part in content
                    if isinstance(part, dict) and part.get("type") == "text"
                    and part.get("text", "").strip()
                )
            else:
                clean_text = ""
        else:
            # Fallback: scrape from response box and strip placeholders
            q_line = question_pos.split('.')[0]
            sep = self.response_box.search("-----", f"{int(q_line)+1}.0", stopindex=tk.END)
            raw_q = self.response_box.get(
                f"{q_line}.0", sep if sep else tk.END
            ).strip()
            clean_text = _re.sub(r'^QUESTION:\s*', '', raw_q, flags=_re.IGNORECASE)
            clean_text = _re.sub(r'\[📎 Image \d+\]', '', clean_text)
            clean_text = _re.sub(r'\[Image\]', '', clean_text)
            clean_text = clean_text.strip()

        # ── 4. Copy clean text to clipboard ───────────────────────────────
        self.clipboard_clear()
        self.clipboard_append(clean_text)

        # ── 5. Extract images → save to ~/Desktop ─────────────────────────
        saved_paths = []
        if matched_msg and isinstance(matched_msg.get("content"), list):
            desktop = os.path.expanduser("~/Desktop")
            img_num = 1
            for part in matched_msg["content"]:
                if not (isinstance(part, dict) and part.get("type") == "image_url"):
                    continue
                url = part.get("image_url", {}).get("url", "")
                if not url.startswith("data:image"):
                    continue
                try:
                    header, b64data = url.split(",", 1)
                    ext = "png" if "png" in header else "jpg"
                    img_bytes = base64.b64decode(b64data)
                    fname = f"copied_question_img_{img_num}.{ext}"
                    fpath = os.path.join(desktop, fname)
                    with open(fpath, "wb") as f:
                        f.write(img_bytes)
                    saved_paths.append(fpath)
                    img_num += 1
                except Exception as e:
                    print(f"⚠️ Could not save image {img_num}: {e}")

        # ── 6. Show action dialog ──────────────────────────────────────────
        self._show_copy_result_dialog(clean_text, saved_paths)

    def _show_copy_result_dialog(self, clean_text: str, saved_paths: list):
        """
        Show a clear action dialog after copying a question.
        Tells the user exactly what to do with the text and images in
        Cursor / Claude.ai / ChatGPT.
        """
        dlg = tk.Toplevel(self)
        dlg.title("📤 Ready to paste into another AI tool")
        dlg.resizable(False, False)
        dlg.grab_set()
        dlg.lift()
        dlg.attributes('-topmost', True)

        # Centre over main window
        self.update_idletasks()
        w, h = 480, 340 if saved_paths else 220
        x = self.winfo_x() + (self.winfo_width()  - w) // 2
        y = self.winfo_y() + (self.winfo_height() - h) // 2
        dlg.geometry(f"{w}x{h}+{x}+{y}")

        pad = dict(padx=16, pady=6)

        # ── Text section ──────────────────────────────────────────────────
        ttk.Label(dlg, text="✅  Text copied to clipboard",
                  font=('Arial', 12, 'bold')).pack(anchor="w", **pad)
        ttk.Label(dlg, text="⌘V to paste in Cursor / Claude.ai / ChatGPT",
                  font=('Arial', 10), foreground='gray').pack(anchor="w", padx=16, pady=(0,4))

        # Preview of copied text
        preview_frame = ttk.Frame(dlg)
        preview_frame.pack(fill="x", padx=16, pady=(0, 8))
        preview = tk.Text(preview_frame, height=3, font=('Arial', 10),
                          wrap="word", state=tk.NORMAL,
                          background='#1e1e1e', foreground='#cccccc',
                          relief="flat", borderwidth=1)
        preview.insert("1.0", clean_text[:200] + ("…" if len(clean_text) > 200 else ""))
        preview.config(state=tk.DISABLED)
        preview.pack(fill="x")

        ttk.Separator(dlg, orient="horizontal").pack(fill="x", padx=16, pady=6)

        if saved_paths:
            # ── Images section ─────────────────────────────────────────────
            ttk.Label(dlg, text=f"🖼  {len(saved_paths)} image(s) saved to Desktop",
                      font=('Arial', 12, 'bold')).pack(anchor="w", **pad)

            for p in saved_paths:
                ttk.Label(dlg, text=f"   • {os.path.basename(p)}",
                          font=('Arial', 10), foreground='#888888').pack(anchor="w", padx=16)

            ttk.Label(dlg,
                      text="👆 Drag image(s) from Finder into Cursor / Claude.ai to attach",
                      font=('Arial', 10), foreground='#aaaaaa',
                      wraplength=440).pack(anchor="w", padx=16, pady=(6, 4))

            ttk.Separator(dlg, orient="horizontal").pack(fill="x", padx=16, pady=6)

            btn_frame = ttk.Frame(dlg)
            btn_frame.pack(pady=(0, 12))

            def open_finder():
                try:
                    subprocess.Popen(["open", "-R", saved_paths[0]])
                except Exception:
                    subprocess.Popen(["open", os.path.dirname(saved_paths[0])])

            ttk.Button(btn_frame, text="📂  Open in Finder",
                       command=open_finder, width=20).pack(side="left", padx=8)
            ttk.Button(btn_frame, text="✓  Done",
                       command=dlg.destroy, width=10).pack(side="left", padx=8)
        else:
            # Text-only — no images
            ttk.Button(dlg, text="✓  Done",
                       command=dlg.destroy, width=12).pack(pady=(0, 12))

        self.status.config(
            text=f"📋 Copied!  {len(saved_paths)} image(s) on Desktop."
            if saved_paths else "📋 Question text copied to clipboard!"
        )

    def _show_all_questions_dialog(self):
        """Show a dialog with all questions for easy bookmarking."""
        dialog = tk.Toplevel(self)
        dialog.title("📋 All Questions - Click to Bookmark")
        dialog.geometry("600x400")
        dialog.transient(self)
        
        # Center on parent
        dialog.update_idletasks()
        x = self.winfo_x() + (self.winfo_width() // 2) - (600 // 2)
        y = self.winfo_y() + (self.winfo_height() // 2) - (400 // 2)
        dialog.geometry(f"+{x}+{y}")
        
        # Find all questions
        questions = []
        search_start = "1.0"
        while True:
            pos = self.response_box.search("QUESTION:", search_start, stopindex=tk.END)
            if not pos:
                break
            
            line_num = pos.split('.')[0]
            line_end = f"{line_num}.end"
            q_text = self.response_box.get(pos, line_end).strip()
            
            # Check if already bookmarked
            is_bookmarked = any(idx.split('.')[0] == line_num for idx, _ in self.bookmarks)
            
            questions.append((pos, q_text, is_bookmarked))
            search_start = f"{int(line_num) + 1}.0"
        
        if not questions:
            ttk.Label(dialog, text="No questions found in chat", font=('Arial', 12)).pack(pady=20)
            return
        
        # Header
        ttk.Label(dialog, text=f"Found {len(questions)} questions. Click to bookmark:", font=('Arial', 11, 'bold')).pack(pady=5)
        
        # Scrollable list
        canvas = tk.Canvas(dialog)
        scrollbar = ttk.Scrollbar(dialog, orient="vertical", command=canvas.yview)
        scrollable_frame = ttk.Frame(canvas)
        
        scrollable_frame.bind("<Configure>", lambda e: canvas.configure(scrollregion=canvas.bbox("all")))
        canvas.create_window((0, 0), window=scrollable_frame, anchor="nw")
        canvas.configure(yscrollcommand=scrollbar.set)
        
        # Enable mouse wheel scrolling (macOS compatible)
        def _on_mousewheel(event):
            # macOS uses smaller delta values, Windows uses 120/-120
            if event.delta:
                # Normalize scroll direction
                scroll_amount = -1 if event.delta > 0 else 1
                canvas.yview_scroll(scroll_amount, "units")
        
        # Bind mouse wheel to canvas and dialog
        canvas.bind("<MouseWheel>", _on_mousewheel)
        dialog.bind("<MouseWheel>", _on_mousewheel)
        # Also bind to scrollable_frame for when mouse is over items
        scrollable_frame.bind("<MouseWheel>", _on_mousewheel)
        
        canvas.pack(side="left", fill="both", expand=True, padx=10)
        scrollbar.pack(side="right", fill="y")
        
        # Add each question
        for i, (pos, q_text, is_bookmarked) in enumerate(questions):
            frame = ttk.Frame(scrollable_frame)
            frame.pack(fill="x", pady=2)
            
            # Bookmark indicator
            indicator = "🔖" if is_bookmarked else "○"
            
            btn = ttk.Button(
                frame, 
                text=f"{indicator} Q{i+1}: {q_text[:60]}{'...' if len(q_text) > 60 else ''}", 
                command=lambda p=pos, t=q_text[:40]: self._bookmark_and_jump(p, t, dialog)
            )
            btn.pack(fill="x")
            
            # Bind mousewheel to each frame and button for scrolling
            frame.bind("<MouseWheel>", _on_mousewheel)
            btn.bind("<MouseWheel>", _on_mousewheel)
        
        # Close button
        ttk.Button(dialog, text="Close", command=dialog.destroy).pack(pady=10)
    
    def _bookmark_and_jump(self, pos: str, text: str, dialog=None):
        """Bookmark a question and jump to it."""
        # Add bookmark if not exists
        line_num = pos.split('.')[0]
        is_bookmarked = any(idx.split('.')[0] == line_num for idx, _ in self.bookmarks)
        
        if not is_bookmarked:
            self._add_bookmark(pos, f"QUESTION: {text}")
        
        # Jump to it
        self.response_box.see(pos)
        self._flash_bookmark(pos)
        
        if dialog:
            dialog.destroy()

    def toggle_optimization_mode(self):
        """
        Toggle Fast Mode (optimization) on/off.
        
        ON: Faster responses, compresses images, summarizes old chat
        OFF (default): Full context sent every time (slower but 100% complete)
        
        Your full chat history is ALWAYS preserved locally either way!
        """
        self.assistant.optimization_mode = not self.assistant.optimization_mode
        
        if self.assistant.optimization_mode:
            self.optimize_btn.config(text="⚡ Fast")
            self.status.config(text="⚡ Fast Mode ON - Optimized for speed")
        else:
            self.optimize_btn.config(text="🐢 Full")
            self.status.config(text="🐢 Full Mode - Sending complete context (slower)")

    # ============================================================================
    # BALANCE / BILLING MANAGEMENT
    # ============================================================================
    # DRAG AND DROP REORDERING
    # ============================================================================
    
    def _setup_drag_drop(self, tree: ttk.Treeview, tree_type: str):
        """
        Set up drag-and-drop reordering for a Treeview.
        tree_type: "tabs" or "chats"
        """
        tree._drag_data = {"item": None, "tree_type": tree_type}
        
        tree.bind("<ButtonPress-1>", lambda e: self._on_drag_start(e, tree))
        tree.bind("<B1-Motion>", lambda e: self._on_drag_motion(e, tree))
        tree.bind("<ButtonRelease-1>", lambda e: self._on_drag_release(e, tree))
    
    def _on_drag_start(self, event, tree: ttk.Treeview):
        """Start dragging an item."""
        item = tree.identify_row(event.y)
        if item:
            tree._drag_data["item"] = item
            tree._drag_data["start_y"] = event.y
            tree._drag_data["is_dragging"] = False  # Will be set to True if actually dragged
    
    def _on_drag_motion(self, event, tree: ttk.Treeview):
        """Show visual feedback while dragging."""
        if not tree._drag_data["item"]:
            return
        
        # Check if mouse has moved enough to be considered a drag (not just a click)
        start_y = tree._drag_data.get("start_y", event.y)
        if abs(event.y - start_y) > 5:  # 5 pixel threshold
            tree._drag_data["is_dragging"] = True
            tree.config(cursor="hand2")
    
    def _on_drag_release(self, event, tree: ttk.Treeview):
        """Drop the item at new position."""
        tree.config(cursor="")  # Reset cursor
        
        if not tree._drag_data["item"]:
            return
        
        dragged_item = tree._drag_data["item"]
        target_item = tree.identify_row(event.y)
        was_dragging = tree._drag_data.get("is_dragging", False)
        
        # Reset drag data
        tree._drag_data["item"] = None
        tree._drag_data["is_dragging"] = False
        
        # Only reorder if actually dragged (not just clicked)
        if not was_dragging:
            return
        
        if not target_item or target_item == dragged_item:
            return
        
        tree_type = tree._drag_data["tree_type"]
        
        if tree_type == "tabs":
            self._reorder_tabs(tree, dragged_item, target_item)
        elif tree_type == "chats":
            self._reorder_chats(tree, dragged_item, target_item)
    
    def _reorder_tabs(self, tree: ttk.Treeview, dragged: str, target: str):
        """Reorder tabs or subtabs."""
        try:
            # Determine if we're moving a tab or subtab
            dragged_is_tab = dragged.startswith("tab_")
            target_is_tab = target.startswith("tab_")
            
            if dragged_is_tab and target_is_tab:
                # Moving a tab to another tab position
                dragged_idx = int(dragged.split("_")[1])
                target_idx = int(target.split("_")[1])
                
                # Reorder in data
                tabs = self.prompt_manager.data["tabs"]
                if 0 <= dragged_idx < len(tabs) and 0 <= target_idx < len(tabs):
                    item = tabs.pop(dragged_idx)
                    tabs.insert(target_idx, item)
                    self.prompt_manager.save_tabs()
                    self.load_tabs()
                    self.status.config(text=f"📋 Moved tab to position {target_idx + 1}")
            
            elif dragged.startswith("sub_") and target.startswith("sub_"):
                # Moving a subtab within the same parent
                dragged_parts = dragged.split("_")
                target_parts = target.split("_")
                
                dragged_tab_idx = int(dragged_parts[1])
                dragged_sub_idx = int(dragged_parts[2])
                target_tab_idx = int(target_parts[1])
                target_sub_idx = int(target_parts[2])
                
                # Only allow reordering within the same tab
                if dragged_tab_idx == target_tab_idx:
                    subtabs = self.prompt_manager.data["tabs"][dragged_tab_idx]["subTabs"]
                    if 0 <= dragged_sub_idx < len(subtabs) and 0 <= target_sub_idx < len(subtabs):
                        item = subtabs.pop(dragged_sub_idx)
                        subtabs.insert(target_sub_idx, item)
                        self.prompt_manager.save_tabs()
                        self.load_tabs()
                        self.status.config(text=f"📋 Moved subtab to position {target_sub_idx + 1}")
                else:
                    self.status.config(text="⚠️ Can only reorder subtabs within same tab")
            
            elif dragged.startswith("sub_") and target_is_tab:
                # Moving subtab to a different tab
                dragged_parts = dragged.split("_")
                dragged_tab_idx = int(dragged_parts[1])
                dragged_sub_idx = int(dragged_parts[2])
                target_tab_idx = int(target.split("_")[1])
                
                if dragged_tab_idx != target_tab_idx:
                    source_subtabs = self.prompt_manager.data["tabs"][dragged_tab_idx]["subTabs"]
                    target_subtabs = self.prompt_manager.data["tabs"][target_tab_idx]["subTabs"]
                    
                    if 0 <= dragged_sub_idx < len(source_subtabs):
                        item = source_subtabs.pop(dragged_sub_idx)
                        target_subtabs.append(item)
                        self.prompt_manager.save_tabs()
                        self.load_tabs()
                        self.status.config(text=f"📋 Moved subtab to tab '{self.prompt_manager.get_tab_name(target_tab_idx)}'")
                        
        except Exception as e:
            print(f"Tab reorder error: {e}")
            self.status.config(text=f"❌ Reorder failed: {e}")
    
    def _reorder_chats(self, tree: ttk.Treeview, dragged: str, target: str):
        """Reorder chat sessions."""
        try:
            if not dragged.startswith("chat_") or not target.startswith("chat_"):
                return
            
            dragged_idx = int(dragged.split("_")[1])
            target_idx = int(target.split("_")[1])
            
            sessions = self.chat_manager.sessions
            if 0 <= dragged_idx < len(sessions) and 0 <= target_idx < len(sessions):
                item = sessions.pop(dragged_idx)
                sessions.insert(target_idx, item)
                self.chat_manager.save()
                self.load_chat_tabs()
                
                # Re-select the moved item
                new_id = f"chat_{target_idx}"
                if tree.exists(new_id):
                    tree.selection_set(new_id)
                    tree.see(new_id)
                
                self.status.config(text=f"💬 Moved chat to position {target_idx + 1}")
                
        except Exception as e:
            print(f"Chat reorder error: {e}")
            self.status.config(text=f"❌ Reorder failed: {e}")

if __name__ == "__main__":
    app = Application()
    style = ttk.Style()
    style.configure('TLabel', background='#343541', foreground='white')
    style.configure('TButton', font=('Arial', 12))
    
    def _restart_self():
        """
        Relaunch the current script using the same Python interpreter and args,
        then terminate this process (after closing Tk and the hotkey listener).
        """
        try:
            # Spawn the new process first
            python = sys.executable
            script = os.path.abspath(sys.argv[0])
            args = [python, script] + sys.argv[1:]
            subprocess.Popen(args)
        except Exception as e:
            print(f"❌ Restart spawn failed: {e}")
            return

        # Try to persist UI prefs before exit (optional but nice)
        try:
            app.save_ui_prefs()
        except Exception:
            pass

        # Stop listener if present
        try:
            listener.stop()
        except Exception:
            pass

        # Destroy Tk and hard-exit (to kill worker threads cleanly)
        try:
            app.destroy()
        except Exception:
            pass
        os._exit(0)


    # Define this AFTER app is created
    def setup_hotkey_listener():
        from pynput import keyboard
        combo_upload_resume = {keyboard.KeyCode(char='2'), keyboard.KeyCode(char='3')}
        combo_focus_chatbox = {keyboard.KeyCode(char='1'), keyboard.KeyCode(char='2')}
        combo_toggle_input_mode = {keyboard.KeyCode(char='3'), keyboard.KeyCode(char='4')}
        combo_listen_external = {keyboard.KeyCode(char='5'), keyboard.KeyCode(char='6')}
        combo_increase_font = {keyboard.Key.cmd,keyboard.Key.ctrl,  keyboard.KeyCode(char='=')}   # Cmd + +
        combo_decrease_font = {keyboard.Key.cmd, keyboard.Key.ctrl, keyboard.KeyCode(char='-')}   # Cmd + -
        combo_pin_window     = {keyboard.Key.cmd, keyboard.KeyCode(char='p')}  # Cmd + P
        combo_restart = {keyboard.Key.cmd, keyboard.Key.shift, keyboard.KeyCode(char='z')}
        combo_default_interview = {keyboard.Key.cmd, keyboard.Key.shift, keyboard.KeyCode(char='i')}  # Cmd+Shift+I: feed default interview instructions





        current_keys = set()
        
        def on_activate_toggle_input():
            print("🎚 Global hotkey 3 + 4: Toggle input device (internal/external)")
            app.focus_force()
            app.lift()
            app.attributes('-topmost', True)
            app.attributes('-topmost', False)
            app.toggle_input_mode()

            
        def on_activate_upload_resume():
            print("📁 Global hotkey 2 + 3: Trigger resume upload")
            app.focus_force()
            app.lift()
            app.attributes('-topmost', True)
            app.attributes('-topmost', False)
            app.upload_resume()

        def on_activate_focus_chatbox():
            print("⌨️ Global hotkey 1 + 2: Focus chat input")
            app.focus_force()
            app.lift()
            app.attributes('-topmost', True)
            app.attributes('-topmost', False)
            app.input_entry.focus_set()
            app.input_entry.mark_set(tk.INSERT, tk.END)

        def on_press(key):
            try:
                # Ignore Caps Lock to avoid crash (canonical/key handling can fail with it)
                if getattr(keyboard, 'Key', None) and key == getattr(keyboard.Key, 'caps_lock', None):
                    return
                if hasattr(key, 'name') and getattr(key, 'name', '') == 'caps_lock':
                    return
            except Exception:
                pass
            try:
                canonical_key = listener.canonical(key)
                if key not in combo_listen_external:
                    hotkey_listen.press(canonical_key)
                hotkey_stop.press(canonical_key)
                hotkey_screenshot.press(canonical_key)
            except Exception as e:
                # Avoid crash on keys that canonical() or HotKey can't handle (e.g. some special keys)
                pass

            if key in (combo_focus_chatbox | combo_upload_resume | combo_toggle_input_mode |
                    combo_listen_external | combo_increase_font | combo_decrease_font |
                    combo_pin_window | combo_restart | combo_default_interview):
                current_keys.add(key)

                if combo_focus_chatbox.issubset(current_keys):
                    on_activate_focus_chatbox()
                elif combo_upload_resume.issubset(current_keys):
                    on_activate_upload_resume()
                elif combo_toggle_input_mode.issubset(current_keys):
                    on_activate_toggle_input()
                elif combo_listen_external.issubset(current_keys):
                    on_activate_listen_external()
                elif combo_increase_font.issubset(current_keys):
                    print("🔎 Global hotkey Cmd + +: Increase font")
                    app.increase_font()
                elif combo_decrease_font.issubset(current_keys):
                    print("🔍 Global hotkey Cmd -: Decrease font")
                    app.decrease_font()
                elif combo_pin_window.issubset(current_keys):
                    print("📌 Global hotkey Cmd + P: Toggle pin")
                    app.toggle_always_on_top()
                elif combo_restart.issubset(current_keys):
                    on_activate_restart()
                elif combo_default_interview.issubset(current_keys):
                    print("📌 Global hotkey Cmd+Shift+I: Feed default interview instructions")
                    app.after(0, lambda: app.apply_default_interview_instructions())


                        

            
        def on_activate_restart():
            print("🔁 Global hotkey Cmd + R: Restarting app...")
            # Do it on a short timer so the print/status can flush
            try:
                app.status.config(text="🔁 Restarting...")
            except Exception:
                pass
            threading.Thread(target=_restart_self, daemon=True).start()

                    
        def on_activate_listen_external():
            if not app.assistant.recorder.is_recording:
                print("🎧 Global hotkey 5 + 6: Start Listening with External Mic")
                app.assistant.recorder.input_mode = "external"
                app.toggle_recording()
            else:
                print("🛑 Global hotkey 5 + 6: Stop and Process (External Mic)")
                app.toggle_recording()

                # Schedule switch back to BlackHole after recording is done
                def revert_input_mode():
                    time.sleep(1.5)  # slight delay to ensure stop finishes
                    app.assistant.recorder.input_mode = "internal"
                    app.status.config(text="🔈 Reverted to Internal Audio (BlackHole) after external session")
                    app.toggle_input_btn.config(text="🔈 Internal Audio (BlackHole)")

                threading.Thread(target=revert_input_mode, daemon=True).start()




            
            
            

        def on_release(key):
            try:
                if getattr(keyboard, 'Key', None) and key == getattr(keyboard.Key, 'caps_lock', None):
                    return
                if hasattr(key, 'name') and getattr(key, 'name', '') == 'caps_lock':
                    return
            except Exception:
                pass
            try:
                canonical_key = listener.canonical(key)
                hotkey_listen.release(canonical_key)
                hotkey_stop.release(canonical_key)
                hotkey_screenshot.release(canonical_key)
            except Exception:
                pass
            current_keys.discard(key)

        def on_activate_listen():
            if not app.assistant.recorder.is_recording:
                print("🎤 Global hotkey `: Start Listening")
                app.toggle_recording()
            else:
                print("🛑 Global hotkey `: Stop & Process")
                app.toggle_recording()

        def on_activate_stop():
            """Stop GPT generation and reset any processing flags."""
            print("🧠 Global hotkey ~: Force stop and reset state.")
            app.stop_output()
            app.is_processing_audio = False            # ✅ Reset processing state
            app.assistant.streaming = False            # ✅ Ensure streaming is stopped
            app.assistant.current_response = ""        # ✅ Clear any lingering response
            app.record_btn.config(text="🎤 Listen", state=tk.NORMAL)  # ✅ Reset button state
            app.stop_btn.config(state=tk.DISABLED)     # ✅ Ensure stop button is disabled
            app.status.config(text="✅ Stopped and Reset. Ready for next input.")


        def on_activate_screenshot():
            print("📸 Global hotkey !: Capturing full monitor under mouse and attaching...")
            try:
                mouse_x, mouse_y = pyautogui.position()
                screens = Quartz.NSScreen.screens()
                target_screen = None

                for screen in screens:
                    frame = screen.frame()
                    x = int(frame.origin.x)
                    y = int(frame.origin.y)
                    width = int(frame.size.width)
                    height = int(frame.size.height)

                    if x <= mouse_x < x + width and y <= mouse_y < y + height:
                        target_screen = (x, y, width, height)
                        break

                if not target_screen:
                    app.status.config(text="❌ No monitor found under mouse")
                    return

                screenshot = pyautogui.screenshot(region=target_screen)
                
                # Compress for faster transmission
                b64_image = compress_image_png(screenshot, max_size=1280)
                compressed_size = len(b64_image) // 1024
                print(f"📸 Screenshot compressed: {compressed_size}KB")

                if not hasattr(app, 'pending_attachments'):
                    app.pending_attachments = []

                app.pending_attachments.append({
                    "type": "image_url",
                    "image_url": {"url": f"data:image/png;base64,{b64_image}"}
                })

                app.status.config(text=f"📎 Screenshot attached ({compressed_size}KB). Press Enter to send.")
                app.input_entry.insert(tk.END, "📎 [Full screen screenshot ready to send] ")

            except Exception as e:
                app.status.config(text=f"❌ Screenshot error: {e}")
                print(f"❌ Screenshot error: {e}")

        hotkey_listen = keyboard.HotKey(keyboard.HotKey.parse('`'), on_activate_listen)
        hotkey_stop = keyboard.HotKey(keyboard.HotKey.parse('~'), on_activate_stop)
        hotkey_screenshot = keyboard.HotKey(keyboard.HotKey.parse('!'), on_activate_screenshot)

        listener = keyboard.Listener(on_press=on_press, on_release=on_release)
        return listener


    listener = setup_hotkey_listener()
    listener.start()

    app.mainloop()
    listener.join()

